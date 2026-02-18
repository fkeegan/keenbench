#!/usr/bin/env python3
import base64
import csv
import copy as pycopy
import datetime
import decimal
import hashlib
import io
import json
import math
import os
import re
import sys
import tempfile
import shutil
import subprocess
import threading
import time
import traceback
from typing import Any, Callable, Dict, List, Optional, Set, Tuple

JSONRPC_VERSION = "2.0"
ERROR_CODE = -32000

MAX_PREVIEW_BYTES = 10 * 1024 * 1024
MAX_PREVIEW_DIM = 2048
MIN_SCALE = 0.25
MAX_SCALE = 2.0
MAX_GRID_ROWS = 200
MAX_GRID_COLS = 50

# Chunk sizes for file maps
XLSX_CHUNK_ROWS = 50
DOCX_CHUNK_CHARS = 4000
PDF_CHUNK_PAGES = 5
TEXT_CHUNK_LINES = 200
TABULAR_CHUNK_ROWS = 100
TABULAR_CACHE_VERSION = 1
TABULAR_TABLE_NAME = "data"
TABULAR_TOTAL_COUNT_COLUMN = "__keenbench_total_row_count"
TABULAR_QUERY_TIMEOUT_MS = 5000
TABULAR_MEMORY_LIMIT_MB = 512
TABULAR_MAX_THREADS = 2

WORKBENCHES_DIR = ""
PREVIEW_RENDERER_PATH = ""
DEBUG_ENABLED = False
LOG_LEVEL = "info"
LOG_LEVELS = {"debug": 10, "info": 20, "warn": 30, "error": 40}
LOG_LOCK = threading.Lock()
WORKER_VERSION = "0.1.0"


def parse_bool(value: Any) -> bool:
    if value is None:
        return False
    text = str(value).strip().lower()
    return text in ("1", "true", "t", "yes", "y", "on")


def parse_int_clamped(value: Any, default: int, minimum: int, maximum: int) -> int:
    try:
        parsed = int(str(value).strip())
    except Exception:
        return default
    if parsed < minimum:
        return minimum
    if parsed > maximum:
        return maximum
    return parsed


def _should_log(level: str) -> bool:
    return LOG_LEVELS.get(level, 100) >= LOG_LEVELS.get(LOG_LEVEL, 20)


def _sanitize_value(value: Any) -> Any:
    if value is None or isinstance(value, (str, int, float, bool)):
        return value
    if isinstance(value, (list, tuple)):
        return [_sanitize_value(item) for item in value]
    if isinstance(value, dict):
        return {str(key): _sanitize_value(val) for key, val in value.items()}
    return str(value)


def _emit_log(level: str, message: str, **fields: Any) -> None:
    if not _should_log(level):
        return
    payload = {
        "ts": datetime.datetime.utcnow().isoformat(timespec="milliseconds") + "Z",
        "level": level,
        "message": message,
        "worker": "pyworker",
        "pid": os.getpid(),
    }
    for key, value in fields.items():
        payload[key] = _sanitize_value(value)
    line = json.dumps(payload, ensure_ascii=True)
    with LOG_LOCK:
        sys.stderr.write(line + "\n")
        sys.stderr.flush()


def log_debug(message: str, **fields: Any) -> None:
    _emit_log("debug", message, **fields)


def log_info(message: str, **fields: Any) -> None:
    _emit_log("info", message, **fields)


def log_warn(message: str, **fields: Any) -> None:
    _emit_log("warn", message, **fields)


def log_error(message: str, **fields: Any) -> None:
    _emit_log("error", message, **fields)


def _split_env_line(line: str) -> Tuple[Optional[str], Optional[str]]:
    if "=" not in line:
        return None, None
    key, value = line.split("=", 1)
    key = key.strip()
    if not key:
        return None, None
    value = value.strip()
    if len(value) >= 2 and ((value[0] == '"' and value[-1] == '"') or (value[0] == "'" and value[-1] == "'")):
        value = value[1:-1]
    return key, value


def _find_upwards(start: str, filename: str) -> str:
    current = start
    while True:
        candidate = os.path.join(current, filename)
        if os.path.isfile(candidate):
            return candidate
        parent = os.path.dirname(current)
        if parent == current:
            return ""
        current = parent


def load_env_file() -> Dict[str, Any]:
    override = os.environ.get("KEENBENCH_ENV_PATH", "").strip()
    path = override
    if not path:
        try:
            cwd = os.getcwd()
        except Exception as err:
            return {"loaded": False, "path": "", "keys": 0, "error": str(err)}
        path = _find_upwards(cwd, ".env")
    if not path:
        return {"loaded": False, "path": "", "keys": 0}
    keys = 0
    try:
        with open(path, "r", encoding="utf-8") as handle:
            for raw in handle:
                line = raw.strip()
                if not line or line.startswith("#"):
                    continue
                if line.startswith("export "):
                    line = line[len("export ") :].strip()
                key, value = _split_env_line(line)
                if not key:
                    continue
                if key in os.environ:
                    continue
                os.environ[key] = value or ""
                keys += 1
        return {"loaded": True, "path": path, "keys": keys}
    except Exception as err:
        return {"loaded": False, "path": path, "keys": keys, "error": str(err)}


def init_config() -> None:
    global WORKBENCHES_DIR, PREVIEW_RENDERER_PATH, DEBUG_ENABLED, LOG_LEVEL
    global TABULAR_QUERY_TIMEOUT_MS, TABULAR_MEMORY_LIMIT_MB, TABULAR_MAX_THREADS
    env_result = load_env_file()
    DEBUG_ENABLED = parse_bool(os.environ.get("KEENBENCH_DEBUG"))
    LOG_LEVEL = "debug" if DEBUG_ENABLED else "info"
    WORKBENCHES_DIR = os.environ.get("KEENBENCH_WORKBENCHES_DIR", "").strip()
    PREVIEW_RENDERER_PATH = os.environ.get("KEENBENCH_PREVIEW_RENDERER_PATH", "").strip()
    TABULAR_QUERY_TIMEOUT_MS = parse_int_clamped(
        os.environ.get("KEENBENCH_TABULAR_QUERY_TIMEOUT_MS"),
        TABULAR_QUERY_TIMEOUT_MS,
        100,
        600_000,
    )
    TABULAR_MEMORY_LIMIT_MB = parse_int_clamped(
        os.environ.get("KEENBENCH_TABULAR_MEMORY_LIMIT_MB"),
        TABULAR_MEMORY_LIMIT_MB,
        64,
        65_536,
    )
    TABULAR_MAX_THREADS = parse_int_clamped(
        os.environ.get("KEENBENCH_TABULAR_MAX_THREADS"),
        TABULAR_MAX_THREADS,
        1,
        64,
    )
    log_info(
        "worker.start",
        version=WORKER_VERSION,
        debug=DEBUG_ENABLED,
        workbenches_dir_set=bool(WORKBENCHES_DIR),
        preview_renderer_set=bool(PREVIEW_RENDERER_PATH),
        tabular_query_timeout_ms=TABULAR_QUERY_TIMEOUT_MS,
        tabular_memory_limit_mb=TABULAR_MEMORY_LIMIT_MB,
        tabular_max_threads=TABULAR_MAX_THREADS,
    )
    if env_result.get("loaded"):
        log_debug(
            "worker.env_loaded",
            path=env_result.get("path"),
            keys=env_result.get("keys"),
        )
    if env_result.get("error"):
        log_error(
            "worker.env_load_failed",
            path=env_result.get("path"),
            error=env_result.get("error"),
        )


class WorkerError(Exception):
    def __init__(self, code: str, message: str):
        super().__init__(message)
        self.code = code
        self.message = message


def send_response(resp: Dict[str, Any]) -> None:
    sys.stdout.write(json.dumps(resp))
    sys.stdout.write("\n")
    sys.stdout.flush()


def send_error(req_id: Any, code: str, message: str) -> None:
    payload = {
        "jsonrpc": JSONRPC_VERSION,
        "id": req_id,
        "error": {
            "code": ERROR_CODE,
            "message": message,
            "data": {"error_code": code},
        },
    }
    send_response(payload)


def send_result(req_id: Any, result: Any) -> None:
    payload = {
        "jsonrpc": JSONRPC_VERSION,
        "id": req_id,
        "result": result,
    }
    send_response(payload)


def require_workbenches_dir() -> str:
    if not WORKBENCHES_DIR:
        raise WorkerError("TOOL_WORKER_UNAVAILABLE", "workbenches dir not set")
    return WORKBENCHES_DIR


def validate_root(root: str) -> None:
    if not root:
        raise WorkerError("SANDBOX_VIOLATION", "missing root")
    if os.path.isabs(root):
        raise WorkerError("SANDBOX_VIOLATION", "absolute root not allowed")
    if root.startswith("..") or ".." in root:
        raise WorkerError("SANDBOX_VIOLATION", "invalid root")
    if "/" in root or "\\" in root:
        raise WorkerError("SANDBOX_VIOLATION", "invalid root")


def validate_write_root(root: str) -> None:
    validate_root(root)
    if root == "draft":
        return
    if root.startswith("draft.") and root.endswith(".staging"):
        return
    raise WorkerError("SANDBOX_VIOLATION", "writes must target draft")


def validate_rel_path(path: str) -> None:
    if not path:
        raise WorkerError("VALIDATION_FAILED", "missing path")
    if os.path.isabs(path):
        raise WorkerError("SANDBOX_VIOLATION", "absolute path not allowed")
    if path.startswith("..") or ".." in path:
        raise WorkerError("SANDBOX_VIOLATION", "invalid path")
    if "/" in path or "\\" in path:
        raise WorkerError("SANDBOX_VIOLATION", "nested paths are not allowed")


def require_workbench_id(params: Dict[str, Any]) -> str:
    workbench_id = (params.get("workbench_id") or "").strip()
    if not workbench_id:
        raise WorkerError("VALIDATION_FAILED", "missing workbench_id")
    if "/" in workbench_id or "\\" in workbench_id:
        raise WorkerError("VALIDATION_FAILED", "invalid workbench_id")
    return workbench_id


def _resolve_rel_path(workbench_id: str, root: str, rel_path: str) -> str:
    base = os.path.join(require_workbenches_dir(), workbench_id, root)
    full = os.path.abspath(os.path.join(base, rel_path))
    base_abs = os.path.abspath(base)
    if not full.startswith(base_abs + os.sep):
        raise WorkerError("SANDBOX_VIOLATION", "path escapes sandbox")
    return full


def resolve_named_path(
    params: Dict[str, Any],
    path_key: str,
    root_key: str = "root",
    default_root: str = "draft",
    require_write: bool = False,
) -> str:
    workbench_id = require_workbench_id(params)
    root = (params.get(root_key) or default_root).strip()
    if require_write:
        validate_write_root(root)
    else:
        validate_root(root)
    rel_path = (params.get(path_key) or "").strip()
    validate_rel_path(rel_path)
    return _resolve_rel_path(workbench_id, root, rel_path)


def resolve_path(params: Dict[str, Any]) -> str:
    return resolve_named_path(params, "path")


def parse_non_negative_index(value: Any, name: str, default: int = 0) -> int:
    if value is None:
        return default
    try:
        idx = int(value)
    except Exception:
        raise WorkerError("VALIDATION_FAILED", f"invalid {name}")
    if idx < 0:
        raise WorkerError("VALIDATION_FAILED", f"invalid {name}")
    return idx


def clamp_scale(value: Any) -> float:
    try:
        scale = float(value)
    except Exception:
        scale = 1.0
    if scale < MIN_SCALE:
        return MIN_SCALE
    if scale > MAX_SCALE:
        return MAX_SCALE
    return scale


def import_module(name: str):
    try:
        return __import__(name)
    except Exception:
        raise WorkerError("TOOL_WORKER_UNAVAILABLE", f"missing dependency: {name}")


def _slug(value: Any) -> str:
    text = str(value or "").strip().lower()
    if not text:
        return ""
    chars: List[str] = []
    for ch in text:
        if ch.isalnum():
            chars.append(ch)
        elif ch in (" ", "-", "_", ":"):
            chars.append("_")
    return "".join(chars).strip("_")


def _asset_selector_text(asset: Any) -> str:
    if isinstance(asset, str):
        return asset.strip()
    if isinstance(asset, dict):
        for key in ("asset_id", "selector", "id", "name", "type"):
            value = asset.get(key)
            if value is not None and str(value).strip():
                return str(value).strip()
    return str(asset)


def _parse_asset_list(params: Dict[str, Any]) -> List[Any]:
    assets = params.get("assets")
    if not isinstance(assets, list) or len(assets) == 0:
        raise WorkerError("VALIDATION_FAILED", "missing assets")
    return assets


def _sha1_of_bytes(data: bytes) -> str:
    h = hashlib.sha1()
    h.update(data)
    return h.hexdigest()


def _coerce_int(value: Any, default: Optional[int] = None) -> Optional[int]:
    if value is None:
        return default
    try:
        return int(value)
    except Exception:
        return default


_HEX_RGB_RE = re.compile(r"^[0-9A-Fa-f]{6}$")


def _coerce_float(value: Any, default: Optional[float] = None) -> Optional[float]:
    if value is None:
        return default
    try:
        number = float(value)
    except Exception:
        return default
    if not math.isfinite(number):
        return default
    return number


def _coerce_bool_strict(value: Any, default: Optional[bool] = None) -> Optional[bool]:
    if value is None:
        return default
    if isinstance(value, bool):
        return value
    if isinstance(value, (int, float)):
        if value == 1:
            return True
        if value == 0:
            return False
        return default
    if isinstance(value, str):
        text = value.strip().lower()
        if text in ("true", "t", "yes", "y", "1", "on"):
            return True
        if text in ("false", "f", "no", "n", "0", "off"):
            return False
    return default


def _normalize_hex_rgb(value: Any) -> Optional[str]:
    if not isinstance(value, str):
        return None
    text = value.strip()
    if text.startswith("#"):
        text = text[1:]
    if len(text) == 8 and _HEX_RGB_RE.match(text[2:]):
        text = text[2:]
    if not _HEX_RGB_RE.match(text):
        return None
    return text.upper()


def _style_warn(scope: str, field: str, value: Any, reason: str, **fields: Any) -> None:
    payload: Dict[str, Any] = {
        "scope": scope,
        "field": field,
        "value": _sanitize_value(value),
        "reason": reason,
    }
    for key, val in fields.items():
        payload[key] = _sanitize_value(val)
    log_warn("style.skipped", **payload)


# ------------------ DOCX ------------------

def docx_apply_ops(params: Dict[str, Any]) -> Dict[str, Any]:
    docx = import_module("docx")
    path = resolve_path(params)
    ops = params.get("ops") or []
    create_new = params.get("create_new", False)
    copy_from = params.get("copy_from")

    if not isinstance(ops, list) or len(ops) == 0:
        raise WorkerError("VALIDATION_FAILED", "missing ops")

    # Handle create_new with optional copy_from
    if create_new:
        if copy_from:
            source_params = dict(params)
            source_params["path"] = copy_from
            source_path = resolve_path(source_params)
            if not os.path.exists(source_path):
                raise WorkerError("FILE_READ_FAILED", f"source file not found: {copy_from}")
            doc = docx.Document(source_path)
            log_info("docx.copy_from", source=copy_from, target=path)
        else:
            doc = docx.Document()
            log_info("docx.create_new", target=path)
    elif os.path.exists(path):
        doc = docx.Document(path)
    else:
        doc = docx.Document()
    for op in ops:
        if not isinstance(op, dict):
            raise WorkerError("VALIDATION_FAILED", "invalid op")
        name = op.get("op")
        if name == "set_paragraphs":
            doc = docx.Document()
            paragraphs = op.get("paragraphs") or []
            if not isinstance(paragraphs, list):
                raise WorkerError("VALIDATION_FAILED", "paragraphs must be list")
            for idx, item in enumerate(paragraphs):
                if not isinstance(item, dict):
                    para = doc.add_paragraph(str(item))
                    continue
                para = doc.add_paragraph("")
                _docx_fill_paragraph(para, item, op_name=name, paragraph_index=idx)
        elif name == "append_paragraph":
            para = doc.add_paragraph("")
            _docx_fill_paragraph(para, op, op_name=name)
        elif name == "replace_text":
            search = op.get("search", "")
            replace = op.get("replace", "")
            match_case = bool(op.get("match_case", False))
            if not search:
                continue
            for para in _docx_iter_replaceable_paragraphs(doc):
                text = para.text
                if match_case:
                    if search in text:
                        para.text = text.replace(search, replace)
                else:
                    lower = text.lower()
                    if search.lower() in lower:
                        para.text = _replace_case_insensitive(text, search, replace)
        else:
            raise WorkerError("VALIDATION_FAILED", f"unsupported op: {name}")
    os.makedirs(os.path.dirname(path), exist_ok=True)
    doc.save(path)
    return {"ok": True}


def _docx_fill_paragraph(para: Any, item: Dict[str, Any], op_name: str, paragraph_index: Optional[int] = None) -> None:
    from docx.shared import Pt, Inches, RGBColor
    from docx.enum.text import WD_ALIGN_PARAGRAPH, WD_COLOR_INDEX

    allowed_keys = {
        "text",
        "style",
        "runs",
        "alignment",
        "space_before",
        "space_after",
        "line_spacing",
        "indent_left",
        "indent_right",
        "indent_first_line",
        "op",
    }
    for key in item.keys():
        if key not in allowed_keys:
            _style_warn(
                "docx.paragraph",
                key,
                item.get(key),
                "unknown style key",
                op=op_name,
                paragraph_index=paragraph_index,
            )

    style_name = item.get("style")
    if style_name is not None:
        text = str(style_name).strip()
        if text:
            try:
                para.style = text
            except Exception as err:
                _style_warn(
                    "docx.paragraph",
                    "style",
                    style_name,
                    "invalid style name",
                    op=op_name,
                    paragraph_index=paragraph_index,
                    error=str(err),
                )

    runs = item.get("runs")
    if isinstance(runs, list):
        for run_index, run_item in enumerate(runs):
            if isinstance(run_item, dict):
                run_text = str(run_item.get("text") or "")
                run = para.add_run(run_text)
                _docx_apply_run_style(
                    run,
                    run_item,
                    WD_COLOR_INDEX,
                    RGBColor,
                    Pt,
                    op_name=op_name,
                    paragraph_index=paragraph_index,
                    run_index=run_index,
                )
            else:
                _style_warn(
                    "docx.run",
                    "runs",
                    run_item,
                    "run entry must be object; coerced to plain text",
                    op=op_name,
                    paragraph_index=paragraph_index,
                    run_index=run_index,
                )
                para.add_run(str(run_item))
    else:
        if runs is not None:
            _style_warn(
                "docx.paragraph",
                "runs",
                runs,
                "runs must be list",
                op=op_name,
                paragraph_index=paragraph_index,
            )
        para.add_run(str(item.get("text") or ""))

    _docx_apply_paragraph_format(
        para,
        item,
        WD_ALIGN_PARAGRAPH,
        Pt,
        Inches,
        op_name=op_name,
        paragraph_index=paragraph_index,
    )


def _docx_apply_run_style(
    run: Any,
    payload: Dict[str, Any],
    wd_color_index: Any,
    rgb_color: Any,
    pt_ctor: Any,
    op_name: str,
    paragraph_index: Optional[int],
    run_index: int,
) -> None:
    allowed_keys = {"text", "font_name", "font_size", "bold", "italic", "underline", "font_color", "highlight_color"}
    for key in payload.keys():
        if key not in allowed_keys:
            _style_warn(
                "docx.run",
                key,
                payload.get(key),
                "unknown style key",
                op=op_name,
                paragraph_index=paragraph_index,
                run_index=run_index,
            )

    font = run.font
    font_name = payload.get("font_name")
    if font_name is not None:
        if isinstance(font_name, str) and font_name.strip():
            font.name = font_name.strip()
        else:
            _style_warn(
                "docx.run",
                "font_name",
                font_name,
                "font_name must be non-empty string",
                op=op_name,
                paragraph_index=paragraph_index,
                run_index=run_index,
            )

    font_size = payload.get("font_size")
    if font_size is not None:
        size_value = _coerce_float(font_size, None)
        if size_value is None or size_value <= 0:
            _style_warn(
                "docx.run",
                "font_size",
                font_size,
                "font_size must be a positive number",
                op=op_name,
                paragraph_index=paragraph_index,
                run_index=run_index,
            )
        else:
            font.size = pt_ctor(size_value)

    for field, attr in (("bold", "bold"), ("italic", "italic"), ("underline", "underline")):
        raw = payload.get(field)
        if raw is None:
            continue
        parsed = _coerce_bool_strict(raw, None)
        if parsed is None:
            _style_warn(
                "docx.run",
                field,
                raw,
                f"{field} must be boolean",
                op=op_name,
                paragraph_index=paragraph_index,
                run_index=run_index,
            )
            continue
        setattr(run, attr, parsed)

    font_color = payload.get("font_color")
    if font_color is not None:
        rgb = _normalize_hex_rgb(font_color)
        if rgb is None:
            _style_warn(
                "docx.run",
                "font_color",
                font_color,
                "font_color must be #RRGGBB",
                op=op_name,
                paragraph_index=paragraph_index,
                run_index=run_index,
            )
        else:
            try:
                font.color.rgb = rgb_color.from_string(rgb)
            except Exception as err:
                _style_warn(
                    "docx.run",
                    "font_color",
                    font_color,
                    "failed to apply font_color",
                    op=op_name,
                    paragraph_index=paragraph_index,
                    run_index=run_index,
                    error=str(err),
                )

    highlight_color = payload.get("highlight_color")
    if highlight_color is not None:
        parsed = _docx_parse_highlight_color(highlight_color, wd_color_index)
        if parsed is None:
            _style_warn(
                "docx.run",
                "highlight_color",
                highlight_color,
                "unsupported highlight_color",
                op=op_name,
                paragraph_index=paragraph_index,
                run_index=run_index,
            )
        else:
            try:
                font.highlight_color = parsed
            except Exception as err:
                _style_warn(
                    "docx.run",
                    "highlight_color",
                    highlight_color,
                    "failed to apply highlight_color",
                    op=op_name,
                    paragraph_index=paragraph_index,
                    run_index=run_index,
                    error=str(err),
                )


def _docx_parse_highlight_color(value: Any, wd_color_index: Any) -> Any:
    if not isinstance(value, str) or not value.strip():
        return None
    key = _slug(value)
    for attr in dir(wd_color_index):
        if not attr.isupper():
            continue
        if _slug(attr) == key:
            try:
                return getattr(wd_color_index, attr)
            except Exception:
                continue
    return None


def _docx_apply_paragraph_format(
    para: Any,
    payload: Dict[str, Any],
    wd_align_paragraph: Any,
    pt_ctor: Any,
    inches_ctor: Any,
    op_name: str,
    paragraph_index: Optional[int],
) -> None:
    paragraph_format = para.paragraph_format
    alignment_value = payload.get("alignment")
    if alignment_value is not None:
        if not isinstance(alignment_value, str):
            _style_warn(
                "docx.paragraph",
                "alignment",
                alignment_value,
                "alignment must be string",
                op=op_name,
                paragraph_index=paragraph_index,
            )
        else:
            alignment_map = {
                "left": wd_align_paragraph.LEFT,
                "center": wd_align_paragraph.CENTER,
                "right": wd_align_paragraph.RIGHT,
                "justify": wd_align_paragraph.JUSTIFY,
            }
            alignment = alignment_map.get(_slug(alignment_value))
            if alignment is None:
                _style_warn(
                    "docx.paragraph",
                    "alignment",
                    alignment_value,
                    "unsupported alignment",
                    op=op_name,
                    paragraph_index=paragraph_index,
                )
            else:
                paragraph_format.alignment = alignment

    for field in ("space_before", "space_after"):
        raw_value = payload.get(field)
        if raw_value is None:
            continue
        parsed = _coerce_float(raw_value, None)
        if parsed is None or parsed < 0:
            _style_warn(
                "docx.paragraph",
                field,
                raw_value,
                f"{field} must be a non-negative number",
                op=op_name,
                paragraph_index=paragraph_index,
            )
            continue
        setattr(paragraph_format, field, pt_ctor(parsed))

    line_spacing = payload.get("line_spacing")
    if line_spacing is not None:
        parsed = _coerce_float(line_spacing, None)
        if parsed is None or parsed <= 0:
            _style_warn(
                "docx.paragraph",
                "line_spacing",
                line_spacing,
                "line_spacing must be a positive number",
                op=op_name,
                paragraph_index=paragraph_index,
            )
        else:
            paragraph_format.line_spacing = parsed

    for src_key, target_attr in (
        ("indent_left", "left_indent"),
        ("indent_right", "right_indent"),
        ("indent_first_line", "first_line_indent"),
    ):
        raw_value = payload.get(src_key)
        if raw_value is None:
            continue
        parsed = _coerce_float(raw_value, None)
        if parsed is None:
            _style_warn(
                "docx.paragraph",
                src_key,
                raw_value,
                f"{src_key} must be a number",
                op=op_name,
                paragraph_index=paragraph_index,
            )
            continue
        setattr(paragraph_format, target_attr, inches_ctor(parsed))


def _replace_case_insensitive(text: str, search: str, replace: str) -> str:
    lower = text.lower()
    target = search.lower()
    out = []
    i = 0
    while i < len(text):
        if lower[i : i + len(target)] == target:
            out.append(replace)
            i += len(target)
        else:
            out.append(text[i])
            i += 1
    return "".join(out)


def _docx_iter_replaceable_paragraphs(doc):
    for para in doc.paragraphs:
        yield para

    for table in doc.tables:
        for row in table.rows:
            for cell in row.cells:
                for para in cell.paragraphs:
                    yield para

    for section in doc.sections:
        for para in section.header.paragraphs:
            yield para
        for para in section.footer.paragraphs:
            yield para


def docx_extract_text(params: Dict[str, Any]) -> Dict[str, Any]:
    """Extract text from docx, optionally filtered by section heading."""
    docx = import_module("docx")
    path = resolve_path(params)
    doc = docx.Document(path)
    section_filter = params.get("section")

    if section_filter is not None:
        # Extract text for a specific section only
        return _docx_extract_section(doc, section_filter)

    lines: List[str] = []
    for para in doc.paragraphs:
        if para.text:
            lines.append(para.text)
    for table in doc.tables:
        for row in table.rows:
            row_vals = []
            for cell in row.cells:
                row_vals.append(cell.text)
            lines.append("\t".join(row_vals))
    return {"text": "\n".join(lines)}


def _docx_extract_section(doc, section_filter) -> Dict[str, Any]:
    """Extract text from a specific section of a docx document."""
    sections = _docx_build_sections(doc)

    # Find matching section by heading text or index
    target = None
    if isinstance(section_filter, int) or (isinstance(section_filter, str) and section_filter.isdigit()):
        idx = int(section_filter)
        if 0 <= idx < len(sections):
            target = sections[idx]
    else:
        section_filter_lower = str(section_filter).lower()
        for s in sections:
            if s["heading"].lower() == section_filter_lower:
                target = s
                break

    if target is None:
        raise WorkerError("VALIDATION_FAILED", f"section not found: {section_filter}")

    section_idx = sections.index(target)
    total_sections = len(sections)

    lines: List[str] = []
    for para in target.get("paragraphs", []):
        if not isinstance(para, dict):
            continue
        text = str(para.get("text") or "").strip()
        if text:
            lines.append(text)
    for table in target.get("tables", []):
        if not isinstance(table, dict):
            continue
        rows = table.get("rows", [])
        if not isinstance(rows, list):
            continue
        for row in rows:
            if not isinstance(row, list):
                continue
            row_vals = [str(cell or "") for cell in row]
            lines.append("\t".join(row_vals))

    result: Dict[str, Any] = {"text": "\n".join(lines)}
    if total_sections > 1:
        result["chunk_info"] = {
            "chunk_index": section_idx,
            "total_chunks": total_sections,
            "has_more": section_idx < total_sections - 1,
            "range": target["heading"],
        }
    return result


def docx_render_page(params: Dict[str, Any]) -> Dict[str, Any]:
    path = resolve_path(params)
    page_index = int(params.get("page_index", 0))
    scale = clamp_scale(params.get("scale", 1.0))
    try:
        pdf_path, tmp_dir = convert_to_pdf(path)
    except WorkerError as err:
        if err.code not in ("TOOL_WORKER_UNAVAILABLE", "FILE_READ_FAILED"):
            raise
        log_info("preview.docx_fallback", path=path, reason=err.message)
        return render_docx_page_fallback(path, page_index, scale)
    except Exception as err:
        log_info("preview.docx_fallback", path=path, reason=str(err))
        return render_docx_page_fallback(path, page_index, scale)
    try:
        return render_pdf_page(pdf_path, page_index, scale)
    finally:
        shutil.rmtree(tmp_dir, ignore_errors=True)


def docx_get_map(params: Dict[str, Any]) -> Dict[str, Any]:
    """Return structural map of a docx file: sections, tables, images, metadata."""
    docx = import_module("docx")
    path = resolve_path(params)
    doc = docx.Document(path)

    # Build sections based on heading styles
    sections: List[Dict[str, Any]] = []
    current_heading: Optional[str] = None
    current_level: int = 0
    current_chars: int = 0
    current_has_tables: bool = False
    current_has_images: bool = False
    current_para_count: int = 0
    total_char_count: int = 0

    def _flush_section():
        nonlocal current_heading, current_level, current_chars
        nonlocal current_has_tables, current_has_images, current_para_count
        if current_heading is None and current_chars == 0:
            return
        section: Dict[str, Any] = {
            "heading": current_heading or "(untitled)",
            "level": current_level,
            "char_count": current_chars,
            "has_tables": current_has_tables,
            "has_images": current_has_images,
        }
        # Add chunks for large sections
        if current_chars > DOCX_CHUNK_CHARS:
            chunk_count = (current_chars + DOCX_CHUNK_CHARS - 1) // DOCX_CHUNK_CHARS
            chunks = []
            for i in range(chunk_count):
                start_char = i * DOCX_CHUNK_CHARS
                end_char = min((i + 1) * DOCX_CHUNK_CHARS, current_chars)
                chunks.append({
                    "index": i,
                    "char_count": end_char - start_char,
                })
            section["chunks"] = chunks
        sections.append(section)
        current_heading = None
        current_level = 0
        current_chars = 0
        current_has_tables = False
        current_has_images = False
        current_para_count = 0

    for para in doc.paragraphs:
        style_name = para.style.name if para.style else ""
        # Detect heading levels
        heading_level = 0
        if style_name.startswith("Heading"):
            try:
                heading_level = int(style_name.replace("Heading", "").strip())
            except (ValueError, TypeError):
                heading_level = 0

        if heading_level > 0:
            _flush_section()
            current_heading = para.text
            current_level = heading_level
        else:
            char_count = len(para.text)
            current_chars += char_count
            total_char_count += char_count
            current_para_count += 1
            # Check for images in paragraph runs
            for run in para.runs:
                if run._element.findall('.//{http://schemas.openxmlformats.org/wordprocessingml/2006/main}drawing') or \
                   run._element.findall('.//{http://schemas.openxmlformats.org/wordprocessingml/2006/main}pict'):
                    current_has_images = True

    _flush_section()

    # Tables summary
    tables_summary: List[Dict[str, Any]] = []
    for idx, table in enumerate(doc.tables):
        rows = len(table.rows)
        cols = len(table.columns) if table.rows else 0
        tables_summary.append({
            "index": idx,
            "rows": rows,
            "cols": cols,
        })
        # Mark the current section as having tables
        # (tables in docx come after the paragraph that references them)

    # Images summary - check inline shapes
    images_summary: List[Dict[str, Any]] = []
    img_idx = 0
    for rel in doc.part.rels.values():
        if "image" in rel.reltype:
            images_summary.append({
                "index": img_idx,
                "alt_text": None,
            })
            img_idx += 1

    # Check headers/footers
    has_headers_footers = False
    for section in doc.sections:
        if section.header and section.header.paragraphs:
            for p in section.header.paragraphs:
                if p.text.strip():
                    has_headers_footers = True
                    break
        if not has_headers_footers and section.footer and section.footer.paragraphs:
            for p in section.footer.paragraphs:
                if p.text.strip():
                    has_headers_footers = True
                    break
        if has_headers_footers:
            break

    return {
        "sections": sections,
        "tables": tables_summary,
        "images": images_summary,
        "has_headers_footers": has_headers_footers,
        "total_char_count": total_char_count,
    }


def docx_get_section_content(params: Dict[str, Any]) -> Dict[str, Any]:
    """Return deterministic structured payload for one docx section."""
    docx = import_module("docx")
    path = resolve_path(params)
    doc = docx.Document(path)
    section_index = parse_non_negative_index(
        params.get("section_index", params.get("section")),
        "section_index",
        default=0,
    )
    sections = _docx_build_sections(doc)
    if section_index >= len(sections):
        raise WorkerError("VALIDATION_FAILED", "invalid section_index")
    return {
        "section_index": section_index,
        "section_count": len(sections),
        "section": sections[section_index],
    }


def _docx_build_sections(doc) -> List[Dict[str, Any]]:
    para_by_element = {id(para._element): para for para in doc.paragraphs}
    table_by_element = {id(table._element): table for table in doc.tables}
    table_indexes = {id(table._element): idx for idx, table in enumerate(doc.tables)}
    image_targets = _docx_image_targets(doc)

    sections: List[Dict[str, Any]] = []
    current: Optional[Dict[str, Any]] = None
    current_image_ids: set[str] = set()

    def _start_section(heading: Optional[str], level: int) -> Dict[str, Any]:
        return {
            "heading": heading or "(untitled)",
            "level": level,
            "paragraphs": [],
            "tables": [],
            "images": [],
        }

    def _flush() -> None:
        nonlocal current, current_image_ids
        if current is None:
            return
        sections.append(current)
        current = None
        current_image_ids = set()

    def _ensure_current() -> Dict[str, Any]:
        nonlocal current
        if current is None:
            current = _start_section("(untitled)", 0)
        return current

    for child in doc.element.body.iterchildren():
        tag = child.tag or ""
        if tag.endswith("}p"):
            para = para_by_element.get(id(child))
            if para is None:
                continue
            heading_level = _docx_heading_level(para)
            if heading_level > 0:
                _flush()
                current = _start_section(para.text, heading_level)
            section = _ensure_current()
            section["paragraphs"].append(_docx_paragraph_payload(para, len(section["paragraphs"])))
            for image_ref in _docx_paragraph_image_refs(para, image_targets):
                rel_id = image_ref.get("rel_id") or ""
                if rel_id and rel_id in current_image_ids:
                    continue
                if rel_id:
                    current_image_ids.add(rel_id)
                section["images"].append(image_ref)
        elif tag.endswith("}tbl"):
            table = table_by_element.get(id(child))
            if table is None:
                continue
            section = _ensure_current()
            table_idx = table_indexes.get(id(child), len(section["tables"]))
            section["tables"].append(_docx_table_payload(table, table_idx))
    _flush()

    if not sections:
        sections.append(
            {
                "heading": "(untitled)",
                "level": 0,
                "paragraphs": [],
                "tables": [],
                "images": [],
            }
        )
    return sections


def _docx_heading_level(para) -> int:
    style_name = para.style.name if para.style else ""
    if not style_name.startswith("Heading"):
        return 0
    suffix = style_name.replace("Heading", "").strip()
    if not suffix:
        return 0
    try:
        return int(suffix)
    except Exception:
        return 0


def _docx_paragraph_payload(para, paragraph_index: int) -> Dict[str, Any]:
    runs: List[Dict[str, Any]] = []
    for run_index, run in enumerate(para.runs):
        color = None
        size_pt = None
        font_name = None
        if run.font:
            try:
                if run.font.color and run.font.color.rgb is not None:
                    color = str(run.font.color.rgb)
            except Exception:
                color = None
            try:
                if run.font.size is not None:
                    size_pt = float(run.font.size.pt)
            except Exception:
                size_pt = None
            try:
                font_name = run.font.name
            except Exception:
                font_name = None
        runs.append(
            {
                "index": run_index,
                "text": run.text or "",
                "bold": run.bold,
                "italic": run.italic,
                "underline": run.underline,
                "size_pt": size_pt,
                "font_name": font_name,
                "color": color,
            }
        )
    return {
        "index": paragraph_index,
        "text": para.text or "",
        "style": para.style.name if para.style else None,
        "runs": runs,
    }


def _docx_table_payload(table, table_index: int) -> Dict[str, Any]:
    rows_out: List[List[str]] = []
    max_cols = 0
    for row in table.rows:
        row_vals: List[str] = []
        for cell in row.cells:
            text = cell.text or ""
            row_vals.append(text)
        if len(row_vals) > max_cols:
            max_cols = len(row_vals)
        rows_out.append(row_vals)
    return {
        "index": table_index,
        "row_count": len(rows_out),
        "col_count": max_cols,
        "rows": rows_out,
    }


def _docx_image_targets(doc) -> Dict[str, Optional[str]]:
    targets: Dict[str, Optional[str]] = {}
    for rel_id in sorted(doc.part.rels.keys()):
        rel = doc.part.rels[rel_id]
        if "image" not in rel.reltype:
            continue
        target_name = None
        try:
            if rel.target_ref:
                target_name = os.path.basename(str(rel.target_ref))
        except Exception:
            target_name = None
        targets[rel_id] = target_name
    return targets


def _docx_paragraph_image_refs(para, image_targets: Dict[str, Optional[str]]) -> List[Dict[str, Any]]:
    refs: List[Dict[str, Any]] = []
    seen: set[str] = set()
    ns_a = "{http://schemas.openxmlformats.org/drawingml/2006/main}"
    ns_r = "{http://schemas.openxmlformats.org/officeDocument/2006/relationships}"
    ns_wp = "{http://schemas.openxmlformats.org/drawingml/2006/wordprocessingDrawing}"
    for run in para.runs:
        alt_text = None
        for doc_pr in run._element.findall(f".//{ns_wp}docPr"):
            alt = (doc_pr.get("descr") or doc_pr.get("title") or "").strip()
            if alt:
                alt_text = alt
                break
        for blip in run._element.findall(f".//{ns_a}blip"):
            rel_id = blip.get(f"{ns_r}embed") or blip.get(f"{ns_r}link")
            if not rel_id or rel_id in seen:
                continue
            seen.add(rel_id)
            refs.append(
                {
                    "rel_id": rel_id,
                    "target": image_targets.get(rel_id),
                    "alt_text": alt_text,
                }
            )
    return refs


def docx_get_styles(params: Dict[str, Any]) -> Dict[str, Any]:
    docx = import_module("docx")
    path = resolve_path(params)
    doc = docx.Document(path)

    style_entries: List[Dict[str, Any]] = []
    paragraph_assets: List[Dict[str, Any]] = []
    character_assets: List[Dict[str, Any]] = []
    table_assets: List[Dict[str, Any]] = []
    for style in sorted(doc.styles, key=lambda s: (str(getattr(s, "style_id", "") or ""), str(getattr(s, "name", "") or ""))):
        name = str(getattr(style, "name", "") or "").strip()
        style_id = str(getattr(style, "style_id", "") or "").strip()
        style_type = _docx_style_type_name(style)
        if not name and not style_id:
            continue
        entry = {
            "asset_id": f"{style_type}:{name or style_id}",
            "name": name or style_id,
            "style_id": style_id or None,
            "type": style_type,
            "builtin": bool(getattr(style, "builtin", False)),
            "base_style": _docx_base_style_name(style),
        }
        style_entries.append(entry)
        if style_type == "paragraph_style":
            paragraph_assets.append(entry)
        elif style_type == "character_style":
            character_assets.append(entry)
        elif style_type == "table_style":
            table_assets.append(entry)

    image_assets: List[Dict[str, Any]] = []
    for rel_id in sorted(doc.part.rels.keys()):
        rel = doc.part.rels[rel_id]
        if "image" not in rel.reltype:
            continue
        filename = None
        try:
            if rel.target_ref:
                filename = os.path.basename(str(rel.target_ref))
        except Exception:
            filename = None
        content_type = None
        digest = None
        size_bytes = None
        try:
            if rel.target_part is not None:
                content_type = getattr(rel.target_part, "content_type", None)
                blob = rel.target_part.blob
                if blob is not None:
                    size_bytes = len(blob)
                    digest = _sha1_of_bytes(blob)
        except Exception:
            pass
        image_assets.append(
            {
                "asset_id": f"image:{rel_id}",
                "type": "image",
                "rel_id": rel_id,
                "filename": filename,
                "content_type": content_type,
                "size_bytes": size_bytes,
                "sha1": digest,
            }
        )

    return {
        "format": "docx",
        "style_count": len(style_entries),
        "styles": style_entries,
        "assets": sorted(style_entries + image_assets, key=lambda a: (str(a.get("type") or ""), str(a.get("asset_id") or ""))),
        "paragraph_styles": paragraph_assets,
        "character_styles": character_assets,
        "table_styles": table_assets,
        "image_assets": image_assets,
        "supported_copy_asset_types": ["paragraph_style", "character_style", "table_style", "image"],
    }


def docx_copy_assets(params: Dict[str, Any]) -> Dict[str, Any]:
    docx = import_module("docx")
    source_path = resolve_named_path(params, "source_path", "source_root", "draft")
    target_path = resolve_named_path(params, "target_path", "target_root", "draft", require_write=True)
    assets = _parse_asset_list(params)

    src = docx.Document(source_path)
    dst = docx.Document(target_path) if os.path.exists(target_path) else docx.Document()
    copied: List[Dict[str, Any]] = []

    for asset in assets:
        asset_type, selector = _docx_parse_asset(asset)
        if asset_type in ("paragraph_style", "character_style", "table_style", "style"):
            src_style = _docx_find_style(src, selector)
            if src_style is None:
                raise WorkerError("VALIDATION_FAILED", f"style not found: {selector}")
            dst_style = _docx_ensure_style(dst, src_style)
            _docx_copy_style_properties(src_style, dst_style)
            copied.append(
                {
                    "type": _docx_style_type_name(src_style),
                    "selector": selector,
                    "name": str(getattr(src_style, "name", "") or selector),
                }
            )
            continue
        if asset_type == "image":
            rel_id, image_part = _docx_resolve_image(src, selector)
            if image_part is None:
                raise WorkerError("VALIDATION_FAILED", f"image not found: {selector}")
            blob = image_part.blob
            if not blob:
                raise WorkerError("VALIDATION_FAILED", f"image has no data: {selector}")
            dst.add_picture(io.BytesIO(blob))
            copied.append(
                {
                    "type": "image",
                    "selector": selector,
                    "rel_id": rel_id,
                    "size_bytes": len(blob),
                    "sha1": _sha1_of_bytes(blob),
                }
            )
            continue
        raise WorkerError("VALIDATION_FAILED", f"unsupported asset type: {asset_type}")

    os.makedirs(os.path.dirname(target_path), exist_ok=True)
    dst.save(target_path)
    return {
        "ok": True,
        "copied_count": len(copied),
        "copied": copied,
    }


def _docx_style_type_name(style: Any) -> str:
    try:
        value = style.type
        text = str(value).strip().lower()
        if "paragraph" in text:
            return "paragraph_style"
        if "character" in text:
            return "character_style"
        if "table" in text:
            return "table_style"
    except Exception:
        pass
    return "style"


def _docx_base_style_name(style: Any) -> Optional[str]:
    try:
        base_style = style.base_style
        if base_style is None:
            return None
        name = str(getattr(base_style, "name", "") or "").strip()
        return name or None
    except Exception:
        return None


def _docx_parse_asset(asset: Any) -> Tuple[str, str]:
    if isinstance(asset, str):
        text = asset.strip()
        if not text:
            raise WorkerError("VALIDATION_FAILED", "asset selector is empty")
        if ":" in text:
            prefix, tail = text.split(":", 1)
            return _slug(prefix), tail.strip()
        return "style", text
    if isinstance(asset, dict):
        asset_type = _slug(asset.get("type") or "")
        selector = str(asset.get("selector") or asset.get("id") or asset.get("name") or "").strip()
        if not selector:
            selector = _asset_selector_text(asset).strip()
            if ":" in selector:
                _, tail = selector.split(":", 1)
                selector = tail.strip()
        if not selector:
            raise WorkerError("VALIDATION_FAILED", "asset selector is empty")
        if not asset_type:
            raw_id = str(asset.get("asset_id") or "").strip()
            if ":" in raw_id:
                prefix, _ = raw_id.split(":", 1)
                asset_type = _slug(prefix)
        if not asset_type:
            asset_type = "style"
        return asset_type, selector
    raise WorkerError("VALIDATION_FAILED", "asset selector must be string or object")


def _docx_find_style(doc, selector: str):
    selector_lower = selector.strip().lower()
    for style in doc.styles:
        name = str(getattr(style, "name", "") or "").strip()
        style_id = str(getattr(style, "style_id", "") or "").strip()
        if name.lower() == selector_lower or style_id.lower() == selector_lower:
            return style
    return None


def _docx_ensure_style(doc, src_style):
    target = _docx_find_style(doc, str(getattr(src_style, "name", "") or ""))
    if target is not None:
        return target
    try:
        target = doc.styles.add_style(src_style.name, src_style.type)
        return target
    except Exception:
        pass
    target = _docx_find_style(doc, str(getattr(src_style, "style_id", "") or ""))
    if target is not None:
        return target
    raise WorkerError("FILE_WRITE_FAILED", f"failed to create style: {getattr(src_style, 'name', '')}")


def _docx_copy_style_properties(src_style, dst_style) -> None:
    scalar_attrs = [
        "hidden",
        "locked",
        "priority",
        "quick_style",
        "unhide_when_used",
        "builtin",
    ]
    for attr in scalar_attrs:
        try:
            setattr(dst_style, attr, getattr(src_style, attr))
        except Exception:
            pass
    try:
        base_name = _docx_base_style_name(src_style)
        if base_name:
            pass
    except Exception:
        pass
    _copy_docx_font(src_style, dst_style)
    _copy_docx_paragraph_format(src_style, dst_style)


def _copy_docx_font(src_style, dst_style) -> None:
    try:
        src_font = src_style.font
        dst_font = dst_style.font
    except Exception:
        return
    for attr in [
        "name",
        "size",
        "bold",
        "italic",
        "underline",
        "all_caps",
        "small_caps",
        "strike",
        "subscript",
        "superscript",
    ]:
        try:
            setattr(dst_font, attr, getattr(src_font, attr))
        except Exception:
            pass
    try:
        if src_font.color is not None and src_font.color.rgb is not None:
            dst_font.color.rgb = src_font.color.rgb
    except Exception:
        pass


def _copy_docx_paragraph_format(src_style, dst_style) -> None:
    try:
        src_pf = src_style.paragraph_format
        dst_pf = dst_style.paragraph_format
    except Exception:
        return
    for attr in [
        "alignment",
        "left_indent",
        "right_indent",
        "first_line_indent",
        "space_before",
        "space_after",
        "line_spacing",
        "line_spacing_rule",
        "keep_together",
        "keep_with_next",
        "page_break_before",
        "widow_control",
    ]:
        try:
            setattr(dst_pf, attr, getattr(src_pf, attr))
        except Exception:
            pass


def _docx_resolve_image(doc, selector: str) -> Tuple[Optional[str], Any]:
    selector = selector.strip()
    if not selector:
        return None, None
    rel_id = selector
    if selector.isdigit():
        index = int(selector)
        image_rels = []
        for rid in sorted(doc.part.rels.keys()):
            rel = doc.part.rels[rid]
            if "image" in rel.reltype:
                image_rels.append((rid, rel))
        if index < 0 or index >= len(image_rels):
            return None, None
        rel_id, rel = image_rels[index]
        return rel_id, rel.target_part
    if not rel_id.startswith("rId") and rel_id.lower().startswith("image:"):
        rel_id = rel_id.split(":", 1)[1].strip()
    rel = doc.part.rels.get(rel_id)
    if rel is None or "image" not in rel.reltype:
        return None, None
    return rel_id, rel.target_part


# ------------------ ODT ------------------

def odt_extract_text(params: Dict[str, Any]) -> Dict[str, Any]:
    odf = import_module("odf")
    from odf import text as odf_text
    from odf import teletype
    path = resolve_path(params)
    doc = odf.opendocument.load(path)
    lines = []
    for para in doc.getElementsByType(odf_text.P):
        content = teletype.extractText(para)
        if content:
            lines.append(content)
    return {"text": "\n".join(lines)}


def odt_render_page(params: Dict[str, Any]) -> Dict[str, Any]:
    path = resolve_path(params)
    page_index = int(params.get("page_index", 0))
    scale = clamp_scale(params.get("scale", 1.0))
    pdf_path, tmp_dir = convert_to_pdf(path)
    try:
        return render_pdf_page(pdf_path, page_index, scale)
    finally:
        shutil.rmtree(tmp_dir, ignore_errors=True)


# ------------------ XLSX ------------------

def _xlsx_sanitize_value(value):
    """Strip XML-illegal control characters from string cell values."""
    if isinstance(value, str):
        from openpyxl.cell.cell import ILLEGAL_CHARACTERS_RE
        return ILLEGAL_CHARACTERS_RE.sub("", value)
    return value


def xlsx_apply_ops(params: Dict[str, Any]) -> Dict[str, Any]:
    openpyxl = import_module("openpyxl")
    from openpyxl.utils import get_column_letter
    from openpyxl.utils.cell import coordinate_from_string, column_index_from_string
    from openpyxl.cell.cell import MergedCell

    path = resolve_path(params)
    ops = params.get("ops") or []
    create_new = params.get("create_new", False)
    copy_from = params.get("copy_from")

    if not isinstance(ops, list) or len(ops) == 0:
        raise WorkerError("VALIDATION_FAILED", "missing ops")

    # Handle create_new with optional copy_from
    if create_new:
        if copy_from:
            # Copy from source file - resolve the source path
            source_params = dict(params)
            source_params["path"] = copy_from
            source_path = resolve_path(source_params)
            if not os.path.exists(source_path):
                raise WorkerError("FILE_READ_FAILED", f"source file not found: {copy_from}")
            wb = openpyxl.load_workbook(source_path)
            log_info("xlsx.copy_from", source=copy_from, target=path)
        else:
            wb = openpyxl.Workbook()
            log_info("xlsx.create_new", target=path)
    elif os.path.exists(path):
        wb = openpyxl.load_workbook(path)
    else:
        wb = openpyxl.Workbook()

    def find_merged_anchor(ws, row: int, col: int) -> Optional[Tuple[int, int, str]]:
        for cell_range in ws.merged_cells.ranges:
            if cell_range.min_row <= row <= cell_range.max_row and cell_range.min_col <= col <= cell_range.max_col:
                return cell_range.min_row, cell_range.min_col, str(cell_range)
        return None

    def set_cell_value(ws, row: int, col: int, value: Any, seen_merges: Optional[set]) -> Optional[Any]:
        cell = ws.cell(row=row, column=col)
        if isinstance(cell, MergedCell):
            anchor = find_merged_anchor(ws, row, col)
            if not anchor:
                return None
            key = anchor[2]
            if seen_merges is not None and key in seen_merges:
                return None
            anchor_cell = ws.cell(row=anchor[0], column=anchor[1], value=_xlsx_sanitize_value(value))
            if seen_merges is not None:
                seen_merges.add(key)
            return anchor_cell
        cell.value = _xlsx_sanitize_value(value)
        return cell

    def parse_column_index(raw: Any, default: str) -> int:
        value = default if raw is None else raw
        if isinstance(value, int):
            if value < 1:
                raise WorkerError("VALIDATION_FAILED", "column index must be >= 1")
            return value
        text = str(value).strip()
        if not text:
            text = default
        try:
            return column_index_from_string(text.upper())
        except Exception as exc:
            raise WorkerError("VALIDATION_FAILED", f"invalid column reference: {value}") from exc

    def row_value(row_values: Tuple[Any, ...], col_index_1_based: int) -> Any:
        idx = col_index_1_based - 1
        if idx < 0 or idx >= len(row_values):
            return None
        return row_values[idx]

    for op in ops:
        if not isinstance(op, dict):
            raise WorkerError("VALIDATION_FAILED", "invalid op")
        name = op.get("op")
        if name == "ensure_sheet":
            sheet_name = op.get("sheet")
            if not sheet_name:
                raise WorkerError("VALIDATION_FAILED", "missing sheet")
            if sheet_name not in wb.sheetnames:
                wb.create_sheet(sheet_name)
        elif name == "set_cells":
            sheet_name = op.get("sheet")
            if not sheet_name:
                raise WorkerError("VALIDATION_FAILED", "missing sheet")
            if sheet_name not in wb.sheetnames:
                wb.create_sheet(sheet_name)
            ws = wb[sheet_name]
            seen_merges: set = set()
            for cell in op.get("cells") or []:
                ref = cell.get("cell") if isinstance(cell, dict) else None
                if not ref:
                    continue
                value = cell.get("value") if isinstance(cell, dict) else None
                col_letter, row_idx = coordinate_from_string(ref)
                base_col = column_index_from_string(col_letter)
                target_cell = set_cell_value(ws, row_idx, base_col, value, seen_merges)
                if target_cell is not None and isinstance(cell, dict) and "style" in cell:
                    _xlsx_apply_inline_style(
                        target_cell,
                        cell.get("style"),
                        op_name=name,
                        sheet=sheet_name,
                        cell_ref=str(target_cell.coordinate),
                    )
        elif name == "set_range":
            sheet_name = op.get("sheet")
            start = op.get("start")
            values = op.get("values") or []
            style = op.get("style")
            if not sheet_name or not start:
                raise WorkerError("VALIDATION_FAILED", "missing sheet/start")
            if sheet_name not in wb.sheetnames:
                wb.create_sheet(sheet_name)
            ws = wb[sheet_name]
            col_letter, row_idx = coordinate_from_string(start)
            base_col = column_index_from_string(col_letter)
            base_row = row_idx
            seen_merges: set = set()
            for r, row in enumerate(values):
                for c, cell_value in enumerate(row):
                    target_cell = set_cell_value(ws, base_row + r, base_col + c, cell_value, seen_merges)
                    if target_cell is not None and style is not None:
                        _xlsx_apply_inline_style(
                            target_cell,
                            style,
                            op_name=name,
                            sheet=sheet_name,
                            cell_ref=str(target_cell.coordinate),
                        )
        elif name == "set_column_widths":
            sheet_name = op.get("sheet")
            columns = op.get("columns")
            if not sheet_name:
                raise WorkerError("VALIDATION_FAILED", "missing sheet")
            if not isinstance(columns, list):
                raise WorkerError("VALIDATION_FAILED", "columns must be list")
            if sheet_name not in wb.sheetnames:
                wb.create_sheet(sheet_name)
            ws = wb[sheet_name]
            for idx, entry in enumerate(columns):
                if not isinstance(entry, dict):
                    _style_warn(
                        "xlsx.column_widths",
                        "columns",
                        entry,
                        "columns entry must be object",
                        op=name,
                        sheet=sheet_name,
                        index=idx,
                    )
                    continue
                column_text = _xlsx_parse_column_letter(entry.get("column"), get_column_letter, column_index_from_string)
                if column_text is None:
                    _style_warn(
                        "xlsx.column_widths",
                        "column",
                        entry.get("column"),
                        "invalid column reference",
                        op=name,
                        sheet=sheet_name,
                        index=idx,
                    )
                    continue
                width = _coerce_float(entry.get("width"), None)
                if width is None or width < 0:
                    _style_warn(
                        "xlsx.column_widths",
                        "width",
                        entry.get("width"),
                        "width must be a non-negative number",
                        op=name,
                        sheet=sheet_name,
                        index=idx,
                        column=column_text,
                    )
                    continue
                ws.column_dimensions[column_text].width = width
        elif name == "set_row_heights":
            sheet_name = op.get("sheet")
            rows = op.get("rows")
            if not sheet_name:
                raise WorkerError("VALIDATION_FAILED", "missing sheet")
            if not isinstance(rows, list):
                raise WorkerError("VALIDATION_FAILED", "rows must be list")
            if sheet_name not in wb.sheetnames:
                wb.create_sheet(sheet_name)
            ws = wb[sheet_name]
            for idx, entry in enumerate(rows):
                if not isinstance(entry, dict):
                    _style_warn(
                        "xlsx.row_heights",
                        "rows",
                        entry,
                        "rows entry must be object",
                        op=name,
                        sheet=sheet_name,
                        index=idx,
                    )
                    continue
                row_index = _xlsx_parse_row_index(entry.get("row"))
                if row_index is None:
                    _style_warn(
                        "xlsx.row_heights",
                        "row",
                        entry.get("row"),
                        "row must be a positive integer",
                        op=name,
                        sheet=sheet_name,
                        index=idx,
                    )
                    continue
                height = _coerce_float(entry.get("height"), None)
                if height is None or height < 0:
                    _style_warn(
                        "xlsx.row_heights",
                        "height",
                        entry.get("height"),
                        "height must be a non-negative number",
                        op=name,
                        sheet=sheet_name,
                        index=idx,
                        row=row_index,
                    )
                    continue
                ws.row_dimensions[row_index].height = height
        elif name == "freeze_panes":
            sheet_name = op.get("sheet")
            if not sheet_name:
                raise WorkerError("VALIDATION_FAILED", "missing sheet")
            if sheet_name not in wb.sheetnames:
                wb.create_sheet(sheet_name)
            ws = wb[sheet_name]
            row = _coerce_int(op.get("row", 0), None)
            column = _coerce_int(op.get("column", 0), None)
            if row is None or row < 0:
                _style_warn("xlsx.freeze_panes", "row", op.get("row"), "row must be a non-negative integer", op=name, sheet=sheet_name)
                continue
            if column is None or column < 0:
                _style_warn("xlsx.freeze_panes", "column", op.get("column"), "column must be a non-negative integer", op=name, sheet=sheet_name)
                continue
            if row == 0 and column == 0:
                ws.freeze_panes = None
            else:
                ws.freeze_panes = f"{get_column_letter(column + 1)}{row + 1}"
        elif name == "summarize_by_category":
            target_sheet_name = op.get("sheet")
            source_sheets = op.get("source_sheets")
            if not target_sheet_name:
                raise WorkerError("VALIDATION_FAILED", "missing sheet")
            if not isinstance(source_sheets, list) or len(source_sheets) == 0:
                raise WorkerError("VALIDATION_FAILED", "missing source_sheets")

            normalized_sources: List[str] = []
            for source_name in source_sheets:
                if not isinstance(source_name, str) or not source_name.strip():
                    raise WorkerError("VALIDATION_FAILED", "invalid source_sheets entry")
                cleaned_name = source_name.strip()
                if cleaned_name not in wb.sheetnames:
                    raise WorkerError("VALIDATION_FAILED", f"unknown sheet: {cleaned_name}")
                normalized_sources.append(cleaned_name)

            category_col_idx = parse_column_index(op.get("category_col"), "B")
            amount_col_idx = parse_column_index(op.get("amount_col"), "C")
            category_header = str(op.get("category_header") or "Category")

            if target_sheet_name not in wb.sheetnames:
                wb.create_sheet(target_sheet_name)
            ws = wb[target_sheet_name]
            if ws.max_row > 0:
                ws.delete_rows(1, ws.max_row)

            totals_by_sheet: Dict[str, Dict[str, float]] = {}
            categories: Set[str] = set()
            for source_name in normalized_sources:
                source_ws = wb[source_name]
                sheet_totals: Dict[str, float] = {}
                for row in source_ws.iter_rows(min_row=2, values_only=True):
                    category_val = row_value(row, category_col_idx)
                    amount_val = row_value(row, amount_col_idx)
                    if category_val is None:
                        continue
                    category = str(category_val).strip()
                    if not category or not isinstance(amount_val, (int, float)):
                        continue
                    sheet_totals[category] = sheet_totals.get(category, 0.0) + float(amount_val)
                    categories.add(category)
                totals_by_sheet[source_name] = sheet_totals

            headers = [category_header] + [f"{sheet_name}_Total" for sheet_name in normalized_sources] + ["Grand_Total"]
            ws.append([_xlsx_sanitize_value(h) for h in headers])

            for category in sorted(categories):
                row_cells: List[Any] = [category]
                grand_total = 0.0
                for source_name in normalized_sources:
                    subtotal = round(totals_by_sheet[source_name].get(category, 0.0), 2)
                    row_cells.append(subtotal)
                    grand_total += subtotal
                row_cells.append(round(grand_total, 2))
                ws.append([_xlsx_sanitize_value(v) for v in row_cells])
        else:
            raise WorkerError("VALIDATION_FAILED", f"unsupported op: {name}")
    os.makedirs(os.path.dirname(path), exist_ok=True)
    wb.save(path)
    return {"ok": True}


def _xlsx_parse_column_letter(value: Any, get_column_letter: Any, column_index_from_string: Any) -> Optional[str]:
    if value is None:
        return None
    if isinstance(value, int):
        if value < 1:
            return None
        return get_column_letter(value)
    text = str(value).strip()
    if not text:
        return None
    if text.isdigit():
        idx = int(text)
        if idx < 1:
            return None
        return get_column_letter(idx)
    try:
        idx = column_index_from_string(text.upper())
    except Exception:
        return None
    return get_column_letter(idx)


def _xlsx_parse_row_index(value: Any) -> Optional[int]:
    row = _coerce_int(value, None)
    if row is None or row < 1:
        return None
    return row


def _xlsx_apply_inline_style(cell: Any, style: Any, op_name: str, sheet: str, cell_ref: str) -> None:
    from openpyxl.styles import Side

    context = {"op": op_name, "sheet": sheet, "cell": cell_ref}
    if style is None:
        return
    if not isinstance(style, dict):
        _style_warn("xlsx.style", "style", style, "style must be object", **context)
        return

    known_keys = {
        "font_name",
        "font_size",
        "font_bold",
        "font_italic",
        "font_color",
        "fill_color",
        "fill_pattern",
        "number_format",
        "h_align",
        "v_align",
        "wrap_text",
        "border_top",
        "border_bottom",
        "border_left",
        "border_right",
    }
    for key in style.keys():
        if key not in known_keys:
            _style_warn("xlsx.style", key, style.get(key), "unknown style key", **context)

    font = pycopy.copy(cell.font)
    font_changed = False
    font_name = style.get("font_name")
    if font_name is not None:
        if isinstance(font_name, str) and font_name.strip():
            font.name = font_name.strip()
            font_changed = True
        else:
            _style_warn("xlsx.style", "font_name", font_name, "font_name must be non-empty string", **context)

    font_size = style.get("font_size")
    if font_size is not None:
        parsed = _coerce_float(font_size, None)
        if parsed is None or parsed <= 0:
            _style_warn("xlsx.style", "font_size", font_size, "font_size must be a positive number", **context)
        else:
            font.size = parsed
            font_changed = True

    for field, attr in (("font_bold", "bold"), ("font_italic", "italic")):
        value = style.get(field)
        if value is None:
            continue
        parsed = _coerce_bool_strict(value, None)
        if parsed is None:
            _style_warn("xlsx.style", field, value, f"{field} must be boolean", **context)
            continue
        setattr(font, attr, parsed)
        font_changed = True

    font_color = style.get("font_color")
    if font_color is not None:
        rgb = _normalize_hex_rgb(font_color)
        if rgb is None:
            _style_warn("xlsx.style", "font_color", font_color, "font_color must be #RRGGBB", **context)
        else:
            font.color = f"FF{rgb}"
            font_changed = True

    if font_changed:
        cell.font = font

    fill = pycopy.copy(cell.fill)
    fill_changed = False
    fill_pattern_map = {
        "none": None,
        "solid": "solid",
        "darkdown": "darkDown",
        "darkgray": "darkGray",
        "darkgrid": "darkGrid",
        "darkhorizontal": "darkHorizontal",
        "darktrellis": "darkTrellis",
        "darkup": "darkUp",
        "darkvertical": "darkVertical",
        "gray0625": "gray0625",
        "gray125": "gray125",
        "lightdown": "lightDown",
        "lightgray": "lightGray",
        "lightgrid": "lightGrid",
        "lighthorizontal": "lightHorizontal",
        "lighttrellis": "lightTrellis",
        "lightup": "lightUp",
        "lightvertical": "lightVertical",
        "mediumgray": "mediumGray",
    }
    fill_pattern = style.get("fill_pattern")
    if fill_pattern is not None:
        if not isinstance(fill_pattern, str):
            _style_warn("xlsx.style", "fill_pattern", fill_pattern, "fill_pattern must be string", **context)
        else:
            parsed_pattern = fill_pattern_map.get(_slug(fill_pattern))
            if _slug(fill_pattern) not in fill_pattern_map:
                _style_warn("xlsx.style", "fill_pattern", fill_pattern, "unsupported fill_pattern", **context)
            else:
                fill.patternType = parsed_pattern
                fill_changed = True

    fill_color = style.get("fill_color")
    if fill_color is not None:
        rgb = _normalize_hex_rgb(fill_color)
        if rgb is None:
            _style_warn("xlsx.style", "fill_color", fill_color, "fill_color must be #RRGGBB", **context)
        else:
            fill.fgColor = f"FF{rgb}"
            if fill_pattern is None:
                pattern_type = str(getattr(fill, "patternType", "") or "").strip()
                if not pattern_type:
                    fill.patternType = "solid"
            fill_changed = True

    if fill_changed:
        cell.fill = fill

    if "number_format" in style:
        number_format = style.get("number_format")
        if isinstance(number_format, str):
            cell.number_format = number_format
        else:
            _style_warn("xlsx.style", "number_format", number_format, "number_format must be string", **context)

    alignment = pycopy.copy(cell.alignment)
    alignment_changed = False
    h_align = style.get("h_align")
    if h_align is not None:
        if not isinstance(h_align, str):
            _style_warn("xlsx.style", "h_align", h_align, "h_align must be string", **context)
        else:
            h_align_map = {
                "left": "left",
                "center": "center",
                "right": "right",
                "justify": "justify",
                "general": "general",
                "fill": "fill",
                "distributed": "distributed",
                "centercontinuous": "centerContinuous",
            }
            parsed = h_align_map.get(_slug(h_align))
            if parsed is None:
                _style_warn("xlsx.style", "h_align", h_align, "unsupported h_align", **context)
            else:
                alignment.horizontal = parsed
                alignment_changed = True

    v_align = style.get("v_align")
    if v_align is not None:
        if not isinstance(v_align, str):
            _style_warn("xlsx.style", "v_align", v_align, "v_align must be string", **context)
        else:
            v_align_map = {
                "top": "top",
                "center": "center",
                "bottom": "bottom",
                "justify": "justify",
                "distributed": "distributed",
            }
            parsed = v_align_map.get(_slug(v_align))
            if parsed is None:
                _style_warn("xlsx.style", "v_align", v_align, "unsupported v_align", **context)
            else:
                alignment.vertical = parsed
                alignment_changed = True

    wrap_text = style.get("wrap_text")
    if wrap_text is not None:
        parsed = _coerce_bool_strict(wrap_text, None)
        if parsed is None:
            _style_warn("xlsx.style", "wrap_text", wrap_text, "wrap_text must be boolean", **context)
        else:
            alignment.wrap_text = parsed
            alignment_changed = True

    if alignment_changed:
        cell.alignment = alignment

    border = pycopy.copy(cell.border)
    border_changed = False
    border_style_map = {
        "none": None,
        "thin": "thin",
        "medium": "medium",
        "thick": "thick",
        "dashed": "dashed",
        "dotted": "dotted",
        "double": "double",
    }
    for field, side_attr in (
        ("border_top", "top"),
        ("border_bottom", "bottom"),
        ("border_left", "left"),
        ("border_right", "right"),
    ):
        if field not in style:
            continue
        value = style.get(field)
        if not isinstance(value, dict):
            _style_warn("xlsx.style", field, value, f"{field} must be object", **context)
            continue
        for key in value.keys():
            if key not in ("style", "color"):
                _style_warn("xlsx.style", f"{field}.{key}", value.get(key), "unknown style key", **context)

        side_style = value.get("style")
        side_color = value.get("color")
        parsed_side_style = None
        has_side_style = False
        if side_style is not None:
            if not isinstance(side_style, str):
                _style_warn("xlsx.style", f"{field}.style", side_style, "border style must be string", **context)
            else:
                key = _slug(side_style)
                if key not in border_style_map:
                    _style_warn(
                        "xlsx.style",
                        f"{field}.style",
                        side_style,
                        "unsupported border style",
                        **context,
                    )
                else:
                    parsed_side_style = border_style_map[key]
                    has_side_style = True

        parsed_side_color = None
        has_side_color = False
        if side_color is not None:
            rgb = _normalize_hex_rgb(side_color)
            if rgb is None:
                _style_warn("xlsx.style", f"{field}.color", side_color, "border color must be #RRGGBB", **context)
            else:
                parsed_side_color = f"FF{rgb}"
                has_side_color = True

        if not has_side_style and not has_side_color:
            continue
        try:
            setattr(border, side_attr, Side(style=parsed_side_style, color=parsed_side_color))
            border_changed = True
        except Exception as err:
            _style_warn("xlsx.style", field, value, "failed to apply border side", error=str(err), **context)

    if border_changed:
        cell.border = border


def xlsx_get_styles(params: Dict[str, Any]) -> Dict[str, Any]:
    openpyxl = import_module("openpyxl")
    path = resolve_path(params)
    wb = openpyxl.load_workbook(path)
    sheet_name = str(params.get("sheet") or "").strip()
    ws = None
    if sheet_name:
        if sheet_name not in wb.sheetnames:
            raise WorkerError("VALIDATION_FAILED", f"unknown sheet: {sheet_name}")
        ws = wb[sheet_name]

    named_styles: List[Dict[str, Any]] = []
    for style in wb.named_styles:
        if isinstance(style, str):
            name = style
            builtin = None
            style_id = None
        else:
            name = getattr(style, "name", "") or ""
            builtin = getattr(style, "builtinId", None) is not None
            style_id = getattr(style, "style_id", None)
        if not name:
            continue
        named_styles.append(
            {
                "asset_id": f"named_style:{name}",
                "type": "named_style",
                "name": name,
                "builtin": builtin,
                "style_id": style_id,
            }
        )
    named_styles.sort(key=lambda item: str(item.get("name") or "").lower())

    fonts: List[Dict[str, Any]] = []
    for idx, font in enumerate(getattr(wb, "_fonts", [])):
        fonts.append(
            {
                "index": idx,
                "name": getattr(font, "name", None),
                "size": float(font.sz) if getattr(font, "sz", None) is not None else None,
                "bold": bool(getattr(font, "b", False)),
                "italic": bool(getattr(font, "i", False)),
                "underline": str(getattr(font, "u", "") or "") or None,
                "color": str(getattr(getattr(font, "color", None), "rgb", "") or "") or None,
            }
        )

    fills: List[Dict[str, Any]] = []
    for idx, fill in enumerate(getattr(wb, "_fills", [])):
        fg = None
        bg = None
        try:
            fg = str(getattr(fill.fgColor, "rgb", "") or "") or None
            bg = str(getattr(fill.bgColor, "rgb", "") or "") or None
        except Exception:
            pass
        fills.append(
            {
                "index": idx,
                "pattern": str(getattr(fill, "patternType", "") or "") or None,
                "fg_color": fg,
                "bg_color": bg,
            }
        )

    borders: List[Dict[str, Any]] = []
    for idx, border in enumerate(getattr(wb, "_borders", [])):
        borders.append(
            {
                "index": idx,
                "left": str(getattr(getattr(border, "left", None), "style", "") or "") or None,
                "right": str(getattr(getattr(border, "right", None), "style", "") or "") or None,
                "top": str(getattr(getattr(border, "top", None), "style", "") or "") or None,
                "bottom": str(getattr(getattr(border, "bottom", None), "style", "") or "") or None,
            }
        )

    number_formats: List[Dict[str, Any]] = []
    for idx, code in enumerate(getattr(wb, "_number_formats", [])):
        if not code:
            continue
        number_formats.append(
            {
                "asset_id": f"number_format:{code}",
                "type": "number_format",
                "index": idx,
                "code": str(code),
            }
        )
    number_formats.sort(key=lambda item: item.get("code", ""))

    cell_style_assets: List[Dict[str, Any]] = []
    if ws is not None:
        for row in ws.iter_rows(min_row=ws.min_row or 1, max_row=ws.max_row or 0, min_col=ws.min_column or 1, max_col=ws.max_column or 0):
            for cell in row:
                if cell.style_id is None or int(cell.style_id) == 0:
                    continue
                selector = f"{ws.title}!{cell.coordinate}"
                cell_style_assets.append(
                    {
                        "asset_id": f"cell_style:{selector}",
                        "type": "cell_style",
                        "sheet": ws.title,
                        "cell": cell.coordinate,
                        "style_id": int(cell.style_id),
                        "named_style": str(cell.style or "") or None,
                        "number_format": str(cell.number_format or "") or None,
                    }
                )
        cell_style_assets.sort(key=lambda item: (str(item.get("sheet") or "").lower(), str(item.get("cell") or "")))

    all_assets = sorted(
        named_styles + number_formats + cell_style_assets,
        key=lambda item: (str(item.get("type") or ""), str(item.get("asset_id") or "")),
    )
    result: Dict[str, Any] = {
        "format": "xlsx",
        "sheet_count": len(wb.sheetnames),
        "sheets": list(wb.sheetnames),
        "style_count": len(named_styles),
        "named_styles": named_styles,
        "number_formats": number_formats,
        "fonts": fonts,
        "fills": fills,
        "borders": borders,
        "assets": all_assets,
        "supported_copy_asset_types": ["named_style", "cell_style", "number_format"],
    }
    if ws is not None:
        result["sheet"] = ws.title
        result["cell_style_assets"] = cell_style_assets
    return result


def xlsx_copy_assets(params: Dict[str, Any]) -> Dict[str, Any]:
    openpyxl = import_module("openpyxl")
    from openpyxl.utils.cell import range_boundaries

    source_path = resolve_named_path(params, "source_path", "source_root", "draft")
    target_path = resolve_named_path(params, "target_path", "target_root", "draft", require_write=True)
    assets = _parse_asset_list(params)

    src_wb = openpyxl.load_workbook(source_path)
    dst_wb = openpyxl.load_workbook(target_path) if os.path.exists(target_path) else openpyxl.Workbook()

    copied: List[Dict[str, Any]] = []
    normalized_assets = [_xlsx_parse_asset(asset) for asset in assets]
    for item in normalized_assets:
        kind = item["type"]
        if kind == "named_style":
            style_name = item["name"]
            _xlsx_copy_named_style(src_wb, dst_wb, style_name)
            copied.append({"type": "named_style", "name": style_name})
            continue
        if kind == "number_format":
            code = item["number_format"]
            targets = _xlsx_resolve_target_cells(dst_wb, item, range_boundaries)
            if len(targets) == 0:
                raise WorkerError("VALIDATION_FAILED", "number_format asset requires target cell or range")
            for cell in targets:
                cell.number_format = code
            copied.append({"type": "number_format", "code": code, "target_count": len(targets)})
            continue
        if kind == "cell_style":
            src_sheet = item["source_sheet"]
            src_cell = item["source_cell"]
            if src_sheet not in src_wb.sheetnames:
                raise WorkerError("VALIDATION_FAILED", f"unknown source sheet: {src_sheet}")
            src_ws = src_wb[src_sheet]
            source_cell = src_ws[src_cell]
            targets = _xlsx_resolve_target_cells(dst_wb, item, range_boundaries)
            if len(targets) == 0:
                raise WorkerError("VALIDATION_FAILED", "cell_style asset requires target cell or range")
            _xlsx_copy_style_to_cells(source_cell, targets, src_wb, dst_wb)
            copied.append(
                {
                    "type": "cell_style",
                    "source": f"{src_sheet}!{src_cell}",
                    "target_count": len(targets),
                }
            )
            continue
        raise WorkerError("VALIDATION_FAILED", f"unsupported asset type: {kind}")

    os.makedirs(os.path.dirname(target_path), exist_ok=True)
    dst_wb.save(target_path)
    return {"ok": True, "copied_count": len(copied), "copied": copied}


def _xlsx_parse_asset(asset: Any) -> Dict[str, Any]:
    if isinstance(asset, str):
        text = asset.strip()
        if not text:
            raise WorkerError("VALIDATION_FAILED", "asset selector is empty")
        if ":" not in text:
            return {"type": "named_style", "name": text}
        prefix, tail = text.split(":", 1)
        kind = _slug(prefix)
        payload = tail.strip()
        if kind == "named_style":
            return {"type": "named_style", "name": payload}
        if kind == "number_format":
            return {"type": "number_format", "number_format": payload}
        if kind == "cell_style":
            # Sheet1!A1->Sheet2!B2 or Sheet1!A1->Sheet2!B2:C3
            source_part, _, target_part = payload.partition("->")
            src_sheet, src_cell = _xlsx_split_sheet_cell(source_part)
            target_sheet, target_ref = _xlsx_split_sheet_cell(target_part) if target_part else ("", "")
            if not target_sheet or not target_ref:
                raise WorkerError("VALIDATION_FAILED", "cell_style selector requires target mapping")
            out = {"type": "cell_style", "source_sheet": src_sheet, "source_cell": src_cell, "target_sheet": target_sheet}
            if ":" in target_ref:
                out["target_range"] = target_ref
            else:
                out["target_cell"] = target_ref
            return out
        raise WorkerError("VALIDATION_FAILED", f"unsupported asset type: {kind}")
    if not isinstance(asset, dict):
        raise WorkerError("VALIDATION_FAILED", "asset selector must be string or object")

    kind = _slug(asset.get("type") or "")
    if not kind:
        raw = str(asset.get("asset_id") or "").strip()
        if ":" in raw:
            kind = _slug(raw.split(":", 1)[0])
    if not kind:
        kind = "named_style"

    if kind == "named_style":
        name = str(asset.get("name") or asset.get("id") or "").strip()
        if not name:
            raise WorkerError("VALIDATION_FAILED", "named_style asset missing name")
        return {"type": "named_style", "name": name}
    if kind == "number_format":
        code = str(asset.get("number_format") or asset.get("code") or asset.get("name") or "").strip()
        if not code:
            raise WorkerError("VALIDATION_FAILED", "number_format asset missing code")
        out = {"type": "number_format", "number_format": code}
        target_sheet = str(asset.get("target_sheet") or "").strip()
        target_cell = str(asset.get("target_cell") or "").strip()
        target_range = str(asset.get("target_range") or "").strip()
        if target_sheet:
            out["target_sheet"] = target_sheet
        if target_cell:
            out["target_cell"] = target_cell
        if target_range:
            out["target_range"] = target_range
        return out
    if kind == "cell_style":
        source_sheet = str(asset.get("source_sheet") or "").strip()
        source_cell = str(asset.get("source_cell") or "").strip()
        target_sheet = str(asset.get("target_sheet") or "").strip()
        target_cell = str(asset.get("target_cell") or "").strip()
        target_range = str(asset.get("target_range") or "").strip()
        if not source_sheet or not source_cell:
            source = str(asset.get("source") or "").strip()
            if "!" in source:
                source_sheet, source_cell = _xlsx_split_sheet_cell(source)
        if not target_sheet or (not target_cell and not target_range):
            target = str(asset.get("target") or "").strip()
            if "!" in target:
                target_sheet, target_ref = _xlsx_split_sheet_cell(target)
                if ":" in target_ref:
                    target_range = target_ref
                else:
                    target_cell = target_ref
        if not source_sheet or not source_cell:
            raise WorkerError("VALIDATION_FAILED", "cell_style asset missing source_sheet/source_cell")
        if not target_sheet or (not target_cell and not target_range):
            raise WorkerError("VALIDATION_FAILED", "cell_style asset missing target_sheet/target_cell or target_range")
        out = {"type": "cell_style", "source_sheet": source_sheet, "source_cell": source_cell, "target_sheet": target_sheet}
        if target_range:
            out["target_range"] = target_range
        else:
            out["target_cell"] = target_cell
        return out
    raise WorkerError("VALIDATION_FAILED", f"unsupported asset type: {kind}")


def _xlsx_split_sheet_cell(value: str) -> Tuple[str, str]:
    text = value.strip()
    if "!" not in text:
        raise WorkerError("VALIDATION_FAILED", f"invalid sheet/cell selector: {value}")
    sheet, cell = text.split("!", 1)
    sheet = sheet.strip().strip("'")
    cell = cell.strip()
    if not sheet or not cell:
        raise WorkerError("VALIDATION_FAILED", f"invalid sheet/cell selector: {value}")
    return sheet, cell


def _xlsx_copy_named_style(src_wb, dst_wb, style_name: str) -> None:
    source_style = None
    for style in src_wb.named_styles:
        if isinstance(style, str):
            if style == style_name:
                source_style = style
                break
            continue
        if str(getattr(style, "name", "") or "") == style_name:
            source_style = style
            break
    if source_style is None:
        raise WorkerError("VALIDATION_FAILED", f"named style not found: {style_name}")

    for existing in dst_wb.named_styles:
        if isinstance(existing, str):
            if existing == style_name:
                return
            continue
        if str(getattr(existing, "name", "") or "") == style_name:
            return

    if isinstance(source_style, str):
        # Built-in style name; applying by name is enough.
        return
    try:
        dst_wb.add_named_style(pycopy.copy(source_style))
    except Exception:
        # If style already exists by id/name, treat as copied.
        pass


def _xlsx_resolve_target_cells(dst_wb, item: Dict[str, Any], range_boundaries) -> List[Any]:
    target_sheet = str(item.get("target_sheet") or "").strip()
    if not target_sheet:
        raise WorkerError("VALIDATION_FAILED", "missing target_sheet")
    if target_sheet not in dst_wb.sheetnames:
        dst_wb.create_sheet(target_sheet)
    ws = dst_wb[target_sheet]

    target_range = str(item.get("target_range") or "").strip()
    target_cell = str(item.get("target_cell") or "").strip()
    if target_range:
        min_col, min_row, max_col, max_row = range_boundaries(target_range)
        out = []
        for r in range(min_row, max_row + 1):
            for c in range(min_col, max_col + 1):
                out.append(ws.cell(row=r, column=c))
        return out
    if target_cell:
        return [ws[target_cell]]
    return []


def _xlsx_copy_style_to_cells(source_cell: Any, targets: List[Any], src_wb: Any, dst_wb: Any) -> None:
    try:
        style_name = str(source_cell.style or "").strip()
    except Exception:
        style_name = ""
    if style_name:
        _xlsx_copy_named_style(src_wb, dst_wb, style_name)
    for target in targets:
        try:
            target._style = pycopy.copy(source_cell._style)
        except Exception:
            pass
        try:
            target.number_format = source_cell.number_format
        except Exception:
            pass
        try:
            target.font = pycopy.copy(source_cell.font)
        except Exception:
            pass
        try:
            target.fill = pycopy.copy(source_cell.fill)
        except Exception:
            pass
        try:
            target.border = pycopy.copy(source_cell.border)
        except Exception:
            pass
        try:
            target.alignment = pycopy.copy(source_cell.alignment)
        except Exception:
            pass
        try:
            target.protection = pycopy.copy(source_cell.protection)
        except Exception:
            pass


def xlsx_extract_text(params: Dict[str, Any]) -> Dict[str, Any]:
    openpyxl = import_module("openpyxl")
    path = resolve_path(params)
    wb = openpyxl.load_workbook(path, data_only=False)
    lines = []
    for sheet in wb.worksheets:
        lines.append(f"# Sheet: {sheet.title}")
        for row in sheet.iter_rows(values_only=True):
            parts = []
            for val in row:
                if val is None:
                    parts.append("")
                else:
                    parts.append(str(val))
            lines.append("\t".join(parts))
    return {"text": "\n".join(lines)}


def xlsx_get_info(params: Dict[str, Any]) -> Dict[str, Any]:
    """Get information about an xlsx file including sheet names and dimensions."""
    openpyxl = import_module("openpyxl")
    path = resolve_path(params)
    wb = openpyxl.load_workbook(path, data_only=False)
    sheets = []
    for sheet in wb.worksheets:
        max_row = sheet.max_row or 0
        max_col = sheet.max_column or 0
        sheets.append({
            "name": sheet.title,
            "row_count": max_row,
            "col_count": max_col,
        })
    return {"sheets": sheets}


def xlsx_read_range(params: Dict[str, Any]) -> Dict[str, Any]:
    """Read a specific range from an xlsx sheet, with optional chunk_info."""
    openpyxl = import_module("openpyxl")
    from openpyxl.utils import get_column_letter
    path = resolve_path(params)
    wb = openpyxl.load_workbook(path, data_only=False)
    sheets = wb.sheetnames
    sheet_name = params.get("sheet") or (sheets[0] if sheets else None)
    if not sheet_name:
        raise WorkerError("VALIDATION_FAILED", "missing sheet")
    if sheet_name not in wb.sheetnames:
        raise WorkerError("VALIDATION_FAILED", f"unknown sheet: {sheet_name}")
    ws = wb[sheet_name]

    cell_range = params.get("range")
    lines = []
    lines.append(f"# Sheet: {sheet_name}")

    sheet_max_row = ws.max_row or 0
    sheet_min_row = ws.min_row or 1
    sheet_min_col = ws.min_column or 1
    sheet_max_col = ws.max_column or 1
    total_rows = max(0, sheet_max_row - sheet_min_row + 1) if sheet_max_row > 0 else 0

    if cell_range:
        # Parse range like "A1:D10"
        from openpyxl.utils import range_boundaries
        min_col, min_row, max_col, max_row = range_boundaries(cell_range)
        for row in ws.iter_rows(min_row=min_row, max_row=max_row,
                                min_col=min_col, max_col=max_col, values_only=True):
            parts = []
            for val in row:
                if val is None:
                    parts.append("")
                else:
                    parts.append(str(val))
            lines.append("\t".join(parts))

        # Compute chunk_info
        total_chunks = max(1, (total_rows + XLSX_CHUNK_ROWS - 1) // XLSX_CHUNK_ROWS)
        chunk_index = (min_row - sheet_min_row) // XLSX_CHUNK_ROWS if total_rows > 0 else 0
        has_more = max_row < sheet_max_row

        result = {"text": "\n".join(lines)}
        if total_chunks > 1:
            result["chunk_info"] = {
                "chunk_index": chunk_index,
                "total_chunks": total_chunks,
                "has_more": has_more,
                "range": cell_range,
            }
        return result
    else:
        # Read all (limited to 200 rows, 50 cols)
        for i, row in enumerate(ws.iter_rows(values_only=True)):
            if i >= MAX_GRID_ROWS:
                lines.append(f"... (truncated at {MAX_GRID_ROWS} rows)")
                break
            parts = []
            for j, val in enumerate(row):
                if j >= MAX_GRID_COLS:
                    break
                if val is None:
                    parts.append("")
                else:
                    parts.append(str(val))
            lines.append("\t".join(parts))

        return {"text": "\n".join(lines)}


def xlsx_render_grid(params: Dict[str, Any]) -> Dict[str, Any]:
    openpyxl = import_module("openpyxl")
    path = resolve_path(params)
    wb = openpyxl.load_workbook(path, data_only=False)
    sheets = wb.sheetnames
    sheet_name = params.get("sheet") or (sheets[0] if sheets else None)
    if not sheet_name:
        raise WorkerError("VALIDATION_FAILED", "missing sheet")
    if sheet_name not in wb.sheetnames:
        raise WorkerError("VALIDATION_FAILED", "unknown sheet")
    ws = wb[sheet_name]
    row_start = int(params.get("row_start", 0))
    col_start = int(params.get("col_start", 0))
    row_count = int(params.get("row_count", MAX_GRID_ROWS))
    col_count = int(params.get("col_count", MAX_GRID_COLS))
    if row_count > MAX_GRID_ROWS:
        row_count = MAX_GRID_ROWS
    if col_count > MAX_GRID_COLS:
        col_count = MAX_GRID_COLS
    merged_skip = set()
    if ws.merged_cells.ranges:
        row_min = row_start + 1
        row_max = row_start + row_count
        col_min = col_start + 1
        col_max = col_start + col_count
        for merged in ws.merged_cells.ranges:
            min_row = merged.min_row
            max_row = merged.max_row
            min_col = merged.min_col
            max_col = merged.max_col
            if max_row < row_min or min_row > row_max or max_col < col_min or min_col > col_max:
                continue
            top_left = (min_row, min_col)
            for r in range(max(min_row, row_min), min(max_row, row_max) + 1):
                for c in range(max(min_col, col_min), min(max_col, col_max) + 1):
                    if (r, c) != top_left:
                        merged_skip.add((r, c))
    cells: List[List[Dict[str, Any]]] = []
    for r in range(row_start, row_start + row_count):
        row_cells = []
        for c in range(col_start, col_start + col_count):
            cell = ws.cell(row=r + 1, column=c + 1)
            value = cell.value
            data_type = cell.data_type
            cell_info: Dict[str, Any] = {"value": None, "type": "blank", "formula": None}
            if (r + 1, c + 1) in merged_skip:
                row_cells.append(cell_info)
                continue
            if data_type == "f":
                cell_info["formula"] = str(value) if value is not None else None
            elif value is None:
                pass
            elif isinstance(value, bool):
                cell_info["value"] = value
                cell_info["type"] = "boolean"
            elif isinstance(value, (int, float)):
                cell_info["value"] = value
                cell_info["type"] = "number"
            else:
                cell_info["value"] = str(value)
                cell_info["type"] = "string"
            row_cells.append(cell_info)
        cells.append(row_cells)
    return {
        "sheets": sheets,
        "row_count": row_count,
        "col_count": col_count,
        "cells": cells,
    }


def xlsx_get_map(params: Dict[str, Any]) -> Dict[str, Any]:
    """Return structural map of an xlsx file: sheets, islands, chunks, flags."""
    openpyxl = import_module("openpyxl")
    from openpyxl.utils import get_column_letter
    path = resolve_path(params)
    wb = openpyxl.load_workbook(path, data_only=False)
    sheets_out: List[Dict[str, Any]] = []

    for ws in wb.worksheets:
        max_row = ws.max_row or 0
        max_col = ws.max_column or 0
        if max_row == 0 or max_col == 0:
            sheets_out.append({
                "name": ws.title,
                "used_range": None,
                "row_count": 0,
                "col_count": 0,
                "islands": [],
                "chunks": [],
                "has_charts": False,
                "has_merged_cells": False,
                "has_conditional_formatting": False,
                "has_formulas": False,
            })
            continue

        # Determine actual used range by scanning for non-empty cells
        min_row = ws.min_row or 1
        min_col = ws.min_column or 1

        used_range = {
            "min_row": min_row,
            "max_row": max_row,
            "min_col": min_col,
            "max_col": max_col,
        }

        # Island detection: scan rows, group consecutive non-blank rows
        islands: List[Dict[str, Any]] = []
        island_start = None
        for r in range(min_row, max_row + 1):
            blank = True
            for c in range(min_col, max_col + 1):
                val = ws.cell(row=r, column=c).value
                if val is not None and str(val).strip() != "":
                    blank = False
                    break
            if not blank:
                if island_start is None:
                    island_start = r
            else:
                if island_start is not None:
                    islands.append(_build_island(ws, island_start, r - 1, min_col, max_col))
                    island_start = None
        if island_start is not None:
            islands.append(_build_island(ws, island_start, max_row, min_col, max_col))

        # Chunking: divide used range into chunks of XLSX_CHUNK_ROWS
        chunks: List[Dict[str, Any]] = []
        chunk_idx = 0
        r = min_row
        while r <= max_row:
            end_r = min(r + XLSX_CHUNK_ROWS - 1, max_row)
            range_str = "{}{}:{}{}".format(
                get_column_letter(min_col), r,
                get_column_letter(max_col), end_r,
            )
            chunks.append({
                "index": chunk_idx,
                "range": range_str,
                "rows": end_r - r + 1,
            })
            chunk_idx += 1
            r = end_r + 1

        # Sheet-level flags
        has_charts = len(ws._charts) > 0 if hasattr(ws, '_charts') else False
        has_merged = len(ws.merged_cells.ranges) > 0
        has_cond_fmt = len(ws.conditional_formatting) > 0 if hasattr(ws, 'conditional_formatting') else False
        has_formulas = False
        for row in ws.iter_rows(min_row=min_row, max_row=max_row,
                                min_col=min_col, max_col=max_col):
            for cell in row:
                if cell.data_type == "f":
                    has_formulas = True
                    break
            if has_formulas:
                break

        sheets_out.append({
            "name": ws.title,
            "used_range": used_range,
            "row_count": max_row - min_row + 1,
            "col_count": max_col - min_col + 1,
            "islands": islands,
            "chunks": chunks,
            "has_charts": has_charts,
            "has_merged_cells": has_merged,
            "has_conditional_formatting": has_cond_fmt,
            "has_formulas": has_formulas,
        })

    return {"sheets": sheets_out}


def _build_island(ws, start_row: int, end_row: int, min_col: int, max_col: int) -> Dict[str, Any]:
    """Build an island descriptor for consecutive non-blank rows."""
    from openpyxl.utils import get_column_letter

    # Determine the actual used column range within this island
    actual_min_col = max_col + 1
    actual_max_col = min_col - 1
    for r in range(start_row, end_row + 1):
        for c in range(min_col, max_col + 1):
            val = ws.cell(row=r, column=c).value
            if val is not None and str(val).strip() != "":
                if c < actual_min_col:
                    actual_min_col = c
                if c > actual_max_col:
                    actual_max_col = c
    if actual_min_col > actual_max_col:
        actual_min_col = min_col
        actual_max_col = max_col

    range_str = "{}{}:{}{}".format(
        get_column_letter(actual_min_col), start_row,
        get_column_letter(actual_max_col), end_row,
    )

    # Detect headers: check if all non-empty cells in the first row are strings
    headers = None
    first_row_all_strings = True
    header_values: List[Optional[str]] = []
    has_any_value = False
    for c in range(actual_min_col, actual_max_col + 1):
        val = ws.cell(row=start_row, column=c).value
        if val is None or str(val).strip() == "":
            header_values.append(None)
        elif isinstance(val, str):
            header_values.append(val)
            has_any_value = True
        else:
            first_row_all_strings = False
            header_values.append(str(val))
            has_any_value = True

    if first_row_all_strings and has_any_value:
        headers = header_values
        label = "header"
    else:
        label = "data"

    return {
        "range": range_str,
        "label": label,
        "row_count": end_row - start_row + 1,
        "col_count": actual_max_col - actual_min_col + 1,
        "headers": headers,
    }


# ------------------ PPTX ------------------

def pptx_apply_ops(params: Dict[str, Any]) -> Dict[str, Any]:
    pptx = import_module("pptx")
    path = resolve_path(params)
    ops = params.get("ops") or []
    create_new = params.get("create_new", False)
    copy_from = params.get("copy_from")

    if not isinstance(ops, list) or len(ops) == 0:
        raise WorkerError("VALIDATION_FAILED", "missing ops")

    # Handle create_new with optional copy_from
    if create_new:
        if copy_from:
            source_params = dict(params)
            source_params["path"] = copy_from
            source_path = resolve_path(source_params)
            if not os.path.exists(source_path):
                raise WorkerError("FILE_READ_FAILED", f"source file not found: {copy_from}")
            prs = pptx.Presentation(source_path)
            log_info("pptx.copy_from", source=copy_from, target=path)
        else:
            prs = pptx.Presentation()
            log_info("pptx.create_new", target=path)
    elif os.path.exists(path):
        prs = pptx.Presentation(path)
    else:
        prs = pptx.Presentation()
    layout_map = {
        "title_and_content": 1,
        "title_only": 5,
        "section_header": 2,
    }
    for op in ops:
        if not isinstance(op, dict):
            raise WorkerError("VALIDATION_FAILED", "invalid op")
        name = op.get("op")
        if name == "add_slide":
            layout_name = op.get("layout") or "title_and_content"
            layout_idx = layout_map.get(layout_name, 1)
            if layout_idx >= len(prs.slide_layouts):
                layout_idx = 0
            slide = prs.slides.add_slide(prs.slide_layouts[layout_idx])
            title_shape = slide.shapes.title if slide.shapes.title else None
            body_shape = _find_body_placeholder(slide)
            title_runs = op.get("title_runs")
            body_runs = op.get("body_runs")
            if title_runs is not None:
                _pptx_set_text_runs(
                    title_shape,
                    title_runs,
                    op_name=name,
                    run_field="title_runs",
                    style_source=op,
                )
            elif op.get("title") is not None:
                _pptx_set_shape_text(
                    title_shape,
                    op.get("title"),
                    op_name=name,
                    style_source=op,
                    target="title",
                )
            if body_runs is not None:
                _pptx_set_text_runs(
                    body_shape,
                    body_runs,
                    op_name=name,
                    run_field="body_runs",
                    style_source=op,
                )
            elif op.get("body") is not None:
                _set_pptx_body_text(
                    body_shape,
                    op.get("body"),
                    op_name=name,
                    style_source=op,
                    target="body",
                )
        elif name == "set_slide_text":
            idx = int(op.get("index", -1))
            if idx < 0 or idx >= len(prs.slides):
                raise WorkerError("VALIDATION_FAILED", "invalid slide index")
            slide = prs.slides[idx]
            title_shape = slide.shapes.title if slide.shapes.title else None
            body_shape = _find_body_placeholder(slide)
            title_runs = op.get("title_runs")
            body_runs = op.get("body_runs")
            if title_runs is not None:
                _pptx_set_text_runs(
                    title_shape,
                    title_runs,
                    op_name=name,
                    run_field="title_runs",
                    style_source=op,
                )
            elif op.get("title") is not None:
                _pptx_set_shape_text(
                    title_shape,
                    op.get("title"),
                    op_name=name,
                    style_source=op,
                    target="title",
                )
            if body_runs is not None:
                _pptx_set_text_runs(
                    body_shape,
                    body_runs,
                    op_name=name,
                    run_field="body_runs",
                    style_source=op,
                )
            elif op.get("body") is not None:
                _set_pptx_body_text(
                    body_shape,
                    op.get("body"),
                    op_name=name,
                    style_source=op,
                    target="body",
                )
        elif name == "append_bullets":
            idx = int(op.get("index", -1))
            if idx < 0 or idx >= len(prs.slides):
                raise WorkerError("VALIDATION_FAILED", "invalid slide index")
            slide = prs.slides[idx]
            bullets = op.get("bullets") or []
            body_shape = _find_body_placeholder(slide)
            if body_shape is None:
                continue
            tf = body_shape.text_frame
            for i, bullet in enumerate(bullets):
                if i == 0 and len(tf.paragraphs) == 1 and not (tf.paragraphs[0].text or "").strip():
                    p = tf.paragraphs[0]
                else:
                    p = tf.add_paragraph()
                p.text = str(bullet)
                p.level = 0
                _pptx_apply_paragraph_format(
                    p,
                    op,
                    op_name=name,
                    target="body",
                )
        else:
            raise WorkerError("VALIDATION_FAILED", f"unsupported op: {name}")
    os.makedirs(os.path.dirname(path), exist_ok=True)
    prs.save(path)
    return {"ok": True}


def _find_body_placeholder(slide):
    title_shape = None
    try:
        title_shape = slide.shapes.title
    except Exception:
        title_shape = None

    for shape in slide.placeholders:
        if shape is title_shape or not getattr(shape, "has_text_frame", False):
            continue
        ph_type = _placeholder_type_value(shape)
        if ph_type in (2, 7):  # BODY or OBJECT (often used for content placeholders)
            return shape

    for shape in slide.placeholders:
        if shape is title_shape or not getattr(shape, "has_text_frame", False):
            continue
        return shape

    for shape in slide.shapes:
        if shape is title_shape or not getattr(shape, "has_text_frame", False):
            continue
        return shape

    return None


def _placeholder_type_value(shape) -> Optional[int]:
    try:
        placeholder_type = shape.placeholder_format.type
        if hasattr(placeholder_type, "value"):
            return int(placeholder_type.value)
        return int(placeholder_type)
    except Exception:
        return None


def _pptx_set_shape_text(shape, value: Any, op_name: str, style_source: Dict[str, Any], target: str) -> None:
    if shape is None or not getattr(shape, "has_text_frame", False):
        return
    shape.text = str(value)
    tf = shape.text_frame
    if len(tf.paragraphs) > 0:
        _pptx_apply_paragraph_format(tf.paragraphs[0], style_source, op_name=op_name, target=target)


def _pptx_set_text_runs(shape, runs: Any, op_name: str, run_field: str, style_source: Dict[str, Any]) -> None:
    if shape is None or not getattr(shape, "has_text_frame", False):
        return
    target = "title" if run_field == "title_runs" else "body"
    if not isinstance(runs, list):
        _style_warn("pptx.run", run_field, runs, f"{run_field} must be list", op=op_name, target=target)
        fallback_key = "title" if run_field == "title_runs" else "body"
        if fallback_key in style_source:
            _pptx_set_shape_text(shape, style_source.get(fallback_key), op_name=op_name, style_source=style_source, target=target)
        return
    tf = shape.text_frame
    tf.clear()
    paragraph = tf.paragraphs[0]
    for run_index, item in enumerate(runs):
        if isinstance(item, dict):
            run = paragraph.add_run()
            run.text = str(item.get("text") or "")
            _pptx_apply_run_style(run, item, op_name=op_name, run_field=run_field, run_index=run_index, target=target)
        else:
            _style_warn(
                "pptx.run",
                run_field,
                item,
                "run entry must be object; coerced to plain text",
                op=op_name,
                run_index=run_index,
                target=target,
            )
            run = paragraph.add_run()
            run.text = str(item)
    _pptx_apply_paragraph_format(paragraph, style_source, op_name=op_name, target=target)


def _set_pptx_body_text(body_shape, body: Any, op_name: str, style_source: Dict[str, Any], target: str) -> None:
    if body_shape is None:
        return
    text = str(body)
    lines = [line.strip() for line in text.splitlines() if line.strip()]
    if len(lines) <= 1:
        body_shape.text = text
        if getattr(body_shape, "has_text_frame", False):
            tf = body_shape.text_frame
            if len(tf.paragraphs) > 0:
                _pptx_apply_paragraph_format(tf.paragraphs[0], style_source, op_name=op_name, target=target)
        return

    tf = body_shape.text_frame
    tf.clear()
    for idx, line in enumerate(lines):
        if idx == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = line
        p.level = 0
        _pptx_apply_paragraph_format(p, style_source, op_name=op_name, target=target)


def _pptx_apply_run_style(run: Any, payload: Dict[str, Any], op_name: str, run_field: str, run_index: int, target: str) -> None:
    from pptx.util import Pt
    from pptx.dml.color import RGBColor

    allowed_keys = {"text", "font_name", "font_size", "bold", "italic", "underline", "font_color", "highlight_color"}
    for key in payload.keys():
        if key not in allowed_keys:
            _style_warn(
                "pptx.run",
                key,
                payload.get(key),
                "unknown style key",
                op=op_name,
                run_field=run_field,
                run_index=run_index,
                target=target,
            )

    font = run.font
    font_name = payload.get("font_name")
    if font_name is not None:
        if isinstance(font_name, str) and font_name.strip():
            font.name = font_name.strip()
        else:
            _style_warn(
                "pptx.run",
                "font_name",
                font_name,
                "font_name must be non-empty string",
                op=op_name,
                run_field=run_field,
                run_index=run_index,
                target=target,
            )

    font_size = payload.get("font_size")
    if font_size is not None:
        parsed = _coerce_float(font_size, None)
        if parsed is None or parsed <= 0:
            _style_warn(
                "pptx.run",
                "font_size",
                font_size,
                "font_size must be a positive number",
                op=op_name,
                run_field=run_field,
                run_index=run_index,
                target=target,
            )
        else:
            font.size = Pt(parsed)

    for field, attr in (("bold", "bold"), ("italic", "italic"), ("underline", "underline")):
        raw = payload.get(field)
        if raw is None:
            continue
        parsed = _coerce_bool_strict(raw, None)
        if parsed is None:
            _style_warn(
                "pptx.run",
                field,
                raw,
                f"{field} must be boolean",
                op=op_name,
                run_field=run_field,
                run_index=run_index,
                target=target,
            )
            continue
        setattr(font, attr, parsed)

    font_color = payload.get("font_color")
    if font_color is not None:
        rgb = _normalize_hex_rgb(font_color)
        if rgb is None:
            _style_warn(
                "pptx.run",
                "font_color",
                font_color,
                "font_color must be #RRGGBB",
                op=op_name,
                run_field=run_field,
                run_index=run_index,
                target=target,
            )
        else:
            try:
                font.color.rgb = RGBColor.from_string(rgb)
            except Exception as err:
                _style_warn(
                    "pptx.run",
                    "font_color",
                    font_color,
                    "failed to apply font_color",
                    op=op_name,
                    run_field=run_field,
                    run_index=run_index,
                    target=target,
                    error=str(err),
                )

    if payload.get("highlight_color") is not None:
        _style_warn(
            "pptx.run",
            "highlight_color",
            payload.get("highlight_color"),
            "highlight_color is not supported by python-pptx",
            op=op_name,
            run_field=run_field,
            run_index=run_index,
            target=target,
        )


def _pptx_apply_paragraph_format(paragraph: Any, payload: Dict[str, Any], op_name: str, target: str) -> None:
    from pptx.enum.text import PP_ALIGN
    from pptx.util import Pt

    alignment = payload.get("alignment")
    if alignment is not None:
        if not isinstance(alignment, str):
            _style_warn("pptx.paragraph", "alignment", alignment, "alignment must be string", op=op_name, target=target)
        else:
            align_map = {
                "left": PP_ALIGN.LEFT,
                "center": PP_ALIGN.CENTER,
                "right": PP_ALIGN.RIGHT,
                "justify": PP_ALIGN.JUSTIFY,
            }
            parsed = align_map.get(_slug(alignment))
            if parsed is None:
                _style_warn("pptx.paragraph", "alignment", alignment, "unsupported alignment", op=op_name, target=target)
            else:
                paragraph.alignment = parsed

    for field in ("space_before", "space_after"):
        raw = payload.get(field)
        if raw is None:
            continue
        parsed = _coerce_float(raw, None)
        if parsed is None or parsed < 0:
            _style_warn(
                "pptx.paragraph",
                field,
                raw,
                f"{field} must be a non-negative number",
                op=op_name,
                target=target,
            )
            continue
        setattr(paragraph, field, Pt(parsed))

    line_spacing = payload.get("line_spacing")
    if line_spacing is not None:
        parsed = _coerce_float(line_spacing, None)
        if parsed is None or parsed <= 0:
            _style_warn(
                "pptx.paragraph",
                "line_spacing",
                line_spacing,
                "line_spacing must be a positive number",
                op=op_name,
                target=target,
            )
        else:
            paragraph.line_spacing = parsed


def pptx_extract_text(params: Dict[str, Any]) -> Dict[str, Any]:
    """Extract text from pptx, optionally for a specific slide."""
    pptx = import_module("pptx")
    path = resolve_path(params)
    prs = pptx.Presentation(path)
    slide_index = params.get("slide_index")
    total_slides = len(prs.slides)

    if slide_index is not None:
        idx = int(slide_index)
        if idx < 0 or idx >= total_slides:
            raise WorkerError("VALIDATION_FAILED", "invalid slide_index")
        slide = prs.slides[idx]
        lines = [f"# Slide {idx + 1}"]
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text = shape.text
                if text:
                    lines.append(text)
        result: Dict[str, Any] = {"text": "\n".join(lines)}
        if total_slides > 1:
            result["chunk_info"] = {
                "chunk_index": idx,
                "total_chunks": total_slides,
                "has_more": idx < total_slides - 1,
                "range": f"slide {idx}",
            }
        return result

    lines = []
    for idx, slide in enumerate(prs.slides, start=1):
        lines.append(f"# Slide {idx}")
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text = shape.text
                if text:
                    lines.append(text)
    return {"text": "\n".join(lines)}


def pptx_render_slide(params: Dict[str, Any]) -> Dict[str, Any]:
    path = resolve_path(params)
    slide_index = int(params.get("slide_index", 0))
    scale = clamp_scale(params.get("scale", 1.0))
    try:
        pdf_path, tmp_dir = convert_to_pdf(path)
    except WorkerError as err:
        if err.code not in ("TOOL_WORKER_UNAVAILABLE", "FILE_READ_FAILED"):
            raise
        log_info("preview.pptx_fallback", path=path, reason=err.message)
        return render_pptx_slide_fallback(path, slide_index, scale)
    except Exception as err:
        log_info("preview.pptx_fallback", path=path, reason=str(err))
        return render_pptx_slide_fallback(path, slide_index, scale)
    try:
        result = render_pdf_page(pdf_path, slide_index, scale)
        slide_count = result.pop("page_count", 0)
        result["slide_count"] = slide_count
        return result
    finally:
        shutil.rmtree(tmp_dir, ignore_errors=True)


def pptx_get_map(params: Dict[str, Any]) -> Dict[str, Any]:
    """Return structural map of a pptx file: slide metadata."""
    pptx = import_module("pptx")
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    path = resolve_path(params)
    prs = pptx.Presentation(path)
    slides_out: List[Dict[str, Any]] = []

    for idx, slide in enumerate(prs.slides):
        title = None
        if slide.shapes.title:
            title = slide.shapes.title.text or None

        layout_name = slide.slide_layout.name if slide.slide_layout else None

        has_images = False
        has_charts = False
        for shape in slide.shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                has_images = True
            if shape.shape_type == MSO_SHAPE_TYPE.CHART:
                has_charts = True
            if has_images and has_charts:
                break

        has_notes = False
        if slide.has_notes_slide:
            notes_slide = slide.notes_slide
            if notes_slide and notes_slide.notes_text_frame:
                notes_text = notes_slide.notes_text_frame.text.strip()
                if notes_text:
                    has_notes = True

        slides_out.append({
            "index": idx,
            "title": title,
            "layout": layout_name,
            "has_images": has_images,
            "has_charts": has_charts,
            "has_notes": has_notes,
        })

    return {
        "slides": slides_out,
        "slide_count": len(prs.slides),
    }


def pptx_get_styles(params: Dict[str, Any]) -> Dict[str, Any]:
    pptx = import_module("pptx")
    path = resolve_path(params)
    prs = pptx.Presentation(path)

    layouts: List[Dict[str, Any]] = []
    seen_layouts: Set[str] = set()
    for idx, layout in enumerate(prs.slide_layouts):
        name = str(getattr(layout, "name", "") or "").strip()
        key = name or f"layout_{idx}"
        if key in seen_layouts:
            continue
        seen_layouts.add(key)
        layouts.append({"asset_id": f"layout:{key}", "type": "layout", "name": key, "index": idx})

    font_usage: Dict[str, Dict[str, Any]] = {}
    text_style_assets: List[Dict[str, Any]] = []
    image_assets: List[Dict[str, Any]] = []

    for slide_index, slide in enumerate(prs.slides):
        for shape_index, shape in enumerate(slide.shapes):
            if getattr(shape, "has_text_frame", False):
                text_runs: List[Dict[str, Any]] = []
                for paragraph_index, paragraph in enumerate(shape.text_frame.paragraphs):
                    for run_index, run in enumerate(paragraph.runs):
                        run_payload = _pptx_run_payload(run, run_index)
                        run_payload["paragraph_index"] = paragraph_index
                        text_runs.append(run_payload)
                        font_name = str(run_payload.get("font_name") or "").strip()
                        if font_name:
                            entry = font_usage.setdefault(
                                font_name,
                                {"family": font_name, "count": 0},
                            )
                            entry["count"] = int(entry.get("count", 0)) + 1
                if len(text_runs) > 0:
                    text_style_assets.append(
                        {
                            "asset_id": f"text_style:{slide_index}:{shape_index}",
                            "type": "text_style",
                            "slide_index": slide_index,
                            "shape_index": shape_index,
                            "shape_name": getattr(shape, "name", None),
                            "run_count": len(text_runs),
                            "runs": text_runs,
                        }
                    )
            if _pptx_shape_type_name(shape) == "PICTURE":
                image_ref = _pptx_image_ref(shape)
                if image_ref is not None:
                    image_assets.append(
                        {
                            "asset_id": f"image:{slide_index}:{shape_index}",
                            "type": "image",
                            "slide_index": slide_index,
                            "shape_index": shape_index,
                            "shape_name": getattr(shape, "name", None),
                            **image_ref,
                        }
                    )

    fonts = sorted(font_usage.values(), key=lambda item: str(item.get("family") or "").lower())
    assets = sorted(layouts + text_style_assets + image_assets, key=lambda item: (str(item.get("type") or ""), str(item.get("asset_id") or "")))

    return {
        "format": "pptx",
        "slide_count": len(prs.slides),
        "layouts": layouts,
        "fonts": fonts,
        "text_style_assets": text_style_assets,
        "image_assets": image_assets,
        "assets": assets,
        "supported_copy_asset_types": ["text_style", "image"],
    }


def pptx_copy_assets(params: Dict[str, Any]) -> Dict[str, Any]:
    pptx = import_module("pptx")
    source_path = resolve_named_path(params, "source_path", "source_root", "draft")
    target_path = resolve_named_path(params, "target_path", "target_root", "draft", require_write=True)
    assets = _parse_asset_list(params)

    src_prs = pptx.Presentation(source_path)
    dst_prs = pptx.Presentation(target_path) if os.path.exists(target_path) else pptx.Presentation()
    copied: List[Dict[str, Any]] = []

    normalized_assets = [_pptx_parse_asset(asset) for asset in assets]
    for item in normalized_assets:
        asset_type = item["type"]
        if asset_type == "text_style":
            src_slide = _pptx_get_slide(src_prs, item["source_slide_index"], "source_slide_index")
            dst_slide = _pptx_get_slide(dst_prs, item["target_slide_index"], "target_slide_index")
            src_shape = _pptx_get_shape(src_slide, item["source_shape_index"], "source_shape_index")
            dst_shape = _pptx_get_shape(dst_slide, item["target_shape_index"], "target_shape_index")
            _pptx_copy_text_style(src_shape, dst_shape)
            copied.append(
                {
                    "type": "text_style",
                    "source": f"{item['source_slide_index']}:{item['source_shape_index']}",
                    "target": f"{item['target_slide_index']}:{item['target_shape_index']}",
                }
            )
            continue
        if asset_type == "image":
            src_slide = _pptx_get_slide(src_prs, item["source_slide_index"], "source_slide_index")
            dst_slide = _pptx_get_slide(dst_prs, item["target_slide_index"], "target_slide_index")
            src_shape = _pptx_get_shape(src_slide, item["source_shape_index"], "source_shape_index")
            image_ref = _pptx_image_ref(src_shape)
            if image_ref is None:
                raise WorkerError("VALIDATION_FAILED", "source shape does not contain an image")
            try:
                blob = src_shape.image.blob
            except Exception:
                raise WorkerError("FILE_READ_FAILED", "failed to read source image blob")
            left = getattr(src_shape, "left", 0)
            top = getattr(src_shape, "top", 0)
            width = getattr(src_shape, "width", None)
            height = getattr(src_shape, "height", None)
            dst_slide.shapes.add_picture(io.BytesIO(blob), left, top, width=width, height=height)
            copied.append(
                {
                    "type": "image",
                    "source": f"{item['source_slide_index']}:{item['source_shape_index']}",
                    "target_slide_index": item["target_slide_index"],
                    "size_bytes": len(blob),
                    "sha1": _sha1_of_bytes(blob),
                }
            )
            continue
        raise WorkerError("VALIDATION_FAILED", f"unsupported asset type: {asset_type}")

    os.makedirs(os.path.dirname(target_path), exist_ok=True)
    dst_prs.save(target_path)
    return {"ok": True, "copied_count": len(copied), "copied": copied}


def pptx_get_slide_content(params: Dict[str, Any]) -> Dict[str, Any]:
    """Return deterministic structured payload for one pptx slide."""
    pptx = import_module("pptx")
    path = resolve_path(params)
    prs = pptx.Presentation(path)
    slide_count = len(prs.slides)
    slide_index = parse_non_negative_index(
        params.get("slide_index", params.get("index")),
        "slide_index",
        default=0,
    )
    if slide_index >= slide_count:
        raise WorkerError("VALIDATION_FAILED", "invalid slide_index")
    detail = str(params.get("detail") or "positioned").strip().lower()
    if detail not in ("positioned", "layout-lite", "layout_lite", "lite"):
        raise WorkerError("VALIDATION_FAILED", "invalid detail")
    slide = prs.slides[slide_index]
    slide_title = slide.shapes.title.text if slide.shapes.title else None
    layout_name = slide.slide_layout.name if slide.slide_layout else None
    shapes = [_pptx_shape_payload(shape, idx) for idx, shape in enumerate(slide.shapes)]
    notes = None
    if slide.has_notes_slide:
        try:
            notes_text = slide.notes_slide.notes_text_frame.text.strip()
            if notes_text:
                notes = notes_text
        except Exception:
            notes = None
    slide_payload: Dict[str, Any] = {
        "index": slide_index,
        "title": slide_title,
        "layout": layout_name,
        "shapes": shapes,
        "notes": notes,
    }
    if detail == "positioned":
        slide_payload["positioned"] = _pptx_positioned_slide_payload(slide, prs)
        slide_payload["render_mode"] = "positioned"
    else:
        slide_payload["render_mode"] = "layout_lite"
    return {
        "slide_index": slide_index,
        "slide_count": slide_count,
        "slide": slide_payload,
    }


def _pptx_parse_asset(asset: Any) -> Dict[str, Any]:
    if isinstance(asset, str):
        text = asset.strip()
        if not text:
            raise WorkerError("VALIDATION_FAILED", "asset selector is empty")
        if ":" not in text:
            raise WorkerError("VALIDATION_FAILED", f"invalid asset selector: {text}")
        kind, payload = text.split(":", 1)
        kind = _slug(kind)
        payload = payload.strip()
        if kind in ("text_style", "image"):
            source, _, target = payload.partition("->")
            src_slide, src_shape = _pptx_parse_slide_shape_ref(source)
            target_slide = src_slide
            target_shape = src_shape
            if target:
                if kind == "text_style":
                    target_slide, target_shape = _pptx_parse_slide_shape_ref(target)
                else:
                    target_slide, _ = _pptx_parse_slide_shape_ref(target + ":0" if ":" not in target else target)
            out = {
                "type": kind,
                "source_slide_index": src_slide,
                "source_shape_index": src_shape,
                "target_slide_index": target_slide,
            }
            if kind == "text_style":
                out["target_shape_index"] = target_shape
            return out
        raise WorkerError("VALIDATION_FAILED", f"unsupported asset type: {kind}")
    if not isinstance(asset, dict):
        raise WorkerError("VALIDATION_FAILED", "asset selector must be string or object")

    kind = _slug(asset.get("type") or "")
    if not kind:
        raw = str(asset.get("asset_id") or "").strip()
        if ":" in raw:
            kind = _slug(raw.split(":", 1)[0])
    if kind not in ("text_style", "image"):
        raise WorkerError("VALIDATION_FAILED", f"unsupported asset type: {kind or 'unknown'}")

    src_slide = _coerce_int(asset.get("source_slide_index"), None)
    src_shape = _coerce_int(asset.get("source_shape_index"), None)
    if src_slide is None or src_shape is None:
        source = str(asset.get("source") or "").strip()
        if source:
            src_slide, src_shape = _pptx_parse_slide_shape_ref(source)
    if src_slide is None or src_shape is None:
        raise WorkerError("VALIDATION_FAILED", f"{kind} asset missing source_slide_index/source_shape_index")

    target_slide = _coerce_int(asset.get("target_slide_index"), src_slide)
    out = {
        "type": kind,
        "source_slide_index": src_slide,
        "source_shape_index": src_shape,
        "target_slide_index": target_slide,
    }
    if kind == "text_style":
        target_shape = _coerce_int(asset.get("target_shape_index"), None)
        if target_shape is None:
            target = str(asset.get("target") or "").strip()
            if target:
                _, target_shape = _pptx_parse_slide_shape_ref(target)
        if target_shape is None:
            raise WorkerError("VALIDATION_FAILED", "text_style asset missing target_shape_index")
        out["target_shape_index"] = target_shape
    return out


def _pptx_parse_slide_shape_ref(value: str) -> Tuple[int, int]:
    text = value.strip()
    if ":" not in text:
        raise WorkerError("VALIDATION_FAILED", f"invalid slide/shape selector: {value}")
    slide_text, shape_text = text.split(":", 1)
    try:
        slide_index = int(slide_text.strip())
        shape_index = int(shape_text.strip())
    except Exception:
        raise WorkerError("VALIDATION_FAILED", f"invalid slide/shape selector: {value}")
    if slide_index < 0 or shape_index < 0:
        raise WorkerError("VALIDATION_FAILED", f"invalid slide/shape selector: {value}")
    return slide_index, shape_index


def _pptx_get_slide(prs, slide_index: int, field_name: str):
    if slide_index < 0 or slide_index >= len(prs.slides):
        raise WorkerError("VALIDATION_FAILED", f"invalid {field_name}")
    return prs.slides[slide_index]


def _pptx_get_shape(slide, shape_index: int, field_name: str):
    if shape_index < 0 or shape_index >= len(slide.shapes):
        raise WorkerError("VALIDATION_FAILED", f"invalid {field_name}")
    return slide.shapes[shape_index]


def _pptx_copy_text_style(src_shape, dst_shape) -> None:
    if not getattr(src_shape, "has_text_frame", False):
        raise WorkerError("VALIDATION_FAILED", "source shape does not contain text")
    if not getattr(dst_shape, "has_text_frame", False):
        raise WorkerError("VALIDATION_FAILED", "target shape does not contain text")

    src_paragraphs = list(src_shape.text_frame.paragraphs)
    if len(src_paragraphs) == 0:
        return
    src_runs = []
    for paragraph in src_paragraphs:
        for run in paragraph.runs:
            src_runs.append(run)
    if len(src_runs) == 0:
        return

    dst_runs = []
    for paragraph in dst_shape.text_frame.paragraphs:
        for run in paragraph.runs:
            dst_runs.append(run)
    if len(dst_runs) == 0:
        return

    template_run = src_runs[0]
    template_font = template_run.font
    for run in dst_runs:
        font = run.font
        if template_font is None or font is None:
            continue
        for attr in ("name", "size", "bold", "italic", "underline"):
            try:
                setattr(font, attr, getattr(template_font, attr))
            except Exception:
                pass
        try:
            if template_font.color is not None and template_font.color.rgb is not None:
                font.color.rgb = template_font.color.rgb
        except Exception:
            pass


def _pptx_shape_payload(shape, shape_index: int) -> Dict[str, Any]:
    shape_type_name = _pptx_shape_type_name(shape)
    payload: Dict[str, Any] = {
        "index": shape_index,
        "name": getattr(shape, "name", None),
        "shape_type": shape_type_name,
        "z_index": shape_index,
        "left": _pptx_length_value(getattr(shape, "left", None)),
        "top": _pptx_length_value(getattr(shape, "top", None)),
        "width": _pptx_length_value(getattr(shape, "width", None)),
        "height": _pptx_length_value(getattr(shape, "height", None)),
        "rotation": _pptx_rotation_value(shape),
    }
    if getattr(shape, "is_placeholder", False):
        try:
            payload["placeholder_type"] = str(shape.placeholder_format.type)
        except Exception:
            payload["placeholder_type"] = None
    text_blocks: List[Dict[str, Any]] = []
    if getattr(shape, "has_text_frame", False):
        try:
            for paragraph_index, paragraph in enumerate(shape.text_frame.paragraphs):
                runs: List[Dict[str, Any]] = []
                for run_index, run in enumerate(paragraph.runs):
                    runs.append(_pptx_run_payload(run, run_index))
                text_blocks.append(
                    {
                        "index": paragraph_index,
                        "level": int(getattr(paragraph, "level", 0)),
                        "text": paragraph.text or "",
                        "runs": runs,
                    }
                )
        except Exception:
            text_blocks = []
    payload["text_blocks"] = text_blocks
    if getattr(shape, "has_table", False):
        table_rows: List[List[str]] = []
        max_cols = 0
        for row in shape.table.rows:
            row_vals = [cell.text or "" for cell in row.cells]
            table_rows.append(row_vals)
            if len(row_vals) > max_cols:
                max_cols = len(row_vals)
        payload["table"] = {
            "row_count": len(table_rows),
            "col_count": max_cols,
            "rows": table_rows,
        }
    image_ref = _pptx_image_ref(shape)
    if image_ref is not None:
        payload["image_ref"] = image_ref
        payload["has_image"] = True
    else:
        payload["has_image"] = False
    return payload


def _pptx_run_payload(run, run_index: int) -> Dict[str, Any]:
    font = run.font
    size_pt = None
    font_name = None
    color = None
    if font is not None:
        try:
            if font.size is not None:
                size_pt = float(font.size.pt)
        except Exception:
            size_pt = None
        try:
            font_name = font.name
        except Exception:
            font_name = None
        try:
            if font.color is not None and font.color.rgb is not None:
                color = str(font.color.rgb)
        except Exception:
            color = None
    return {
        "index": run_index,
        "text": run.text or "",
        "bold": font.bold if font else None,
        "italic": font.italic if font else None,
        "underline": font.underline if font else None,
        "size_pt": size_pt,
        "font_name": font_name,
        "color": color,
    }


def _pptx_shape_type_name(shape: Any) -> str:
    shape_type = getattr(shape, "shape_type", None)
    if shape_type is None:
        return "UNKNOWN"
    shape_type_name = str(shape_type)
    if hasattr(shape_type, "name"):
        shape_type_name = shape_type.name
    return str(shape_type_name)


def _pptx_rotation_value(shape: Any) -> Optional[float]:
    try:
        rotation = getattr(shape, "rotation", None)
        if rotation is None:
            return None
        return float(rotation)
    except Exception:
        return None


def _pptx_image_ref(shape: Any) -> Optional[Dict[str, Any]]:
    try:
        image = getattr(shape, "image", None)
        if image is None:
            return None
        blob = image.blob
        return {
            "filename": getattr(image, "filename", None),
            "ext": getattr(image, "ext", None),
            "content_type": getattr(image, "content_type", None),
            "size_bytes": len(blob) if blob is not None else None,
            "sha1": _sha1_of_bytes(blob) if blob is not None else None,
        }
    except Exception:
        return None


def _pptx_positioned_slide_payload(slide, presentation) -> Dict[str, Any]:
    slide_width = _pptx_length_value(getattr(presentation, "slide_width", None))
    slide_height = _pptx_length_value(getattr(presentation, "slide_height", None))
    if slide_width is None:
        slide_width = 1
    if slide_height is None:
        slide_height = 1

    positioned_shapes: List[Dict[str, Any]] = []
    missing_features: List[str] = []
    font_families: Set[str] = set()
    for shape_index, shape in enumerate(slide.shapes):
        item = _pptx_positioned_shape_payload(shape, shape_index, slide_width, slide_height)
        if item is None:
            missing_features.append(f"shape_{shape_index}_unsupported")
            continue
        positioned_shapes.append(item)
        for run in item.get("text_runs", []) or []:
            font_name = str(run.get("font_name") or "").strip()
            if font_name:
                font_families.add(font_name)

    return {
        "coordinate_space": "slide_ratio",
        "slide_size": {"width": slide_width, "height": slide_height, "unit": "emu"},
        "positioned_shapes": positioned_shapes,
        "font_families": sorted(font_families),
        "missing_features": sorted(set(missing_features)),
    }


def _pptx_positioned_shape_payload(shape, shape_index: int, slide_width: int, slide_height: int) -> Optional[Dict[str, Any]]:
    left = _pptx_length_value(getattr(shape, "left", None))
    top = _pptx_length_value(getattr(shape, "top", None))
    width = _pptx_length_value(getattr(shape, "width", None))
    height = _pptx_length_value(getattr(shape, "height", None))
    if left is None or top is None or width is None or height is None:
        return None

    bounds = {
        "x": float(left) / float(slide_width) if slide_width > 0 else 0.0,
        "y": float(top) / float(slide_height) if slide_height > 0 else 0.0,
        "width": float(width) / float(slide_width) if slide_width > 0 else 0.0,
        "height": float(height) / float(slide_height) if slide_height > 0 else 0.0,
        "unit": "slide_ratio",
    }

    text_runs: List[Dict[str, Any]] = []
    text_blocks: List[Dict[str, Any]] = []
    plain_lines: List[str] = []
    if getattr(shape, "has_text_frame", False):
        try:
            for paragraph_index, paragraph in enumerate(shape.text_frame.paragraphs):
                paragraph_runs: List[Dict[str, Any]] = []
                for run_index, run in enumerate(paragraph.runs):
                    payload = _pptx_run_payload(run, run_index)
                    payload["paragraph_index"] = paragraph_index
                    paragraph_runs.append(payload)
                    text_runs.append(payload)
                block_text = paragraph.text or ""
                plain_lines.append(block_text)
                text_blocks.append(
                    {
                        "index": paragraph_index,
                        "level": int(getattr(paragraph, "level", 0)),
                        "text": block_text,
                        "runs": paragraph_runs,
                    }
                )
        except Exception:
            pass

    payload: Dict[str, Any] = {
        "index": shape_index,
        "z_index": shape_index,
        "name": getattr(shape, "name", None),
        "shape_type": _pptx_shape_type_name(shape),
        "rotation": _pptx_rotation_value(shape),
        "bounds": bounds,
        "frame_emu": {"x": left, "y": top, "width": width, "height": height},
        "x": float(left),
        "y": float(top),
        "w": float(width),
        "h": float(height),
        "text_blocks": text_blocks,
        "text_runs": text_runs,
        "text": "\n".join([line for line in plain_lines if str(line).strip()]),
    }
    if getattr(shape, "is_placeholder", False):
        try:
            payload["placeholder_type"] = str(shape.placeholder_format.type)
        except Exception:
            payload["placeholder_type"] = None
    image_ref = _pptx_image_ref(shape)
    if image_ref is not None:
        payload["image_ref"] = image_ref
        payload["has_image"] = True
    else:
        payload["has_image"] = False
    return payload


def _pptx_length_value(length_value: Any) -> Optional[int]:
    if length_value is None:
        return None
    try:
        return int(length_value)
    except Exception:
        return None


# ------------------ PDF ------------------

def pdf_extract_text(params: Dict[str, Any]) -> Dict[str, Any]:
    """Extract text from PDF, with optional chunk_info for page-range reads."""
    pypdf = import_module("pypdf")
    path = resolve_path(params)
    reader = pypdf.PdfReader(path)
    pages = params.get("pages")
    total_pages = len(reader.pages)
    text_parts: List[str] = []

    if pages:
        page_nums = [int(p) for p in pages]
        for page_num in page_nums:
            idx = page_num - 1
            if idx < 0 or idx >= total_pages:
                raise WorkerError("VALIDATION_FAILED", "invalid page index")
            text_parts.append(reader.pages[idx].extract_text() or "")

        total_chunks = max(1, (total_pages + PDF_CHUNK_PAGES - 1) // PDF_CHUNK_PAGES)
        min_page = min(page_nums)
        max_page = max(page_nums)
        chunk_index = (min_page - 1) // PDF_CHUNK_PAGES
        has_more = max_page < total_pages

        result: Dict[str, Any] = {"text": "\n".join(text_parts)}
        if total_chunks > 1:
            result["chunk_info"] = {
                "chunk_index": chunk_index,
                "total_chunks": total_chunks,
                "has_more": has_more,
                "range": f"{min_page}-{max_page}",
            }
        return result
    else:
        for page in reader.pages:
            text_parts.append(page.extract_text() or "")
        return {"text": "\n".join(text_parts)}


def pdf_get_info(params: Dict[str, Any]) -> Dict[str, Any]:
    """Get information about a PDF file including page count."""
    pypdf = import_module("pypdf")
    path = resolve_path(params)
    reader = pypdf.PdfReader(path)
    return {"page_count": len(reader.pages)}


def pdf_get_map(params: Dict[str, Any]) -> Dict[str, Any]:
    """Return structural map of a PDF file: page count, TOC, forms, annotations, chunks."""
    pypdf = import_module("pypdf")
    path = resolve_path(params)
    reader = pypdf.PdfReader(path)
    page_count = len(reader.pages)

    # Extract TOC from outlines/bookmarks
    toc: List[Dict[str, Any]] = []
    has_toc = False
    try:
        outlines = reader.outline
        if outlines:
            _extract_toc_entries(outlines, reader, toc)
            has_toc = len(toc) > 0
    except Exception:
        pass

    # Check for forms
    has_forms = False
    try:
        if reader.get_fields():
            has_forms = True
    except Exception:
        pass

    # Check for annotations
    has_annotations = False
    for page in reader.pages:
        annots = page.get("/Annots")
        if annots:
            has_annotations = True
            break

    # Page chunks
    chunks: List[Dict[str, Any]] = []
    chunk_idx = 0
    p = 1
    while p <= page_count:
        end_p = min(p + PDF_CHUNK_PAGES - 1, page_count)
        chunks.append({
            "index": chunk_idx,
            "pages": f"{p}-{end_p}",
        })
        chunk_idx += 1
        p = end_p + 1

    return {
        "page_count": page_count,
        "has_toc": has_toc,
        "toc": toc,
        "has_forms": has_forms,
        "has_annotations": has_annotations,
        "chunks": chunks,
    }


def _extract_toc_entries(outlines: Any, reader: Any, toc: List[Dict[str, Any]]) -> None:
    """Recursively extract TOC entries from PDF outlines."""
    if isinstance(outlines, list):
        for item in outlines:
            _extract_toc_entries(item, reader, toc)
    else:
        try:
            title = outlines.title if hasattr(outlines, 'title') else str(outlines)
            page_num = None
            if hasattr(outlines, 'page'):
                # Try to get page number
                try:
                    dest_page = outlines.page
                    if dest_page is not None:
                        page_idx = reader.pages.index(dest_page)
                        page_num = page_idx + 1
                except (ValueError, AttributeError):
                    pass
            toc.append({"title": title, "page": page_num})
        except Exception:
            pass


def pdf_render_page(params: Dict[str, Any]) -> Dict[str, Any]:
    path = resolve_path(params)
    page_index = int(params.get("page_index", 0))
    scale = clamp_scale(params.get("scale", 1.0))
    return render_pdf_page(path, page_index, scale)


def render_pdf_page(path: str, page_index: int, scale: float) -> Dict[str, Any]:
    fitz = import_module("fitz")
    doc = fitz.open(path)
    if page_index < 0 or page_index >= doc.page_count:
        raise WorkerError("VALIDATION_FAILED", "invalid page index")
    scaled_down = False
    png_bytes, scaled_down = _render_with_limits(doc, page_index, scale)
    return {
        "bytes_base64": base64.b64encode(png_bytes).decode("ascii"),
        "page_count": doc.page_count,
        "mime_type": "image/png",
        "scaled_down": scaled_down,
    }


def _render_with_limits(doc, page_index: int, scale: float) -> Tuple[bytes, bool]:
    fitz = import_module("fitz")
    scaled_down = False
    current_scale = scale
    for _ in range(3):
        mat = fitz.Matrix(current_scale, current_scale)
        pix = doc.load_page(page_index).get_pixmap(matrix=mat)
        png_bytes = pix.tobytes("png")
        if pix.width <= MAX_PREVIEW_DIM and pix.height <= MAX_PREVIEW_DIM and len(png_bytes) <= MAX_PREVIEW_BYTES:
            return png_bytes, scaled_down
        scaled_down = True
        scale_factor = min(
            MAX_PREVIEW_DIM / float(pix.width),
            MAX_PREVIEW_DIM / float(pix.height),
            1.0,
        )
        if len(png_bytes) > MAX_PREVIEW_BYTES:
            scale_factor = min(scale_factor, (MAX_PREVIEW_BYTES / float(len(png_bytes))) ** 0.5)
        new_scale = max(MIN_SCALE, current_scale * scale_factor)
        if new_scale >= current_scale:
            break
        current_scale = new_scale
    return png_bytes, True


def find_preview_renderer() -> str:
    if PREVIEW_RENDERER_PATH:
        if not os.path.isfile(PREVIEW_RENDERER_PATH):
            raise WorkerError("TOOL_WORKER_UNAVAILABLE", "preview renderer not found")
        if not os.access(PREVIEW_RENDERER_PATH, os.X_OK):
            raise WorkerError("TOOL_WORKER_UNAVAILABLE", "preview renderer not executable")
        log_debug("preview.renderer_override", path=PREVIEW_RENDERER_PATH)
        return PREVIEW_RENDERER_PATH
    soffice = shutil.which("soffice")
    if soffice is None:
        raise WorkerError("TOOL_WORKER_UNAVAILABLE", "LibreOffice not available")
    log_debug("preview.renderer_auto", path=soffice)
    return soffice


def convert_to_pdf(path: str) -> Tuple[str, str]:
    renderer = find_preview_renderer()
    tmp_dir = tempfile.mkdtemp(prefix="keenbench-preview-")
    cmd = [
        renderer,
        "--headless",
        "--convert-to",
        "pdf",
        "--outdir",
        tmp_dir,
        path,
    ]
    try:
        subprocess.run(cmd, stdout=subprocess.DEVNULL, stderr=subprocess.DEVNULL, check=True)
    except Exception:
        shutil.rmtree(tmp_dir, ignore_errors=True)
        raise WorkerError("FILE_READ_FAILED", "conversion failed")
    base = os.path.splitext(os.path.basename(path))[0]
    pdf_path = os.path.join(tmp_dir, base + ".pdf")
    if not os.path.exists(pdf_path):
        shutil.rmtree(tmp_dir, ignore_errors=True)
        raise WorkerError("FILE_READ_FAILED", "conversion failed")
    return pdf_path, tmp_dir


def render_docx_page_fallback(path: str, page_index: int, scale: float) -> Dict[str, Any]:
    docx = import_module("docx")
    doc = docx.Document(path)
    lines: List[str] = []
    for para in doc.paragraphs:
        text = (para.text or "").strip()
        if text:
            lines.append(text)
    for table in doc.tables:
        for row in table.rows:
            row_vals = []
            for cell in row.cells:
                row_vals.append((cell.text or "").strip())
            if any(val for val in row_vals):
                lines.append(" | ".join(row_vals))
    return render_text_page_preview(
        lines=lines,
        page_index=page_index,
        scale=scale,
        base_width=1240,
        base_height=1754,
        lines_per_page=36,
    )


def render_pptx_slide_fallback(path: str, slide_index: int, scale: float) -> Dict[str, Any]:
    pptx = import_module("pptx")
    prs = pptx.Presentation(path)
    slide_count = len(prs.slides)
    if slide_index < 0 or slide_index >= slide_count:
        raise WorkerError("VALIDATION_FAILED", "invalid slide_index")
    slide = prs.slides[slide_index]
    lines: List[str] = []
    title_shape = slide.shapes.title if slide.shapes.title else None
    if title_shape is not None:
        title = (title_shape.text or "").strip()
        if title:
            lines.append(title)
    for shape in slide.shapes:
        if title_shape is not None and shape == title_shape:
            continue
        text = getattr(shape, "text", "")
        if not text:
            continue
        for raw_line in str(text).splitlines():
            clean = raw_line.strip()
            if clean:
                lines.append(clean)
    render = render_text_page_preview(
        lines=lines,
        page_index=0,
        scale=scale,
        base_width=1280,
        base_height=720,
        lines_per_page=16,
    )
    return {
        "bytes_base64": render["bytes_base64"],
        "slide_count": slide_count,
        "mime_type": render["mime_type"],
        "scaled_down": render["scaled_down"],
    }


def render_text_page_preview(
    lines: List[str],
    page_index: int,
    scale: float,
    base_width: int,
    base_height: int,
    lines_per_page: int,
) -> Dict[str, Any]:
    import_module("PIL.Image")
    from PIL import Image, ImageDraw, ImageFont

    source_lines = [line for line in lines if line is not None and str(line).strip() != ""]
    if not source_lines:
        source_lines = ["(no preview text available)"]
    total_pages = max(1, (len(source_lines) + lines_per_page - 1) // lines_per_page)
    if page_index < 0 or page_index >= total_pages:
        raise WorkerError("VALIDATION_FAILED", "invalid page index")

    start = page_index * lines_per_page
    end = min(start + lines_per_page, len(source_lines))
    page_lines = [str(line) for line in source_lines[start:end]]

    width = max(320, int(base_width * scale))
    height = max(240, int(base_height * scale))
    scaled_down = False
    if width > MAX_PREVIEW_DIM or height > MAX_PREVIEW_DIM:
        ratio = min(MAX_PREVIEW_DIM / float(width), MAX_PREVIEW_DIM / float(height))
        width = max(1, int(width * ratio))
        height = max(1, int(height * ratio))
        scaled_down = True

    image = Image.new("RGB", (width, height), color=(255, 255, 255))
    draw = ImageDraw.Draw(image)
    try:
        font = ImageFont.truetype("DejaVuSans.ttf", max(14, int(18 * scale)))
    except Exception:
        font = ImageFont.load_default()

    margin_x = max(16, int(24 * scale))
    margin_y = max(16, int(24 * scale))
    max_text_width = max(10, width - (margin_x * 2))
    line_spacing = max(4, int(6 * scale))
    y = margin_y

    for line in page_lines:
        wrapped = _wrap_line_for_width(draw, line, font, max_text_width)
        for visual_line in wrapped:
            draw.text((margin_x, y), visual_line, fill=(34, 34, 34), font=font)
            bbox = draw.textbbox((margin_x, y), visual_line, font=font)
            y += (bbox[3] - bbox[1]) + line_spacing
            if y >= height - margin_y:
                break
        if y >= height - margin_y:
            break

    out = tempfile.SpooledTemporaryFile()
    image.save(out, format="PNG")
    out.seek(0)
    png_bytes = out.read()
    out.close()
    if len(png_bytes) > MAX_PREVIEW_BYTES:
        raise WorkerError("FILE_READ_FAILED", "preview too large")

    return {
        "bytes_base64": base64.b64encode(png_bytes).decode("ascii"),
        "page_count": total_pages,
        "mime_type": "image/png",
        "scaled_down": scaled_down,
    }


def _wrap_line_for_width(draw, line: str, font, max_width: int) -> List[str]:
    text = line.strip()
    if text == "":
        return [""]
    words = text.split()
    if not words:
        return [text]
    wrapped: List[str] = []
    current = words[0]
    for word in words[1:]:
        candidate = current + " " + word
        bbox = draw.textbbox((0, 0), candidate, font=font)
        if bbox[2] - bbox[0] <= max_width:
            current = candidate
            continue
        wrapped.append(current)
        current = word
    wrapped.append(current)
    return wrapped


# ------------------ IMAGES ------------------

def image_get_metadata(params: Dict[str, Any]) -> Dict[str, Any]:
    path = resolve_path(params)
    size_bytes = os.path.getsize(path)
    ext = os.path.splitext(path)[1].lower()
    if ext == ".svg":
        width, height = _read_svg_dimensions(path)
        return {
            "format": "svg",
            "width": width,
            "height": height,
            "size_bytes": size_bytes,
            "color_depth": None,
        }
    import_module("PIL.Image")
    from PIL import Image
    with Image.open(path) as img:
        width, height = img.size
        depth_map = {
            "1": 1,
            "L": 8,
            "P": 8,
            "RGB": 24,
            "RGBA": 32,
            "CMYK": 32,
            "I": 32,
            "F": 32,
        }
        depth = depth_map.get(img.mode)
        return {
            "format": (img.format or ext.replace(".", "")),
            "width": width,
            "height": height,
            "size_bytes": size_bytes,
            "color_depth": depth,
        }


def image_render(params: Dict[str, Any]) -> Dict[str, Any]:
    path = resolve_path(params)
    ext = os.path.splitext(path)[1].lower()
    if ext == ".svg":
        png_bytes = _render_svg(path)
        return {
            "bytes_base64": base64.b64encode(png_bytes).decode("ascii"),
            "mime_type": "image/png",
            "scaled_down": False,
        }
    with open(path, "rb") as f:
        data = f.read()
    return {
        "bytes_base64": base64.b64encode(data).decode("ascii"),
        "mime_type": _guess_mime_from_ext(ext),
        "scaled_down": False,
    }


def _render_svg(path: str) -> bytes:
    try:
        cairosvg = import_module("cairosvg")
        return cairosvg.svg2png(url=path)
    except WorkerError:
        raise
    except Exception:
        raise WorkerError("TOOL_WORKER_UNAVAILABLE", "svg renderer not available")


def _read_svg_dimensions(path: str) -> Tuple[Optional[int], Optional[int]]:
    try:
        import xml.etree.ElementTree as ET
        tree = ET.parse(path)
        root = tree.getroot()
        width = root.attrib.get("width")
        height = root.attrib.get("height")
        def to_int(val):
            if not val:
                return None
            for suffix in ["px", "pt"]:
                if val.endswith(suffix):
                    val = val[:-len(suffix)]
            try:
                return int(float(val))
            except Exception:
                return None
        return to_int(width), to_int(height)
    except Exception:
        return None, None


def _guess_mime_from_ext(ext: str) -> str:
    if ext == ".png":
        return "image/png"
    if ext == ".jpg" or ext == ".jpeg":
        return "image/jpeg"
    if ext == ".gif":
        return "image/gif"
    if ext == ".webp":
        return "image/webp"
    if ext == ".svg":
        return "image/svg+xml"
    return "application/octet-stream"


# ------------------ TEXT ------------------

def text_get_map(params: Dict[str, Any]) -> Dict[str, Any]:
    """Return structural map of a text file: line count, char count, chunks."""
    path = resolve_path(params)
    with open(path, "r", encoding="utf-8", errors="replace") as f:
        content = f.read()
    lines = content.split("\n")
    line_count = len(lines)
    char_count = len(content)

    # Chunks of TEXT_CHUNK_LINES lines
    chunks: List[Dict[str, Any]] = []
    chunk_idx = 0
    ln = 1
    while ln <= line_count:
        end_ln = min(ln + TEXT_CHUNK_LINES - 1, line_count)
        chunks.append({
            "index": chunk_idx,
            "lines": f"{ln}-{end_ln}",
        })
        chunk_idx += 1
        ln = end_ln + 1

    return {
        "line_count": line_count,
        "char_count": char_count,
        "chunks": chunks,
    }


def text_read_lines(params: Dict[str, Any]) -> Dict[str, Any]:
    """Read a range of lines from a text file with chunk_info."""
    path = resolve_path(params)
    with open(path, "r", encoding="utf-8", errors="replace") as f:
        all_lines = f.readlines()
    total_lines = len(all_lines)
    line_start = int(params.get("line_start", 1))
    line_count = int(params.get("line_count", TEXT_CHUNK_LINES))

    if line_start < 1:
        line_start = 1
    if line_start > total_lines:
        return {
            "text": "",
            "chunk_info": {
                "chunk_index": 0,
                "total_chunks": max(1, (total_lines + TEXT_CHUNK_LINES - 1) // TEXT_CHUNK_LINES),
                "has_more": False,
                "range": f"{line_start}-{line_start}",
            },
        }

    end_line = min(line_start + line_count - 1, total_lines)
    text = "".join(all_lines[line_start - 1 : end_line])

    total_chunks = max(1, (total_lines + TEXT_CHUNK_LINES - 1) // TEXT_CHUNK_LINES)
    chunk_index = (line_start - 1) // TEXT_CHUNK_LINES

    return {
        "text": text,
        "chunk_info": {
            "chunk_index": chunk_index,
            "total_chunks": total_chunks,
            "has_more": end_line < total_lines,
            "range": f"{line_start}-{end_line}",
        },
    }


# ------------------ TABULAR (CSV) ------------------

def _tabular_source_signature(path: str) -> Dict[str, Any]:
    stat = os.stat(path)
    hasher = hashlib.sha256()
    with open(path, "rb") as handle:
        while True:
            chunk = handle.read(1024 * 1024)
            if not chunk:
                break
            hasher.update(chunk)
    return {
        "size_bytes": int(stat.st_size),
        "mtime_ns": int(getattr(stat, "st_mtime_ns", int(stat.st_mtime * 1_000_000_000))),
        "sha256": hasher.hexdigest(),
    }


def _tabular_cache_key(rel_path: str) -> str:
    normalized = str(rel_path or "").strip().lower()
    return hashlib.sha256(normalized.encode("utf-8")).hexdigest()


def _tabular_cache_paths(workbench_id: str, rel_path: str) -> Tuple[str, str]:
    workbench_base = os.path.abspath(os.path.join(require_workbenches_dir(), workbench_id))
    tabular_dir = os.path.abspath(os.path.join(workbench_base, "meta", "tabular"))
    if not tabular_dir.startswith(workbench_base + os.sep):
        raise WorkerError("SANDBOX_VIOLATION", "invalid tabular cache path")
    key = _tabular_cache_key(rel_path)
    return (
        os.path.join(tabular_dir, f"{key}.duckdb"),
        os.path.join(tabular_dir, f"{key}.meta.json"),
    )


def _tabular_load_metadata(path: str) -> Optional[Dict[str, Any]]:
    if not os.path.isfile(path):
        return None
    try:
        with open(path, "r", encoding="utf-8") as handle:
            data = json.load(handle)
        if isinstance(data, dict):
            return data
    except Exception:
        return None
    return None


def _tabular_write_metadata(path: str, payload: Dict[str, Any]) -> None:
    os.makedirs(os.path.dirname(path), exist_ok=True)
    tmp_path = path + ".tmp"
    with open(tmp_path, "w", encoding="utf-8") as handle:
        json.dump(payload, handle, ensure_ascii=True, sort_keys=True)
    os.replace(tmp_path, path)


def _tabular_cache_is_fresh(
    metadata: Dict[str, Any],
    rel_path: str,
    source_sig: Dict[str, Any],
    db_path: str,
) -> bool:
    if metadata.get("version") != TABULAR_CACHE_VERSION:
        return False
    source = metadata.get("source")
    if not isinstance(source, dict):
        return False
    if source.get("path") != rel_path:
        return False
    try:
        size_bytes = int(source.get("size_bytes", -1))
        mtime_ns = int(source.get("mtime_ns", -1))
        sig_size = int(source_sig.get("size_bytes", -2))
        sig_mtime = int(source_sig.get("mtime_ns", -2))
    except Exception:
        return False
    if size_bytes != sig_size:
        return False
    if mtime_ns != sig_mtime:
        return False
    if str(source.get("sha256") or "") != str(source_sig.get("sha256") or ""):
        return False
    if not os.path.isfile(db_path):
        return False
    columns = metadata.get("columns")
    if not isinstance(columns, list):
        return False
    return True


def _tabular_remove_cache_files(db_path: str, metadata_path: str) -> None:
    for suffix in ("", ".wal"):
        target = db_path + suffix
        if os.path.exists(target):
            try:
                os.remove(target)
            except Exception:
                pass
    if os.path.exists(metadata_path):
        try:
            os.remove(metadata_path)
        except Exception:
            pass


def _tabular_can_decode(path: str, encoding: str) -> bool:
    import codecs

    try:
        decoder = codecs.getincrementaldecoder(encoding)(errors="strict")
        with open(path, "rb") as handle:
            while True:
                chunk = handle.read(1024 * 1024)
                if not chunk:
                    break
                decoder.decode(chunk)
        decoder.decode(b"", final=True)
        return True
    except Exception:
        return False


def _tabular_normalize_confidence(value: Any) -> float:
    try:
        conf = float(value)
    except Exception:
        return 0.0
    if conf > 1.0:
        conf = conf / 100.0
    if conf < 0:
        return 0.0
    if conf > 1:
        return 1.0
    return conf


def _tabular_detect_encoding(path: str) -> Tuple[str, float]:
    bom = b""
    try:
        with open(path, "rb") as handle:
            bom = handle.read(3)
    except Exception:
        bom = b""
    if bom.startswith(b"\xef\xbb\xbf"):
        try:
            if _tabular_can_decode(path, "utf-8-sig"):
                return "utf-8-sig", 1.0
        except Exception:
            pass
    try:
        if _tabular_can_decode(path, "utf-8"):
            return "utf-8", 1.0
    except Exception:
        pass
    try:
        if _tabular_can_decode(path, "utf-8-sig"):
            return "utf-8-sig", 1.0
    except Exception:
        pass

    charset_normalizer = import_module("charset_normalizer")
    try:
        matches = charset_normalizer.from_path(path)
        best = matches.best()
    except Exception as err:
        raise WorkerError("FILE_READ_FAILED", f"encoding detection failed: {err}")
    if best is None or not getattr(best, "encoding", None):
        raise WorkerError("FILE_READ_FAILED", "encoding detection failed")
    confidence = _tabular_normalize_confidence(getattr(best, "coherence", None))
    if confidence <= 0:
        confidence = _tabular_normalize_confidence(getattr(best, "percent_coherence", None))
    return str(best.encoding), confidence


def _tabular_open_connection(db_path: str, read_only: bool = False):
    duckdb = import_module("duckdb")
    try:
        conn = duckdb.connect(database=db_path, read_only=read_only)
    except TypeError:
        conn = duckdb.connect(db_path, read_only=read_only)
    _tabular_apply_connection_guardrails(conn)
    return conn


def _tabular_try_set_statement(conn: Any, statements: List[str]) -> bool:
    for statement in statements:
        try:
            conn.execute(statement)
            return True
        except Exception:
            continue
    return False


def _tabular_apply_connection_guardrails(conn: Any) -> None:
    memory_limit_mb = max(64, int(TABULAR_MEMORY_LIMIT_MB))
    max_threads = max(1, int(TABULAR_MAX_THREADS))
    timeout_ms = max(100, int(TABULAR_QUERY_TIMEOUT_MS))

    _tabular_try_set_statement(
        conn,
        [
            f"SET memory_limit='{memory_limit_mb}MB'",
            f"PRAGMA memory_limit='{memory_limit_mb}MB'",
        ],
    )
    _tabular_try_set_statement(
        conn,
        [
            f"SET threads={max_threads}",
            f"PRAGMA threads={max_threads}",
        ],
    )
    timeout_configured = _tabular_try_set_statement(
        conn,
        [
            f"SET statement_timeout='{timeout_ms}ms'",
            f"PRAGMA statement_timeout='{timeout_ms}ms'",
            f"SET max_execution_time='{timeout_ms}ms'",
            f"PRAGMA max_execution_time='{timeout_ms}ms'",
        ],
    )
    if not timeout_configured and not callable(getattr(conn, "interrupt", None)):
        log_error(
            "tabular.timeout_guardrail_unavailable",
            timeout_ms=timeout_ms,
        )


def _tabular_run_with_timeout(
    conn: Any,
    operation: Callable[[], Any],
    *,
    timeout_ms: Optional[int] = None,
    error_code: str = "FILE_READ_FAILED",
):
    timeout = int(timeout_ms if timeout_ms is not None else TABULAR_QUERY_TIMEOUT_MS)
    if timeout <= 0:
        return operation()

    done = threading.Event()
    timed_out = {"value": False}

    def _interrupt() -> None:
        if done.is_set():
            return
        timed_out["value"] = True
        interrupt = getattr(conn, "interrupt", None)
        if callable(interrupt):
            try:
                interrupt()
            except Exception:
                pass

    timer = threading.Timer(timeout / 1000.0, _interrupt)
    timer.daemon = True
    timer.start()
    try:
        return operation()
    except Exception:
        if timed_out["value"]:
            raise WorkerError(error_code, f"query timed out after {timeout}ms")
        raise
    finally:
        done.set()
        timer.cancel()


def _tabular_invalidate_cache_for_rel_path(workbench_id: str, rel_path: str) -> None:
    clean_rel_path = str(rel_path or "").strip()
    if not clean_rel_path or not clean_rel_path.lower().endswith(".csv"):
        return
    try:
        db_path, metadata_path = _tabular_cache_paths(workbench_id, clean_rel_path)
    except Exception:
        return
    _tabular_remove_cache_files(db_path, metadata_path)


def _tabular_normalize_key(value: Any) -> str:
    return "".join(ch for ch in str(value or "").lower() if ch.isalnum())


def _tabular_sniff_csv(conn: Any, source_path: str) -> Dict[str, Any]:
    try:
        cursor = conn.execute("SELECT * FROM sniff_csv(?, sample_size=-1)", [source_path])
        row = cursor.fetchone()
        if row is None:
            return {"delimiter": ",", "quote_char": '"', "has_header": True}
        values: Dict[str, Any] = {}
        for idx, desc in enumerate(cursor.description):
            if idx >= len(row):
                continue
            values[_tabular_normalize_key(desc[0])] = row[idx]
        has_header = values.get("hasheader")
        if has_header is None:
            has_header = values.get("header")
        delimiter = values.get("delimiter") or ","
        quote_char = values.get("quote") or '"'
        delimiter_text = str(delimiter)
        quote_text = str(quote_char)
        return {
            "delimiter": delimiter_text[0] if delimiter_text else ",",
            "quote_char": quote_text[0] if quote_text else '"',
            "has_header": bool(has_header) if has_header is not None else True,
        }
    except Exception:
        return {"delimiter": ",", "quote_char": '"', "has_header": True}


def _tabular_inferred_type(db_type: Any) -> str:
    t = str(db_type or "").strip().upper()
    if not t:
        return "string"
    integer_types = {
        "TINYINT",
        "SMALLINT",
        "INTEGER",
        "INT",
        "BIGINT",
        "HUGEINT",
        "UTINYINT",
        "USMALLINT",
        "UINTEGER",
        "UBIGINT",
    }
    if t in integer_types:
        return "integer"
    if any(x in t for x in ("DECIMAL", "NUMERIC", "DOUBLE", "FLOAT", "REAL")):
        return "float"
    if t in ("BOOLEAN", "BOOL"):
        return "boolean"
    if any(x in t for x in ("DATE", "TIME")):
        return "date"
    if any(x in t for x in ("CHAR", "TEXT", "STRING", "VARCHAR", "UUID", "JSON")):
        return "string"
    return "string"


def _tabular_collect_schema(conn: Any) -> List[Dict[str, Any]]:
    rows = conn.execute(f"PRAGMA table_info('{TABULAR_TABLE_NAME}')").fetchall()
    out: List[Dict[str, Any]] = []
    for idx, row in enumerate(rows):
        name = str(row[1])
        db_type = str(row[2] or "")
        out.append(
            {
                "name": name,
                "index": idx,
                "duckdb_type": db_type,
                "inferred_type": _tabular_inferred_type(db_type),
            }
        )
    return out


def _tabular_validate_source(params: Dict[str, Any]) -> Dict[str, Any]:
    workbench_id = require_workbench_id(params)
    root = (params.get("root") or "draft").strip()
    validate_root(root)
    rel_path = (params.get("path") or "").strip()
    validate_rel_path(rel_path)
    if not rel_path.lower().endswith(".csv"):
        raise WorkerError("VALIDATION_FAILED", "tabular tools only support .csv files")
    source_path = _resolve_rel_path(workbench_id, root, rel_path)
    if not os.path.isfile(source_path):
        raise WorkerError("FILE_READ_FAILED", f"file not found: {rel_path}")
    return {
        "workbench_id": workbench_id,
        "root": root,
        "rel_path": rel_path,
        "source_path": source_path,
    }


def _tabular_rebuild_cache(source: Dict[str, Any], db_path: str, metadata_path: str, source_sig: Dict[str, Any]) -> Dict[str, Any]:
    os.makedirs(os.path.dirname(db_path), exist_ok=True)
    _tabular_remove_cache_files(db_path, metadata_path)

    encoding, confidence = _tabular_detect_encoding(source["source_path"])
    conn = _tabular_open_connection(db_path, read_only=False)
    try:
        sniff = _tabular_sniff_csv(conn, source["source_path"])
        conn.execute(f"DROP TABLE IF EXISTS {TABULAR_TABLE_NAME}")
        conn.execute(
            f"CREATE TABLE {TABULAR_TABLE_NAME} AS "
            "SELECT * FROM read_csv_auto(?, sample_size=-1, ignore_errors=false, all_varchar=false, encoding=?, delim=?, quote=?, header=?)",
            [
                source["source_path"],
                encoding,
                str(sniff.get("delimiter", ","))[:1],
                str(sniff.get("quote_char", '"'))[:1],
                bool(sniff.get("has_header", True)),
            ],
        )
        row_count = int(conn.execute(f"SELECT COUNT(*) FROM {TABULAR_TABLE_NAME}").fetchone()[0] or 0)
        columns = _tabular_collect_schema(conn)
    except WorkerError:
        raise
    except Exception as err:
        raise WorkerError("FILE_READ_FAILED", f"failed to load csv: {err}")
    finally:
        conn.close()

    metadata = {
        "version": TABULAR_CACHE_VERSION,
        "source": {
            "root": source["root"],
            "path": source["rel_path"],
            "size_bytes": int(source_sig["size_bytes"]),
            "mtime_ns": int(source_sig["mtime_ns"]),
            "sha256": str(source_sig.get("sha256") or ""),
        },
        "format": "csv",
        "delimiter": sniff.get("delimiter", ","),
        "quote_char": sniff.get("quote_char", '"'),
        "has_header": bool(sniff.get("has_header", True)),
        "encoding_detected": encoding,
        "encoding_confidence": confidence,
        "row_count": row_count,
        "column_count": len(columns),
        "columns": columns,
    }
    _tabular_write_metadata(metadata_path, metadata)
    return metadata


def _tabular_ensure_cache(params: Dict[str, Any]) -> Dict[str, Any]:
    source = _tabular_validate_source(params)
    source_sig = _tabular_source_signature(source["source_path"])
    db_path, metadata_path = _tabular_cache_paths(source["workbench_id"], source["rel_path"])
    metadata = _tabular_load_metadata(metadata_path)
    if not isinstance(metadata, dict) or not _tabular_cache_is_fresh(metadata, source["rel_path"], source_sig, db_path):
        metadata = _tabular_rebuild_cache(source, db_path, metadata_path, source_sig)
    return {
        "db_path": db_path,
        "metadata_path": metadata_path,
        "metadata": metadata,
        "source": source,
    }


def _tabular_quote_ident(identifier: str) -> str:
    return '"' + identifier.replace('"', '""') + '"'


def _tabular_quote_literal(value: str) -> str:
    return "'" + value.replace("'", "''") + "'"


def _tabular_count_distinct_aggregates(conn: Any, columns: List[Dict[str, Any]]) -> Dict[str, Dict[str, int]]:
    if len(columns) == 0:
        return {}

    select_items: List[str] = []
    names: List[str] = []
    for index, col in enumerate(columns):
        name = str(col.get("name") or "")
        ident = _tabular_quote_ident(name)
        select_items.append(f"COUNT({ident}) AS non_null_{index}")
        select_items.append(f"COALESCE(approx_count_distinct({ident}), 0) AS distinct_{index}")
        names.append(name)

    row = conn.execute(f"SELECT {', '.join(select_items)} FROM {TABULAR_TABLE_NAME}").fetchone()
    values = list(row or [])

    out: Dict[str, Dict[str, int]] = {}
    for index, name in enumerate(names):
        value_index = index * 2
        non_null_count = int(values[value_index] or 0) if value_index < len(values) else 0
        distinct_estimate = int(values[value_index + 1] or 0) if value_index + 1 < len(values) else 0
        out[name] = {
            "non_null_count": non_null_count,
            "distinct_estimate": distinct_estimate,
        }
    return out


def _tabular_json_value(value: Any) -> Any:
    if value is None or isinstance(value, (bool, int, str)):
        return value
    if isinstance(value, float):
        return value if math.isfinite(value) else None
    if isinstance(value, decimal.Decimal):
        if not value.is_finite():
            return None
        as_int = value.to_integral_value()
        if value == as_int:
            return int(as_int)
        return float(value)
    if isinstance(value, (datetime.datetime, datetime.date, datetime.time)):
        return value.isoformat()
    if isinstance(value, bytes):
        return base64.b64encode(value).decode("ascii")
    if isinstance(value, (list, tuple)):
        return [_tabular_json_value(item) for item in value]
    if isinstance(value, dict):
        return {str(key): _tabular_json_value(val) for key, val in value.items()}
    item = getattr(value, "item", None)
    if callable(item):
        try:
            return _tabular_json_value(item())
        except Exception:
            pass
    return str(value)


def _tabular_chunks(total_rows: int, chunk_rows: int = TABULAR_CHUNK_ROWS) -> List[Dict[str, Any]]:
    chunks: List[Dict[str, Any]] = []
    if total_rows <= 0:
        return chunks
    start = 1
    index = 0
    while start <= total_rows:
        end = min(start + chunk_rows - 1, total_rows)
        chunks.append({"index": index, "rows": f"{start}-{end}"})
        index += 1
        start = end + 1
    return chunks


def _tabular_selected_columns(params: Dict[str, Any], schema: List[Dict[str, Any]]) -> List[Dict[str, Any]]:
    schema_by_name: Dict[str, Dict[str, Any]] = {}
    for col in schema:
        schema_by_name[str(col.get("name"))] = col

    requested = params.get("columns")
    if requested is None:
        return list(schema)
    if not isinstance(requested, list) or len(requested) == 0:
        raise WorkerError("VALIDATION_FAILED", "columns must be a non-empty list")

    selected: List[Dict[str, Any]] = []
    seen: Set[str] = set()
    for raw in requested:
        if not isinstance(raw, str) or not raw.strip():
            raise WorkerError("VALIDATION_FAILED", "invalid column name")
        name = raw.strip()
        if name in seen:
            continue
        col = schema_by_name.get(name)
        if col is None:
            raise WorkerError("VALIDATION_FAILED", f"unknown column: {name}")
        selected.append(col)
        seen.add(name)
    if len(selected) == 0:
        raise WorkerError("VALIDATION_FAILED", "missing columns")
    return selected


def _tabular_parse_positive_int(value: Any, name: str, default: int) -> int:
    if value is None:
        return default
    try:
        parsed = int(value)
    except Exception:
        raise WorkerError("VALIDATION_FAILED", f"invalid {name}")
    if parsed <= 0:
        raise WorkerError("VALIDATION_FAILED", f"invalid {name}")
    return parsed


def _tabular_parse_non_negative_int(value: Any, name: str, default: int) -> int:
    if value is None:
        return default
    try:
        parsed = int(value)
    except Exception:
        raise WorkerError("VALIDATION_FAILED", f"invalid {name}")
    if parsed < 0:
        raise WorkerError("VALIDATION_FAILED", f"invalid {name}")
    return parsed


def _tabular_strip_sql_literals_and_comments(sql: str) -> str:
    out: List[str] = []
    state = "normal"
    i = 0
    while i < len(sql):
        ch = sql[i]
        nxt = sql[i + 1] if i + 1 < len(sql) else ""
        if state == "normal":
            if ch == "'" and state == "normal":
                state = "single_quote"
                out.append(" ")
                i += 1
                continue
            if ch == '"' and state == "normal":
                state = "double_quote"
                out.append(" ")
                i += 1
                continue
            if ch == "-" and nxt == "-":
                state = "line_comment"
                i += 2
                continue
            if ch == "/" and nxt == "*":
                state = "block_comment"
                i += 2
                continue
            out.append(ch)
            i += 1
            continue
        if state == "single_quote":
            if ch == "'" and nxt == "'":
                i += 2
                continue
            if ch == "'":
                state = "normal"
            i += 1
            continue
        if state == "double_quote":
            if ch == '"' and nxt == '"':
                i += 2
                continue
            if ch == '"':
                state = "normal"
            i += 1
            continue
        if state == "line_comment":
            if ch == "\n":
                state = "normal"
                out.append("\n")
            i += 1
            continue
        if state == "block_comment":
            if ch == "*" and nxt == "/":
                state = "normal"
                i += 2
                continue
            if ch == "\n":
                out.append("\n")
            i += 1
            continue
    return "".join(out)


# Keywords that, as SQL statements, could modify data or escape the sandbox.
_DANGEROUS_KEYWORDS = frozenset({
    "insert", "update", "delete", "drop", "alter", "create", "replace",
    "truncate", "copy", "attach", "detach", "pragma", "set", "call",
    "vacuum", "analyze", "grant", "revoke", "comment", "install",
    "load", "export", "import", "merge",
})

# Tokens after which the next bare word is most likely an identifier (column
# name, table name, alias) rather than a SQL statement keyword.  This lets
# column names like "Comment" or "Update" pass through the blocklist.
# Defence-in-depth note: the DuckDB connection is opened read-only, so even
# if a dangerous keyword slips through here the database will reject writes.
_IDENT_CONTEXT_TOKENS = frozenset({
    # SQL clauses that introduce identifiers
    "select", "from", "join", "on", "as", "by", "where", "having",
    "and", "or", "not", "in", "between", "like", "ilike", "is",
    "when", "then", "else", "distinct", "all", "over", "partition",
    "using", "limit", "offset", "case", "end", "with", "recursive",
    "filter", "within", "respect", "ignore", "nulls", "asc", "desc",
    "group", "order", "into",
    # Punctuation that typically precedes identifiers / expressions
    ",", ".", "(", ")", "=", "<", ">", "!", "+", "-", "*", "/",
    "|", "&", "^", "~", "%",
})

_TOKEN_SPLIT_RE = re.compile(r"(\s+|[,;.()\[\]=<>!+\-*/|&^~%])")


def _tabular_find_dangerous_keyword(lowered_sql: str) -> Optional[str]:
    """Return a dangerous keyword found in *statement* position, or None.

    Splits the (already-stripped, lowered) SQL into tokens and skips keywords
    that appear right after a token that normally introduces an identifier
    (SELECT, FROM, comma, dot, operators, …).
    """
    parts = _TOKEN_SPLIT_RE.split(lowered_sql)
    tokens = [p for p in parts if p.strip()]
    prev = None
    for token in tokens:
        if token in _DANGEROUS_KEYWORDS:
            if prev is not None and prev in _IDENT_CONTEXT_TOKENS:
                prev = token
                continue
            return token
        prev = token
    return None


def _tabular_validate_query(raw_query: Any) -> str:
    if raw_query is None:
        raise WorkerError("VALIDATION_FAILED", "missing query")
    query = str(raw_query).strip()
    if not query:
        raise WorkerError("VALIDATION_FAILED", "missing query")

    checked = _tabular_strip_sql_literals_and_comments(query).strip()
    if not checked:
        raise WorkerError("VALIDATION_FAILED", "invalid query")

    semicolon_count = checked.count(";")
    if semicolon_count > 1:
        raise WorkerError("VALIDATION_FAILED", "query must be a single statement")
    if semicolon_count == 1:
        if not checked.endswith(";"):
            raise WorkerError("VALIDATION_FAILED", "query must be a single statement")
        if not query.rstrip().endswith(";"):
            raise WorkerError("VALIDATION_FAILED", "query must end with ';'")
        checked = checked[:-1].strip()
        query = query.rstrip()[:-1].strip()
    if not checked:
        raise WorkerError("VALIDATION_FAILED", "invalid query")

    first_token = checked.split(None, 1)[0].lower()
    if first_token not in ("select", "with"):
        raise WorkerError("VALIDATION_FAILED", "only SELECT/CTE queries are allowed")

    lowered = checked.lower()
    dangerous_kw = _tabular_find_dangerous_keyword(lowered)
    if dangerous_kw is not None:
        raise WorkerError("VALIDATION_FAILED", "query must be read-only")
    if re.search(
        r"\b(read_[a-z0-9_]+|[a-z0-9_]+_scan|glob|open_url|httpfs|read_csv|read_parquet|read_json)\s*\(",
        lowered,
    ):
        raise WorkerError("VALIDATION_FAILED", "path-reading functions are not allowed")

    return query


def _tabular_type_family(inferred_type: str) -> str:
    if inferred_type in ("integer", "float"):
        return "numeric"
    if inferred_type == "boolean":
        return "boolean"
    if inferred_type == "date":
        return "date"
    if inferred_type == "string":
        return "string"
    return "other"


def tabular_get_map(params: Dict[str, Any]) -> Dict[str, Any]:
    cache = _tabular_ensure_cache(params)
    metadata = cache["metadata"]
    row_count = int(metadata.get("row_count", 0))
    columns = metadata.get("columns") or []
    out_columns: List[Dict[str, Any]] = []
    for col in columns:
        out_columns.append(
            {
                "name": str(col.get("name") or ""),
                "index": int(col.get("index", 0)),
                "inferred_type": str(col.get("inferred_type") or "string"),
            }
        )
    return {
        "format": "csv",
        "delimiter": str(metadata.get("delimiter") or ","),
        "quote_char": str(metadata.get("quote_char") or '"'),
        "encoding_detected": str(metadata.get("encoding_detected") or ""),
        "encoding_confidence": _tabular_normalize_confidence(metadata.get("encoding_confidence")),
        "has_header": bool(metadata.get("has_header", True)),
        "row_count": row_count,
        "column_count": len(out_columns),
        "columns": out_columns,
        "chunks": _tabular_chunks(row_count, TABULAR_CHUNK_ROWS),
    }


def tabular_describe(params: Dict[str, Any]) -> Dict[str, Any]:
    cache = _tabular_ensure_cache(params)
    metadata = cache["metadata"]
    schema = metadata.get("columns") or []
    row_count = int(metadata.get("row_count", 0))

    conn = _tabular_open_connection(cache["db_path"], read_only=True)
    try:
        counts_by_name = _tabular_count_distinct_aggregates(conn, schema)
        columns_out: List[Dict[str, Any]] = []
        for col in schema:
            col_name = str(col.get("name") or "")
            counts = counts_by_name.get(col_name) or {}
            non_null_count = int(counts.get("non_null_count", 0))
            distinct_estimate = int(counts.get("distinct_estimate", 0))
            columns_out.append(
                {
                    "name": col_name,
                    "index": int(col.get("index", 0)),
                    "inferred_type": str(col.get("inferred_type") or "string"),
                    "nullable": non_null_count < row_count,
                    "non_null_count": non_null_count,
                    "distinct_estimate": distinct_estimate,
                }
            )
    finally:
        conn.close()

    return {
        "row_count": row_count,
        "column_count": len(schema),
        "columns": columns_out,
    }


def tabular_get_stats(params: Dict[str, Any]) -> Dict[str, Any]:
    cache = _tabular_ensure_cache(params)
    metadata = cache["metadata"]
    schema = metadata.get("columns") or []
    selected = _tabular_selected_columns(params, schema)
    row_count = int(metadata.get("row_count", 0))

    conn = _tabular_open_connection(cache["db_path"], read_only=True)
    try:
        counts_by_name = _tabular_count_distinct_aggregates(conn, selected)

        selected_meta: List[Dict[str, Any]] = []
        out_columns: List[Dict[str, Any]] = []
        numeric_positions: List[int] = []
        date_positions: List[int] = []
        boolean_positions: List[int] = []
        text_positions: List[int] = []

        for index, col in enumerate(selected):
            col_name = str(col.get("name") or "")
            inferred_type = str(col.get("inferred_type") or "string")
            counts = counts_by_name.get(col_name) or {}
            non_null_count = int(counts.get("non_null_count", 0))
            distinct_estimate = int(counts.get("distinct_estimate", 0))
            item: Dict[str, Any] = {
                "name": col_name,
                "type": inferred_type,
                "non_null_count": non_null_count,
                "distinct_estimate": distinct_estimate,
            }
            selected_meta.append(
                {
                    "index": index,
                    "name": col_name,
                    "ident": _tabular_quote_ident(col_name),
                    "family": _tabular_type_family(inferred_type),
                }
            )
            out_columns.append(item)

        for meta in selected_meta:
            if meta["family"] == "numeric":
                numeric_positions.append(meta["index"])
            elif meta["family"] == "date":
                date_positions.append(meta["index"])
            elif meta["family"] == "boolean":
                boolean_positions.append(meta["index"])
            else:
                text_positions.append(meta["index"])

        if len(numeric_positions) > 0:
            select_items: List[str] = []
            for index in numeric_positions:
                ident = str(selected_meta[index]["ident"])
                select_items.extend(
                    [
                        f"MIN({ident}) AS min_{index}",
                        f"MAX({ident}) AS max_{index}",
                        f"AVG({ident}) AS mean_{index}",
                        f"SUM({ident}) AS sum_{index}",
                        f"STDDEV_SAMP({ident}) AS stddev_{index}",
                    ]
                )
            row = conn.execute(f"SELECT {', '.join(select_items)} FROM {TABULAR_TABLE_NAME}").fetchone()
            values = list(row or [])
            for position, index in enumerate(numeric_positions):
                base = position * 5
                item = out_columns[index]
                item["min"] = _tabular_json_value(values[base]) if base < len(values) else None
                item["max"] = _tabular_json_value(values[base + 1]) if base + 1 < len(values) else None
                item["mean"] = _tabular_json_value(values[base + 2]) if base + 2 < len(values) else None
                item["sum"] = _tabular_json_value(values[base + 3]) if base + 3 < len(values) else None
                item["stddev"] = _tabular_json_value(values[base + 4]) if base + 4 < len(values) else None

        if len(date_positions) > 0:
            select_items = []
            for index in date_positions:
                ident = str(selected_meta[index]["ident"])
                select_items.extend(
                    [
                        f"MIN({ident}) AS min_{index}",
                        f"MAX({ident}) AS max_{index}",
                    ]
                )
            row = conn.execute(f"SELECT {', '.join(select_items)} FROM {TABULAR_TABLE_NAME}").fetchone()
            values = list(row or [])
            for position, index in enumerate(date_positions):
                base = position * 2
                item = out_columns[index]
                item["min"] = _tabular_json_value(values[base]) if base < len(values) else None
                item["max"] = _tabular_json_value(values[base + 1]) if base + 1 < len(values) else None

        if len(boolean_positions) > 0:
            select_items = []
            for index in boolean_positions:
                ident = str(selected_meta[index]["ident"])
                select_items.extend(
                    [
                        f"SUM(CASE WHEN {ident} IS TRUE THEN 1 ELSE 0 END) AS true_count_{index}",
                        f"SUM(CASE WHEN {ident} IS FALSE THEN 1 ELSE 0 END) AS false_count_{index}",
                    ]
                )
            row = conn.execute(f"SELECT {', '.join(select_items)} FROM {TABULAR_TABLE_NAME}").fetchone()
            values = list(row or [])
            for position, index in enumerate(boolean_positions):
                base = position * 2
                item = out_columns[index]
                item["true_count"] = int(values[base] or 0) if base < len(values) else 0
                item["false_count"] = int(values[base + 1] or 0) if base + 1 < len(values) else 0

        if len(text_positions) > 0:
            length_select_items: List[str] = []
            common_select_items: List[str] = []
            for index in text_positions:
                col_name = str(selected_meta[index]["name"])
                ident = str(selected_meta[index]["ident"])
                col_literal = _tabular_quote_literal(col_name)
                length_select_items.extend(
                    [
                        f"MIN(LENGTH(CAST({ident} AS VARCHAR))) AS min_length_{index}",
                        f"MAX(LENGTH(CAST({ident} AS VARCHAR))) AS max_length_{index}",
                    ]
                )
                common_select_items.append(
                    f"SELECT {col_literal} AS column_name, value, count FROM ("
                    f"SELECT CAST({ident} AS VARCHAR) AS value, COUNT(*) AS count "
                    f"FROM {TABULAR_TABLE_NAME} WHERE {ident} IS NOT NULL "
                    "GROUP BY 1 ORDER BY 2 DESC, 1 ASC LIMIT 5)"
                )

            length_row = conn.execute(f"SELECT {', '.join(length_select_items)} FROM {TABULAR_TABLE_NAME}").fetchone()
            length_values = list(length_row or [])
            for position, index in enumerate(text_positions):
                base = position * 2
                item = out_columns[index]
                item["min_length"] = _tabular_json_value(length_values[base]) if base < len(length_values) else None
                item["max_length"] = (
                    _tabular_json_value(length_values[base + 1]) if base + 1 < len(length_values) else None
                )
                item["most_common"] = []

            if len(common_select_items) > 0:
                common_rows = conn.execute(" UNION ALL ".join(common_select_items)).fetchall()
                text_index_by_name: Dict[str, int] = {}
                for index in text_positions:
                    text_index_by_name[str(selected_meta[index]["name"])] = index
                for row in common_rows:
                    if len(row) < 3:
                        continue
                    col_name = str(row[0] or "")
                    index = text_index_by_name.get(col_name)
                    if index is None:
                        continue
                    out_columns[index]["most_common"].append(
                        {
                            "value": _tabular_json_value(row[1]),
                            "count": int(row[2] or 0),
                        }
                    )
                for index in text_positions:
                    out_columns[index]["most_common"].sort(
                        key=lambda item: (
                            -int(item.get("count", 0)),
                            str(item.get("value") or ""),
                        )
                    )
    finally:
        conn.close()

    return {
        "row_count": row_count,
        "columns": out_columns,
    }


def tabular_read_rows(params: Dict[str, Any]) -> Dict[str, Any]:
    cache = _tabular_ensure_cache(params)
    metadata = cache["metadata"]
    schema = metadata.get("columns") or []
    selected = _tabular_selected_columns(params, schema)
    row_start = _tabular_parse_positive_int(params.get("row_start"), "row_start", 1)
    row_count_requested = _tabular_parse_positive_int(params.get("row_count"), "row_count", TABULAR_CHUNK_ROWS)
    offset = row_start - 1
    select_list = ", ".join(_tabular_quote_ident(str(col["name"])) for col in selected)

    conn = _tabular_open_connection(cache["db_path"], read_only=True)
    try:
        rows = conn.execute(
            f"SELECT {select_list} FROM {TABULAR_TABLE_NAME} ORDER BY rowid LIMIT ? OFFSET ?",
            [row_count_requested, offset],
        ).fetchall()
    finally:
        conn.close()

    row_values = [[_tabular_json_value(v) for v in row] for row in rows]
    total_rows = int(metadata.get("row_count", 0))
    return {
        "columns": [str(col.get("name") or "") for col in selected],
        "column_types": [str(col.get("inferred_type") or "string") for col in selected],
        "rows": row_values,
        "row_start": row_start,
        "row_count": len(row_values),
        "total_rows": total_rows,
        "has_more": (offset + len(row_values)) < total_rows,
    }


def tabular_query(params: Dict[str, Any]) -> Dict[str, Any]:
    cache = _tabular_ensure_cache(params)
    query = _tabular_validate_query(params.get("query"))
    window_rows = _tabular_parse_positive_int(params.get("window_rows"), "window_rows", TABULAR_CHUNK_ROWS)
    window_offset = _tabular_parse_non_negative_int(params.get("window_offset"), "window_offset", 0)
    total_column_ident = _tabular_quote_ident(TABULAR_TOTAL_COUNT_COLUMN)

    conn = _tabular_open_connection(cache["db_path"], read_only=True)
    start = time.time()
    try:
        def _read_window_with_total() -> Tuple[List[List[Any]], List[str], List[str], int]:
            cursor = conn.execute(
                f"SELECT q.*, COUNT(*) OVER() AS {total_column_ident} FROM ({query}) AS q LIMIT ? OFFSET ?",
                [window_rows, window_offset],
            )
            rows_local = cursor.fetchall()
            description = list(cursor.description or [])
            if len(description) == 0:
                return [], [], [], 0

            data_description = description[:-1]
            columns_local = [str(desc[0]) for desc in data_description]
            column_types_local = [
                _tabular_inferred_type(desc[1] if len(desc) > 1 else "")
                for desc in data_description
            ]

            total_local = 0
            trimmed_rows: List[List[Any]] = []
            for row in rows_local:
                values = list(row)
                if len(values) == 0:
                    trimmed_rows.append([])
                    continue
                total_local = int(values[-1] or 0)
                trimmed_rows.append(values[:-1])
            return trimmed_rows, columns_local, column_types_local, total_local

        rows, columns, column_types, total_row_count = _tabular_run_with_timeout(
            conn,
            _read_window_with_total,
            error_code="FILE_READ_FAILED",
        )

        if len(rows) == 0 and window_offset > 0:
            def _count_rows() -> int:
                row = conn.execute(f"SELECT COUNT(*) AS total_row_count FROM ({query}) AS q").fetchone()
                return int(row[0] or 0) if row else 0

            total_row_count = _tabular_run_with_timeout(conn, _count_rows, error_code="FILE_READ_FAILED")
    except WorkerError:
        raise
    except Exception as err:
        raise WorkerError("FILE_READ_FAILED", f"query failed: {err}")
    finally:
        conn.close()

    row_values = [[_tabular_json_value(v) for v in row] for row in rows]
    elapsed_ms = int((time.time() - start) * 1000)
    return {
        "columns": columns,
        "column_types": column_types,
        "rows": row_values,
        "row_count": len(row_values),
        "total_row_count": total_row_count,
        "window_rows": window_rows,
        "window_offset": window_offset,
        "has_more": (window_offset + len(row_values)) < total_row_count,
        "query_elapsed_ms": elapsed_ms,
    }


def _tabular_csv_cell_value(value: Any) -> Any:
    if value is None:
        return ""
    if isinstance(value, (datetime.datetime, datetime.date, datetime.time)):
        return value.isoformat()
    if isinstance(value, bytes):
        return base64.b64encode(value).decode("ascii")
    return value


def _tabular_normalize_sheet_name(value: Any) -> str:
    raw = str(value or "").strip() or "Sheet1"
    sanitized = raw.replace("\\", "_").replace("/", "_").replace("?", "_").replace("*", "_")
    sanitized = sanitized.replace("[", "_").replace("]", "_").replace(":", "_")
    return sanitized[:31] or "Sheet1"


def _tabular_export_csv(conn: Any, query: str, target_path: str) -> Tuple[int, int]:
    def _export() -> Tuple[int, int]:
        cursor = conn.execute(query)
        columns = [str(desc[0]) for desc in cursor.description]
        os.makedirs(os.path.dirname(target_path), exist_ok=True)
        row_count = 0
        with open(target_path, "w", encoding="utf-8", newline="") as handle:
            writer = csv.writer(handle)
            writer.writerow(columns)
            while True:
                batch = cursor.fetchmany(1000)
                if not batch:
                    break
                for row in batch:
                    writer.writerow([_tabular_csv_cell_value(value) for value in row])
                    row_count += 1
        return row_count, len(columns)

    return _tabular_run_with_timeout(conn, _export, error_code="FILE_WRITE_FAILED")


def _tabular_export_xlsx(conn: Any, query: str, target_path: str, sheet_name: str) -> Tuple[int, int]:
    openpyxl = import_module("openpyxl")
    def _export() -> Tuple[int, int]:
        cursor = conn.execute(query)
        columns = [str(desc[0]) for desc in cursor.description]
        wb = openpyxl.Workbook()
        ws = wb.active
        ws.title = _tabular_normalize_sheet_name(sheet_name)
        ws.append([_xlsx_sanitize_value(c) for c in columns])
        row_count = 0
        while True:
            batch = cursor.fetchmany(1000)
            if not batch:
                break
            for row in batch:
                ws.append([_xlsx_sanitize_value(_tabular_json_value(value)) for value in row])
                row_count += 1
        os.makedirs(os.path.dirname(target_path), exist_ok=True)
        wb.save(target_path)
        return row_count, len(columns)

    return _tabular_run_with_timeout(conn, _export, error_code="FILE_WRITE_FAILED")


def _tabular_export_query_and_warnings(params: Dict[str, Any]) -> Tuple[str, List[str]]:
    query = params.get("query")
    if query is None or not str(query).strip():
        return f"SELECT * FROM {TABULAR_TABLE_NAME} ORDER BY rowid", []
    export_query = _tabular_validate_query(query)
    warnings: List[str] = []
    normalized_for_order = _tabular_strip_sql_literals_and_comments(str(query)).lower()
    if not re.search(r"\border\s+by\b", normalized_for_order):
        warnings.append("query_has_no_order_by; output order may vary")
    return export_query, warnings


def _tabular_query_rows_for_export(conn: Any, query: str) -> Tuple[List[str], List[List[Any]]]:
    def _read_rows() -> Tuple[List[str], List[List[Any]]]:
        cursor = conn.execute(query)
        columns = [str(desc[0]) for desc in (cursor.description or [])]
        rows: List[List[Any]] = []
        while True:
            batch = cursor.fetchmany(1000)
            if not batch:
                break
            for row in batch:
                rows.append([_tabular_json_value(value) for value in row])
        return columns, rows

    return _tabular_run_with_timeout(conn, _read_rows, error_code="FILE_READ_FAILED")


def _xlsx_parse_anchor_cell(
    value: Any,
    coordinate_from_string: Any,
    column_index_from_string: Any,
    get_column_letter: Any,
) -> Tuple[int, int, str]:
    raw = str(value or "").strip() or "A1"
    try:
        col_letter, row_index = coordinate_from_string(raw.upper())
        col_index = int(column_index_from_string(str(col_letter).replace("$", "").upper()))
    except Exception:
        raise WorkerError("VALIDATION_FAILED", "invalid start_cell")
    if row_index < 1 or col_index < 1:
        raise WorkerError("VALIDATION_FAILED", "invalid start_cell")
    normalized = f"{get_column_letter(col_index)}{row_index}"
    return row_index, col_index, normalized


def _xlsx_written_range(
    start_row: int,
    start_col: int,
    row_count: int,
    column_count: int,
    get_column_letter: Any,
) -> str:
    if row_count <= 0 or column_count <= 0:
        return ""
    start_ref = f"{get_column_letter(start_col)}{start_row}"
    end_ref = f"{get_column_letter(start_col + column_count - 1)}{start_row + row_count - 1}"
    if start_ref == end_ref:
        return start_ref
    return f"{start_ref}:{end_ref}"


def _xlsx_sheet_has_data(ws: Any) -> bool:
    max_row = int(getattr(ws, "max_row", 0) or 0)
    max_col = int(getattr(ws, "max_column", 0) or 0)
    if max_row <= 0 or max_col <= 0:
        return False
    for row in ws.iter_rows(min_row=1, max_row=max_row, min_col=1, max_col=max_col):
        for cell in row:
            if cell.value is not None:
                return True
    return False


def _xlsx_clear_sheet_contents(ws: Any) -> None:
    try:
        merged_ranges = list(ws.merged_cells.ranges)
    except Exception:
        merged_ranges = []
    for merged in merged_ranges:
        try:
            ws.unmerge_cells(str(merged))
        except Exception:
            continue
    max_row = int(getattr(ws, "max_row", 0) or 0)
    max_col = int(getattr(ws, "max_column", 0) or 0)
    if max_row > 0:
        ws.delete_rows(1, max_row)
    if max_col > 0:
        ws.delete_cols(1, max_col)


def _xlsx_clear_rectangle(
    ws: Any,
    start_row: int,
    start_col: int,
    row_count: int,
    column_count: int,
) -> None:
    if row_count <= 0 or column_count <= 0:
        return
    end_row = start_row + row_count - 1
    end_col = start_col + column_count - 1
    for row in ws.iter_rows(min_row=start_row, max_row=end_row, min_col=start_col, max_col=end_col):
        for cell in row:
            cell.value = None


def _xlsx_write_tabular_block(
    ws: Any,
    start_row: int,
    start_col: int,
    columns: List[str],
    rows: List[List[Any]],
    include_header: bool,
) -> int:
    current_row = start_row
    written_rows = 0
    if include_header:
        for index, column_name in enumerate(columns):
            ws.cell(row=current_row, column=start_col + index, value=_xlsx_sanitize_value(column_name))
        current_row += 1
        written_rows += 1
    for values in rows:
        for index, value in enumerate(values):
            ws.cell(row=current_row, column=start_col + index, value=_xlsx_sanitize_value(value))
        current_row += 1
        written_rows += 1
    return written_rows


def _xlsx_save_workbook_atomic(wb: Any, target_path: str) -> None:
    target_dir = os.path.dirname(target_path) or "."
    os.makedirs(target_dir, exist_ok=True)
    fd, tmp_path = tempfile.mkstemp(prefix=".keenbench-xlsx-", suffix=".tmp.xlsx", dir=target_dir)
    os.close(fd)
    try:
        wb.save(tmp_path)
        os.replace(tmp_path, target_path)
    except Exception as err:
        try:
            os.remove(tmp_path)
        except Exception:
            pass
        raise WorkerError("FILE_WRITE_FAILED", f"failed to write xlsx: {err}")


def tabular_export(params: Dict[str, Any]) -> Dict[str, Any]:
    cache = _tabular_ensure_cache(params)
    workbench_id = str(cache["source"]["workbench_id"])
    target_path = resolve_named_path(params, "target_path", "target_root", "draft", require_write=True)
    target_path_rel = (params.get("target_path") or "").strip()
    target_root = (params.get("target_root") or "draft").strip()
    format_name = str(params.get("format") or "").strip().lower()
    if format_name not in ("csv", "xlsx"):
        raise WorkerError("VALIDATION_FAILED", "unsupported format")
    expected_ext = "." + format_name
    if not target_path_rel.lower().endswith(expected_ext):
        raise WorkerError("VALIDATION_FAILED", "target_path extension must match format")
    export_query, warnings = _tabular_export_query_and_warnings(params)
    sheet_name = _tabular_normalize_sheet_name(params.get("sheet"))

    conn = _tabular_open_connection(cache["db_path"], read_only=True)
    try:
        if format_name == "csv":
            row_count, column_count = _tabular_export_csv(conn, export_query, target_path)
        else:
            row_count, column_count = _tabular_export_xlsx(conn, export_query, target_path, sheet_name)
    except WorkerError:
        raise
    except Exception as err:
        raise WorkerError("FILE_WRITE_FAILED", f"export failed: {err}")
    finally:
        conn.close()

    if format_name == "csv" and target_root == "draft":
        _tabular_invalidate_cache_for_rel_path(workbench_id, target_path_rel)

    result = {
        "target_path": target_path_rel,
        "format": format_name,
        "row_count": row_count,
        "column_count": column_count,
        "warnings": warnings,
    }
    if format_name == "xlsx":
        result["sheet"] = sheet_name
    return result


def tabular_update_from_export(params: Dict[str, Any]) -> Dict[str, Any]:
    openpyxl = import_module("openpyxl")
    from openpyxl.utils import get_column_letter
    from openpyxl.utils.cell import coordinate_from_string, column_index_from_string

    cache = _tabular_ensure_cache(params)
    target_path = resolve_named_path(params, "target_path", "target_root", "draft", require_write=True)
    target_path_rel = (params.get("target_path") or "").strip()
    if not target_path_rel.lower().endswith(".xlsx"):
        raise WorkerError("VALIDATION_FAILED", "target_path extension must be .xlsx")

    raw_sheet_name = str(params.get("sheet") or "").strip()
    if not raw_sheet_name:
        raise WorkerError("VALIDATION_FAILED", "missing sheet")
    sheet_name = _tabular_normalize_sheet_name(raw_sheet_name)

    mode = str(params.get("mode") or "").strip().lower()
    if mode not in ("replace_sheet", "append_rows", "write_range"):
        raise WorkerError("VALIDATION_FAILED", "invalid mode")

    start_cell = str(params.get("start_cell") or "A1").strip() or "A1"
    include_header = _coerce_bool_strict(params.get("include_header"), True)
    create_workbook_if_missing = _coerce_bool_strict(params.get("create_workbook_if_missing"), True)
    create_sheet_if_missing = _coerce_bool_strict(params.get("create_sheet_if_missing"), True)
    clear_target_range = _coerce_bool_strict(params.get("clear_target_range"), False)
    if clear_target_range and mode != "write_range":
        raise WorkerError("VALIDATION_FAILED", "clear_target_range is only supported for write_range mode")

    export_query, warnings = _tabular_export_query_and_warnings(params)

    conn = _tabular_open_connection(cache["db_path"], read_only=True)
    try:
        columns, rows = _tabular_query_rows_for_export(conn, export_query)
    except WorkerError:
        raise
    except Exception as err:
        raise WorkerError("FILE_READ_FAILED", f"query failed: {err}")
    finally:
        conn.close()

    workbook_exists = os.path.isfile(target_path)
    if workbook_exists:
        try:
            wb = openpyxl.load_workbook(target_path)
        except WorkerError:
            raise
        except Exception as err:
            raise WorkerError("FILE_READ_FAILED", f"failed to read xlsx: {err}")
    else:
        if not create_workbook_if_missing:
            raise WorkerError("VALIDATION_FAILED", "target workbook not found")
        wb = openpyxl.Workbook()

    ws = wb[sheet_name] if sheet_name in wb.sheetnames else None
    if ws is None:
        if not create_sheet_if_missing:
            raise WorkerError("VALIDATION_FAILED", f"sheet not found: {sheet_name}")
        if len(wb.sheetnames) == 1:
            default_ws = wb[wb.sheetnames[0]]
            if not _xlsx_sheet_has_data(default_ws):
                try:
                    default_ws.title = sheet_name
                    ws = default_ws
                except Exception:
                    ws = None
        if ws is None:
            ws = wb.create_sheet(sheet_name)

    write_start_row = 1
    write_start_col = 1
    include_header_for_write = include_header
    if mode == "replace_sheet":
        _xlsx_clear_sheet_contents(ws)
    elif mode == "append_rows":
        sheet_has_data = _xlsx_sheet_has_data(ws)
        write_start_row = int(getattr(ws, "max_row", 0) or 0) + 1 if sheet_has_data else 1
        if sheet_has_data and include_header:
            include_header_for_write = False
            warnings.append("header_skipped_on_append; sheet already has data")
    else:
        write_start_row, write_start_col, _ = _xlsx_parse_anchor_cell(
            start_cell,
            coordinate_from_string,
            column_index_from_string,
            get_column_letter,
        )
        if clear_target_range:
            total_write_rows = len(rows) + (1 if include_header_for_write else 0)
            _xlsx_clear_rectangle(ws, write_start_row, write_start_col, total_write_rows, len(columns))

    try:
        written_rows = _xlsx_write_tabular_block(
            ws,
            write_start_row,
            write_start_col,
            columns,
            rows,
            include_header_for_write,
        )
        written_range = _xlsx_written_range(
            write_start_row,
            write_start_col,
            written_rows,
            len(columns),
            get_column_letter,
        )
        _xlsx_save_workbook_atomic(wb, target_path)
    except WorkerError:
        raise
    except Exception as err:
        raise WorkerError("FILE_WRITE_FAILED", f"failed to write xlsx: {err}")

    return {
        "target_path": target_path_rel,
        "sheet": sheet_name,
        "mode": mode,
        "row_count": len(rows),
        "column_count": len(columns),
        "written_range": written_range,
        "warnings": warnings,
    }


# ------------------ INFO ------------------

def worker_get_info(params: Dict[str, Any]) -> Dict[str, Any]:
    return {
        "ok": True,
        "worker": "pyworker",
        "version": WORKER_VERSION,
    }


METHODS = {
    "WorkerGetInfo": (worker_get_info, "read"),
    "DocxApplyOps": (docx_apply_ops, "write"),
    "DocxGetStyles": (docx_get_styles, "read"),
    "DocxCopyAssets": (docx_copy_assets, "write"),
    "XlsxApplyOps": (xlsx_apply_ops, "write"),
    "XlsxGetStyles": (xlsx_get_styles, "read"),
    "XlsxCopyAssets": (xlsx_copy_assets, "write"),
    "PptxApplyOps": (pptx_apply_ops, "write"),
    "PptxGetStyles": (pptx_get_styles, "read"),
    "PptxCopyAssets": (pptx_copy_assets, "write"),
    "DocxExtractText": (docx_extract_text, "read"),
    "DocxGetSectionContent": (docx_get_section_content, "read"),
    "OdtExtractText": (odt_extract_text, "read"),
    "XlsxExtractText": (xlsx_extract_text, "read"),
    "XlsxGetInfo": (xlsx_get_info, "read"),
    "XlsxReadRange": (xlsx_read_range, "read"),
    "PptxExtractText": (pptx_extract_text, "read"),
    "PptxGetSlideContent": (pptx_get_slide_content, "read"),
    "PdfExtractText": (pdf_extract_text, "read"),
    "PdfGetInfo": (pdf_get_info, "read"),
    "ImageGetMetadata": (image_get_metadata, "read"),
    "PdfRenderPage": (pdf_render_page, "read"),
    "DocxRenderPage": (docx_render_page, "read"),
    "OdtRenderPage": (odt_render_page, "read"),
    "PptxRenderSlide": (pptx_render_slide, "read"),
    "XlsxRenderGrid": (xlsx_render_grid, "read"),
    "ImageRender": (image_render, "read"),
    # File map methods (M4)
    "XlsxGetMap": (xlsx_get_map, "read"),
    "DocxGetMap": (docx_get_map, "read"),
    "PptxGetMap": (pptx_get_map, "read"),
    "PdfGetMap": (pdf_get_map, "read"),
    "TextGetMap": (text_get_map, "read"),
    "TextReadLines": (text_read_lines, "read"),
    "TabularGetMap": (tabular_get_map, "read"),
    "TabularDescribe": (tabular_describe, "read"),
    "TabularGetStats": (tabular_get_stats, "read"),
    "TabularReadRows": (tabular_read_rows, "read"),
    "TabularQuery": (tabular_query, "read"),
    "TabularExport": (tabular_export, "write"),
    "TabularUpdateFromExport": (tabular_update_from_export, "write"),
}


def _extract_context(params: Dict[str, Any]) -> Dict[str, Any]:
    ctx: Dict[str, Any] = {}
    for key in [
        "workbench_id",
        "root",
        "path",
        "page_index",
        "scale",
        "slide_index",
        "sheet",
        "row_start",
        "row_count",
        "col_start",
        "col_count",
        "range",
        "section",
        "section_index",
        "pages",
        "line_start",
        "line_count",
        "query",
        "window_rows",
        "window_offset",
        "columns",
        "format",
        "source_path",
        "target_path",
        "source_root",
        "target_root",
        "detail",
    ]:
        if key in params:
            ctx[key] = params.get(key)
    assets = params.get("assets")
    if isinstance(assets, list):
        ctx["assets_count"] = len(assets)
    ops = params.get("ops")
    if isinstance(ops, list):
        ctx["ops_count"] = len(ops)
        op_names = []
        for op in ops[:5]:
            if isinstance(op, dict) and op.get("op"):
                op_names.append(str(op.get("op")))
        if op_names:
            ctx["ops_kinds"] = op_names
        if len(ops) > 5:
            ctx["ops_truncated"] = True
    return ctx


def handle_request(line: str) -> None:
    try:
        req = json.loads(line)
    except Exception as err:
        log_debug("request.invalid_json", error=str(err), size_bytes=len(line))
        return
    req_id = req.get("id")
    if req.get("jsonrpc") != JSONRPC_VERSION or req_id is None:
        log_debug(
            "request.invalid_jsonrpc",
            id=req_id,
            jsonrpc=req.get("jsonrpc"),
            method=req.get("method"),
        )
        return
    method = req.get("method")
    if method not in METHODS:
        log_error("request.unknown_method", id=req_id, method=method)
        send_error(req_id, "VALIDATION_FAILED", f"unknown method: {method}")
        return
    params = req.get("params") or {}
    func, mode = METHODS[method]
    ctx = _extract_context(params)
    start = time.time()
    log_debug("request.start", id=req_id, method=method, mode=mode, **ctx)
    try:
        if mode == "write":
            root = (params.get("root") or "draft").strip()
            validate_write_root(root)
        result = func(params)
        send_result(req_id, result)
        duration_ms = int((time.time() - start) * 1000)
        log_info("request.ok", id=req_id, method=method, duration_ms=duration_ms, **ctx)
    except WorkerError as err:
        duration_ms = int((time.time() - start) * 1000)
        log_error(
            "request.error",
            id=req_id,
            method=method,
            error_code=err.code,
            error=err.message,
            duration_ms=duration_ms,
            **ctx,
        )
        send_error(req_id, err.code, err.message)
    except Exception as err:
        code = "FILE_WRITE_FAILED" if mode == "write" else "FILE_READ_FAILED"
        duration_ms = int((time.time() - start) * 1000)
        fields = {
            "id": req_id,
            "method": method,
            "error_code": code,
            "error": str(err),
            "duration_ms": duration_ms,
        }
        fields.update(ctx)
        if DEBUG_ENABLED:
            fields["traceback"] = traceback.format_exc()
        log_error("request.exception", **fields)
        send_error(req_id, code, str(err))


def main() -> None:
    init_config()
    for line in sys.stdin:
        if not line.strip():
            continue
        handle_request(line)


if __name__ == "__main__":
    main()
