#!/usr/bin/env python3
"""Generate real test fixture files for KeenBench.

Replaces plaintext stubs in engine/testdata/office/ with real office files,
creates synthetic test data in engine/testdata/synthetic/, and syncs copies
to app/integration_test/support/office/.

Usage:
    engine/tools/pyworker/.venv/bin/python scripts/testdata/generate_fixtures.py
"""

import os
import shutil
import subprocess
import sys

from docx import Document
from docx.shared import Inches, Pt
from openpyxl import Workbook
from pptx import Presentation
from pptx.util import Inches as PptxInches, Pt as PptxPt
from odf.opendocument import OpenDocumentText
from odf.text import H, P
from PIL import Image, ImageDraw, ImageFont

ROOT = os.path.abspath(os.path.join(os.path.dirname(__file__), "..", ".."))
OFFICE_DIR = os.path.join(ROOT, "engine", "testdata", "office")
SYNTHETIC_DIR = os.path.join(ROOT, "engine", "testdata", "synthetic")
SUPPORT_OFFICE_DIR = os.path.join(
    ROOT, "app", "integration_test", "support", "office"
)


def generate_simple_docx():
    """Create simple.docx with headings and paragraphs."""
    doc = Document()
    doc.add_heading("Introduction", level=1)
    doc.add_paragraph(
        "This quarterly review covers the key achievements and challenges "
        "encountered during the period. The team made significant progress "
        "on several fronts, including product development, customer "
        "acquisition, and operational efficiency improvements."
    )
    doc.add_paragraph(
        "Revenue grew by 15% compared to the previous quarter, driven "
        "primarily by expansion into new markets. Customer satisfaction "
        "scores remained above target at 4.2 out of 5.0."
    )
    doc.add_heading("Details", level=1)
    doc.add_paragraph(
        "The engineering team completed three major releases during the "
        "quarter, each delivering substantial new functionality. The most "
        "notable release introduced automated reporting, which reduced "
        "manual processing time by approximately 40%."
    )
    path = os.path.join(OFFICE_DIR, "simple.docx")
    doc.save(path)
    print(f"  Created {path}")


def generate_multi_sheet_xlsx():
    """Create multi-sheet.xlsx with Revenue, Expenses, Summary sheets."""
    wb = Workbook()

    # Revenue sheet
    ws_rev = wb.active
    ws_rev.title = "Revenue"
    ws_rev.append(["Month", "Amount"])
    revenue_data = [
        ("January", 12500),
        ("February", 13200),
        ("March", 11800),
        ("April", 14100),
        ("May", 15300),
        ("June", 14700),
    ]
    for row in revenue_data:
        ws_rev.append(row)

    # Expenses sheet
    ws_exp = wb.create_sheet("Expenses")
    ws_exp.append(["Month", "Category", "Amount"])
    expense_data = [
        ("January", "Salaries", 8000),
        ("January", "Rent", 2500),
        ("February", "Salaries", 8000),
        ("February", "Rent", 2500),
        ("March", "Salaries", 8200),
        ("March", "Rent", 2500),
        ("April", "Salaries", 8200),
        ("April", "Marketing", 3100),
    ]
    for row in expense_data:
        ws_exp.append(row)

    # Summary sheet
    ws_sum = wb.create_sheet("Summary")
    ws_sum.append(["Category", "Total"])
    ws_sum.append(["Revenue", 81600])
    ws_sum.append(["Expenses", 43000])
    ws_sum.append(["Net", 38600])

    path = os.path.join(OFFICE_DIR, "multi-sheet.xlsx")
    wb.save(path)
    print(f"  Created {path}")


def generate_slides_pptx():
    """Create slides.pptx with 3 slides."""
    prs = Presentation()

    # Slide 1: Title slide
    slide_layout = prs.slide_layouts[0]  # Title Slide
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = "Q4 Review"
    slide.placeholders[1].text = "Annual Performance Summary"

    # Slide 2: Content with bullet points
    slide_layout = prs.slide_layouts[1]  # Title and Content
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = "Key Highlights"
    body = slide.placeholders[1]
    tf = body.text_frame
    tf.text = "Revenue increased 15% quarter-over-quarter"
    for bullet in [
        "Customer base grew to 1,200 active accounts",
        "Three major product releases shipped on schedule",
        "Employee satisfaction score: 4.3/5.0",
    ]:
        p = tf.add_paragraph()
        p.text = bullet

    # Slide 3: Blank with text box
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)
    from pptx.util import Emu

    txBox = slide.shapes.add_textbox(
        PptxInches(1), PptxInches(2), PptxInches(6), PptxInches(2)
    )
    tf = txBox.text_frame
    tf.text = (
        "Thank you for reviewing. For questions, contact the analytics "
        "team at analytics@example.com."
    )

    path = os.path.join(OFFICE_DIR, "slides.pptx")
    prs.save(path)
    print(f"  Created {path}")


def generate_report_pdf():
    """Create report.pdf via DOCX → LibreOffice conversion."""
    doc = Document()
    doc.add_heading("Annual Report", level=1)
    doc.add_paragraph(
        "This annual report provides a comprehensive overview of the "
        "organization's performance over the past fiscal year. It covers "
        "financial results, operational highlights, and strategic "
        "initiatives undertaken during the period."
    )
    doc.add_paragraph(
        "The organization achieved record revenue of 4.2 million EUR, "
        "representing a 12% increase over the prior year. Operating costs "
        "were kept in check at 3.1 million EUR, resulting in a healthy "
        "operating margin of 26%."
    )
    doc.add_heading("Financial Summary", level=2)

    table = doc.add_table(rows=4, cols=3)
    table.style = "Table Grid"
    headers = ["Category", "Current Year", "Prior Year"]
    for i, h in enumerate(headers):
        table.rows[0].cells[i].text = h
    data = [
        ("Revenue", "4,200,000", "3,750,000"),
        ("Expenses", "3,100,000", "2,900,000"),
        ("Net Income", "1,100,000", "850,000"),
    ]
    for r, row_data in enumerate(data, 1):
        for c, val in enumerate(row_data):
            table.rows[r].cells[c].text = val

    doc.add_heading("Outlook", level=2)
    doc.add_paragraph(
        "Looking ahead, the organization plans to invest significantly in "
        "technology infrastructure and talent acquisition. The projected "
        "revenue target for the next fiscal year is 4.8 million EUR, which "
        "represents a 14% growth target."
    )
    doc.add_paragraph(
        "Key strategic priorities include expanding into two new geographic "
        "markets, launching three new product lines, and achieving ISO 27001 "
        "certification for information security management."
    )

    # Use OFFICE_DIR for temp docx (snap LibreOffice can't access /tmp)
    docx_path = os.path.join(OFFICE_DIR, "_report_tmp.docx")
    pdf_dst = os.path.join(OFFICE_DIR, "report.pdf")
    try:
        doc.save(docx_path)
        result = subprocess.run(
            [
                "/snap/bin/libreoffice",
                "--headless",
                "--convert-to",
                "pdf",
                "--outdir",
                OFFICE_DIR,
                docx_path,
            ],
            capture_output=True,
            text=True,
            timeout=60,
        )
        # LibreOffice exits 0 even on failure; check for actual PDF
        tmp_pdf = os.path.join(OFFICE_DIR, "_report_tmp.pdf")
        if os.path.exists(tmp_pdf):
            os.replace(tmp_pdf, pdf_dst)
            print(f"  Created {pdf_dst}")
        else:
            print(f"  ERROR: LibreOffice conversion failed: {result.stderr}")
            sys.exit(1)
    finally:
        if os.path.exists(docx_path):
            os.remove(docx_path)


def generate_notes_odt():
    """Create notes.odt with meeting notes content."""
    doc = OpenDocumentText()

    h = H(outlinelevel=1, text="Meeting Notes")
    doc.text.addElement(h)

    p1 = P(
        text="The project kickoff meeting was held on Monday at 10:00 AM. "
        "All team leads were present, and the agenda covered sprint planning, "
        "resource allocation, and timeline review for the upcoming quarter."
    )
    doc.text.addElement(p1)

    p2 = P(
        text="Action items were assigned as follows: Maria will finalize the "
        "design mockups by Friday, Carlos will set up the CI/CD pipeline by "
        "next Wednesday, and Ahmed will complete the database schema review "
        "by end of week."
    )
    doc.text.addElement(p2)

    p3 = P(
        text="The next meeting is scheduled for Thursday at 2:00 PM. Each "
        "team lead should prepare a brief status update covering progress, "
        "blockers, and any resource needs."
    )
    doc.text.addElement(p3)

    path = os.path.join(OFFICE_DIR, "notes.odt")
    doc.save(path)
    print(f"  Created {path}")


def generate_chart_png():
    """Create chart.png as a 400x300 bar chart image."""
    width, height = 400, 300
    img = Image.new("RGB", (width, height), "#FFFFFF")
    draw = ImageDraw.Draw(img)

    # Chart area
    margin_left, margin_bottom, margin_top, margin_right = 60, 40, 30, 20
    chart_left = margin_left
    chart_bottom = height - margin_bottom
    chart_top = margin_top
    chart_right = width - margin_right

    # Draw axes
    draw.line(
        [(chart_left, chart_top), (chart_left, chart_bottom)], fill="#333333", width=2
    )
    draw.line(
        [(chart_left, chart_bottom), (chart_right, chart_bottom)],
        fill="#333333",
        width=2,
    )

    # Bar data
    labels = ["Q1", "Q2", "Q3", "Q4"]
    values = [125, 132, 118, 153]
    colors = ["#4A90D9", "#50C878", "#FFB347", "#FF6B6B"]
    max_val = max(values)

    bar_count = len(values)
    bar_area_width = chart_right - chart_left
    bar_width = bar_area_width // (bar_count * 2)
    spacing = bar_width

    for i, (label, value, color) in enumerate(zip(labels, values, colors)):
        bar_height = int((value / max_val) * (chart_bottom - chart_top - 10))
        x0 = chart_left + spacing + i * (bar_width + spacing)
        y0 = chart_bottom - bar_height
        x1 = x0 + bar_width
        y1 = chart_bottom

        draw.rectangle([x0, y0, x1, y1], fill=color)

        # Label below bar
        lw = draw.textlength(label)
        draw.text(
            (x0 + (bar_width - lw) / 2, y1 + 5), label, fill="#333333"
        )

        # Value above bar
        val_text = str(value)
        vw = draw.textlength(val_text)
        draw.text(
            (x0 + (bar_width - vw) / 2, y0 - 15), val_text, fill="#333333"
        )

    # Y-axis labels
    for tick_val in [0, 50, 100, 150]:
        y = chart_bottom - int((tick_val / max_val) * (chart_bottom - chart_top - 10))
        tick_text = str(tick_val)
        tw = draw.textlength(tick_text)
        draw.text((chart_left - tw - 8, y - 6), tick_text, fill="#333333")
        draw.line(
            [(chart_left - 3, y), (chart_left, y)], fill="#333333", width=1
        )

    # Title
    title = "Quarterly Revenue (k EUR)"
    tw = draw.textlength(title)
    draw.text(((width - tw) / 2, 5), title, fill="#333333")

    path = os.path.join(OFFICE_DIR, "chart.png")
    img.save(path)
    print(f"  Created {path}")


def generate_projects_csv():
    """Create projects.csv with project data matching data.csv employees."""
    path = os.path.join(SYNTHETIC_DIR, "projects.csv")
    with open(path, "w") as f:
        f.write("Project,Start Date,Status\n")
        f.write("Atlas,2026-01-15,Active\n")
        f.write("Beacon,2026-02-01,Planning\n")
    print(f"  Created {path}")


def generate_budget_notes_txt():
    """Create budget_notes.txt with monthly budget info."""
    path = os.path.join(SYNTHETIC_DIR, "budget_notes.txt")
    with open(path, "w") as f:
        f.write(
            "Monthly budget: Entertainment 500 EUR, Groceries 800 EUR, "
            "Transport 200 EUR, Bills 2000 EUR.\n"
        )
    print(f"  Created {path}")


def generate_client_data_csv():
    """Create client_data.csv with one client record."""
    path = os.path.join(SYNTHETIC_DIR, "client_data.csv")
    with open(path, "w") as f:
        f.write("Company,Date,Product,Quantity,Unit Price\n")
        f.write("Acme Corp,2026-03-01,Widget,50,24.99\n")
    print(f"  Created {path}")


def generate_invoice_template_docx():
    """Create invoice_template.docx with placeholder fields."""
    doc = Document()
    doc.add_heading("Invoice", level=1)
    doc.add_paragraph("Company: {{company}}")
    doc.add_paragraph("Date: {{date}}")
    doc.add_paragraph("")

    table = doc.add_table(rows=2, cols=4)
    table.style = "Table Grid"
    for i, header in enumerate(["Item", "Quantity", "Unit Price", "Subtotal"]):
        table.rows[0].cells[i].text = header
    table.rows[1].cells[0].text = "{{items_table}}"

    doc.add_paragraph("")
    doc.add_paragraph("Total: {{total}}")

    path = os.path.join(SYNTHETIC_DIR, "invoice_template.docx")
    doc.save(path)
    print(f"  Created {path}")


def generate_quarterly_data_xlsx():
    """Create quarterly_data.xlsx with 4 quarter sheets, deterministic amounts."""
    wb = Workbook()

    quarters = {
        "Q1": [
            ("2026-01-05", "Office Supplies", 120.50, "Printer paper and toner"),
            ("2026-01-12", "Travel", 450.00, "Train tickets to Madrid"),
            ("2026-01-20", "Software", 299.99, "Annual IDE license"),
            ("2026-02-03", "Office Supplies", 85.00, "Desk organizers"),
            ("2026-02-14", "Marketing", 1200.00, "Online ad campaign"),
            ("2026-02-28", "Travel", 380.00, "Hotel for conference"),
            ("2026-03-05", "Software", 150.00, "Cloud storage upgrade"),
            ("2026-03-10", "Office Supplies", 65.50, "Whiteboard markers"),
            ("2026-03-15", "Utilities", 210.00, "Electricity bill"),
            ("2026-03-20", "Marketing", 800.00, "Print brochures"),
            ("2026-03-25", "Travel", 520.00, "Flight to Barcelona"),
            ("2026-03-31", "Utilities", 180.00, "Internet service"),
        ],
        "Q2": [
            ("2026-04-02", "Office Supplies", 95.00, "Notebooks and pens"),
            ("2026-04-10", "Travel", 610.00, "Flight to Lisbon"),
            ("2026-04-18", "Software", 499.99, "Project management tool"),
            ("2026-05-01", "Marketing", 950.00, "Social media ads"),
            ("2026-05-12", "Utilities", 195.00, "Electricity bill"),
            ("2026-05-20", "Office Supplies", 110.00, "Ergonomic mouse"),
            ("2026-06-01", "Travel", 340.00, "Train to Seville"),
            ("2026-06-10", "Software", 75.00, "Design tool subscription"),
            ("2026-06-15", "Marketing", 1500.00, "Trade show booth"),
            ("2026-06-22", "Utilities", 205.00, "Internet and phone"),
            ("2026-06-28", "Office Supplies", 48.00, "Cable management"),
            ("2026-06-30", "Travel", 425.00, "Hotel for workshop"),
        ],
        "Q3": [
            ("2026-07-05", "Software", 299.99, "Renewed IDE license"),
            ("2026-07-12", "Travel", 780.00, "Flight to Berlin"),
            ("2026-07-20", "Office Supplies", 130.00, "Monitor stand"),
            ("2026-08-01", "Marketing", 1100.00, "Email campaign platform"),
            ("2026-08-10", "Utilities", 220.00, "Electricity bill"),
            ("2026-08-18", "Travel", 290.00, "Train to Valencia"),
            ("2026-08-25", "Office Supplies", 72.00, "USB hubs"),
            ("2026-09-02", "Software", 180.00, "Security audit tool"),
            ("2026-09-10", "Marketing", 650.00, "Promotional merchandise"),
            ("2026-09-15", "Utilities", 198.00, "Internet service"),
            ("2026-09-22", "Travel", 510.00, "Hotel for summit"),
            ("2026-09-30", "Office Supplies", 55.00, "Desk lamp"),
        ],
        "Q4": [
            ("2026-10-03", "Marketing", 2000.00, "Year-end campaign"),
            ("2026-10-12", "Travel", 690.00, "Flight to Paris"),
            ("2026-10-20", "Software", 350.00, "Analytics platform"),
            ("2026-11-01", "Office Supplies", 145.00, "Keyboard and mouse set"),
            ("2026-11-10", "Utilities", 230.00, "Electricity bill"),
            ("2026-11-18", "Travel", 415.00, "Train to Milan"),
            ("2026-11-25", "Marketing", 750.00, "Holiday promotions"),
            ("2026-12-01", "Software", 125.00, "Backup service renewal"),
            ("2026-12-08", "Office Supplies", 88.00, "Headset"),
            ("2026-12-15", "Utilities", 210.00, "Internet and phone"),
            ("2026-12-20", "Travel", 560.00, "Hotel for year-end meeting"),
            ("2026-12-28", "Marketing", 400.00, "Thank-you gifts for clients"),
        ],
    }

    first = True
    for quarter, rows in quarters.items():
        if first:
            ws = wb.active
            ws.title = quarter
            first = False
        else:
            ws = wb.create_sheet(quarter)

        ws.append(["Date", "Category", "Amount", "Description"])
        for row in rows:
            ws.append(list(row))

    path = os.path.join(SYNTHETIC_DIR, "quarterly_data.xlsx")
    wb.save(path)
    print(f"  Created {path}")


def sync_support_office():
    """Copy all files from engine/testdata/office/ to app/integration_test/support/office/."""
    os.makedirs(SUPPORT_OFFICE_DIR, exist_ok=True)
    for name in os.listdir(OFFICE_DIR):
        src = os.path.join(OFFICE_DIR, name)
        dst = os.path.join(SUPPORT_OFFICE_DIR, name)
        shutil.copy2(src, dst)
    print(f"  Synced {OFFICE_DIR} → {SUPPORT_OFFICE_DIR}")


def main():
    os.makedirs(OFFICE_DIR, exist_ok=True)
    os.makedirs(SYNTHETIC_DIR, exist_ok=True)

    print("Generating office fixtures...")
    generate_simple_docx()
    generate_multi_sheet_xlsx()
    generate_slides_pptx()
    generate_report_pdf()
    generate_notes_odt()
    generate_chart_png()

    print("\nGenerating synthetic data...")
    generate_projects_csv()
    generate_budget_notes_txt()
    generate_client_data_csv()
    generate_invoice_template_docx()
    generate_quarterly_data_xlsx()

    print("\nSyncing to support/office/...")
    sync_support_office()

    print("\nDone.")


if __name__ == "__main__":
    main()
