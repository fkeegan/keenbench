#!/usr/bin/env python3
"""Generate Workbench Context soap-opera fixtures for manual QA.

Writes binary fixtures into:
  - docs/test/fixtures/workbench-context/

Writes large payloads (gitignored) into:
  - artifacts/testdata/

Run with the tool-worker venv so dependencies are available:
  engine/tools/pyworker/.venv/bin/python scripts/testdata/generate_workbench_context_fixtures.py
"""

from __future__ import annotations

import os
import subprocess
import sys

try:
    from docx import Document
    from openpyxl import Workbook
    from pptx import Presentation
    from pptx.util import Inches as PptxInches
    from PIL import Image, ImageDraw, ImageFont
except Exception as exc:  # pragma: no cover
    print("ERROR: Missing dependencies for fixture generation.")
    print("Run using: engine/tools/pyworker/.venv/bin/python ...")
    print(f"Import error: {exc}")
    sys.exit(1)


ROOT = os.path.abspath(os.path.join(os.path.dirname(__file__), "..", ".."))
FIXTURE_DIR = os.path.join(ROOT, "docs", "test", "fixtures", "workbench-context")
ARTIFACTS_DIR = os.path.join(ROOT, "artifacts", "testdata")


def _read_fixture_text(name: str) -> str:
    path = os.path.join(FIXTURE_DIR, name)
    with open(path, "r", encoding="utf-8") as f:
        return f.read().strip()


def _write_text(path: str, content: str) -> None:
    os.makedirs(os.path.dirname(path), exist_ok=True)
    with open(path, "w", encoding="utf-8") as f:
        f.write(content)


def generate_company_overview_docx() -> str:
    content = _read_fixture_text("keenbench_company_context_v1.txt")
    doc = Document()
    doc.add_heading("KeenBench Company Overview (Fictional QA)", level=1)
    doc.add_paragraph(
        "This document is test data for KeenBench Workbench Context file-mode processing."
    )
    doc.add_heading("Source Text", level=2)
    for line in content.splitlines():
        doc.add_paragraph(line)
    path = os.path.join(FIXTURE_DIR, "keenbench_company_overview.docx")
    doc.save(path)
    return path


def generate_engineering_department_brief_docx() -> str:
    content = _read_fixture_text("keenbench_department_engineering_brief.txt")
    doc = Document()
    doc.add_heading("KeenBench Engineering Department Brief (Fictional QA)", level=1)
    doc.add_heading("Source Text", level=2)
    for line in content.splitlines():
        doc.add_paragraph(line)
    path = os.path.join(FIXTURE_DIR, "keenbench_engineering_department_brief.docx")
    doc.save(path)
    return path


def generate_company_metrics_xlsx() -> str:
    wb = Workbook()

    ws = wb.active
    ws.title = "KPIs"
    ws.append(["QA marker", "KEENBENCH_METRICS_XLSX=V1"])
    ws.append([])
    ws.append(["Metric", "2026-01", "2026-02", "Target"])
    ws.append(["Weekly active workbenches", 12, 19, 25])
    ws.append(["Drafts published", 7, 11, 15])
    ws.append(["Median review time (min)", 6.5, 5.2, 5.0])

    ws2 = wb.create_sheet("Pipeline")
    ws2.append(["Company", "Stage", "Owner", "Notes"])
    ws2.append(["Design Partner A", "Active", "Priya Desai", "Prefers auditability"])
    ws2.append(["Design Partner B", "Active", "Mira Kwon", "Security-conscious"])
    ws2.append(["Design Partner C", "Prospect", "Priya Desai", "Ops-heavy team"])

    path = os.path.join(FIXTURE_DIR, "keenbench_company_metrics.xlsx")
    wb.save(path)
    return path


def generate_pitch_deck_pptx() -> str:
    prs = Presentation()

    # Slide 1: Title
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = "KeenBench"
    slide.placeholders[1].text = "Fictional Pitch Deck (QA)"

    # Slide 2: Problem
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Problem"
    tf = slide.placeholders[1].text_frame
    tf.text = "Turning source material into deliverables is slow and error-prone"
    for bullet in [
        "Inputs are messy: notes, PDFs, spreadsheets, slide decks",
        "Edits are hard to review and hard to audit",
        "Teams need consistency across documents",
    ]:
        p = tf.add_paragraph()
        p.text = bullet

    # Slide 3: Solution
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Solution"
    tf = slide.placeholders[1].text_frame
    tf.text = "KeenBench: assisted drafting with guardrails"
    for bullet in [
        "Workbench holds files + persistent context",
        "Workshop produces a Draft (not auto-publish)",
        "Review shows offline diffs before Publish",
        "Explicit egress consent before any model call",
    ]:
        p = tf.add_paragraph()
        p.text = bullet

    # Slide 4: Team
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Team (Fictional)"
    tf = slide.placeholders[1].text_frame
    tf.text = "Mira Kwon - CEO (Ops; also Sales/Growth)"
    for bullet in [
        "Ethan Park - CTO (Engineering)",
        "Priya Desai - Sales Rep (Sales/Growth)",
        "Luis Alvarez - Lead Engineer (Engineering)",
        "QA marker: KEENBENCH_PITCH_DECK=PPTX_V1",
    ]:
        p = tf.add_paragraph()
        p.text = bullet

    # Slide 5: Footer marker
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
    tx_box = slide.shapes.add_textbox(PptxInches(0.8), PptxInches(3.0), PptxInches(8.0), PptxInches(1.0))
    tx_box.text_frame.text = "QA marker: KEENBENCH_PITCH_DECK=PPTX_V1"

    path = os.path.join(FIXTURE_DIR, "keenbench_pitch_deck.pptx")
    prs.save(path)
    return path


def generate_logo_png() -> str:
    width, height = 512, 512
    img = Image.new("RGB", (width, height), "#5B7FC2")
    draw = ImageDraw.Draw(img)
    font = ImageFont.load_default()
    text = "KEENBENCH"
    bbox = draw.textbbox((0, 0), text, font=font)
    tw = bbox[2] - bbox[0]
    th = bbox[3] - bbox[1]
    x = (width - tw) // 2
    y = (height - th) // 2
    draw.text((x, y), text, fill="#FFFFFF", font=font)

    path = os.path.join(FIXTURE_DIR, "keenbench_logo.png")
    img.save(path)
    return path


def generate_unknown_bin() -> str:
    payload = b"KEENBENCH\x00UNKNOWN_BIN\x01\x02\x03\x04\x05\n"
    path = os.path.join(FIXTURE_DIR, "keenbench_unknown.bin")
    with open(path, "wb") as f:
        f.write(payload)
    return path


def generate_security_overview_pdf() -> str:
    content = _read_fixture_text("keenbench_security_overview.txt")
    doc = Document()
    doc.add_heading("KeenBench Security & Data Handling (Fictional QA)", level=1)
    for line in content.splitlines():
        doc.add_paragraph(line)

    tmp_docx = os.path.join(FIXTURE_DIR, "_keenbench_security_overview_tmp.docx")
    pdf_dst = os.path.join(FIXTURE_DIR, "keenbench_security_overview.pdf")
    try:
        doc.save(tmp_docx)
        result = subprocess.run(
            [
                "/snap/bin/libreoffice",
                "--headless",
                "--convert-to",
                "pdf",
                "--outdir",
                FIXTURE_DIR,
                tmp_docx,
            ],
            capture_output=True,
            text=True,
            timeout=60,
        )
        tmp_pdf = os.path.join(FIXTURE_DIR, "_keenbench_security_overview_tmp.pdf")
        if os.path.exists(tmp_pdf):
            os.replace(tmp_pdf, pdf_dst)
        else:
            print("ERROR: LibreOffice conversion failed.")
            print(result.stderr.strip() or result.stdout.strip())
            raise RuntimeError("LibreOffice did not produce expected PDF output")
    finally:
        if os.path.exists(tmp_docx):
            os.remove(tmp_docx)
    return pdf_dst


def generate_context_clutter_payload() -> str:
    alpha = _read_fixture_text("keenbench_situation_alpha.txt")
    filler = "token "
    # 90k repeats -> ~540k chars -> estimated weight ~135k tokens -> context_share ~0.675 on 200k context
    payload = (
        alpha
        + "\n\n"
        + "Clutter payload filler (do not edit):\n"
        + (filler * 90_000)
        + "\n"
    )
    path = os.path.join(ARTIFACTS_DIR, "keenbench_situation_clutter_payload.md")
    _write_text(path, payload)
    return path


def generate_oversize_csv() -> str:
    target_bytes = 26 * 1024 * 1024  # > 25MB limit
    path = os.path.join(ARTIFACTS_DIR, "keenbench_oversize_roster.csv")
    os.makedirs(os.path.dirname(path), exist_ok=True)
    header = "name,role,department,location,availability,qa_marker\n"
    row = "Mira Kwon,CEO,Ops,Remote,Full-time,KEENBENCH_OVERSIZE_TEST\n"
    written = 0
    with open(path, "w", encoding="utf-8") as f:
        f.write(header)
        written += len(header.encode("utf-8"))
        while written < target_bytes:
            f.write(row)
            written += len(row.encode("utf-8"))
    return path


def main() -> int:
    os.makedirs(FIXTURE_DIR, exist_ok=True)
    os.makedirs(ARTIFACTS_DIR, exist_ok=True)

    created = []
    created.append(generate_company_overview_docx())
    created.append(generate_engineering_department_brief_docx())
    created.append(generate_company_metrics_xlsx())
    created.append(generate_pitch_deck_pptx())
    created.append(generate_security_overview_pdf())
    created.append(generate_logo_png())
    created.append(generate_unknown_bin())
    created.append(generate_context_clutter_payload())
    created.append(generate_oversize_csv())

    print("Created:")
    for path in created:
        rel = os.path.relpath(path, ROOT)
        size = os.path.getsize(path)
        print(f"  - {rel} ({size} bytes)")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())

